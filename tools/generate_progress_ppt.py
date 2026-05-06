from pathlib import Path
import json
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / 'PGA_TEAM_DiTing_日本强震数据_进展汇报.pptx'

BASE_DIR = Path('/tmp/japan_build_test3')
STATION_CACHE_CSV = BASE_DIR / 'station_cache_d9a7d4d6.csv'
SPLIT_EVENT_CSV = BASE_DIR / 'split_events.csv'
SPLIT_STATION_CSV = BASE_DIR / 'split_stations.csv'
IMG_FULL_SCALE = ROOT / 'ppt_figures/data_overview_2024_overfit_ppt/full_dataset_scale_ppt.png'
IMG_MAG = ROOT / 'ppt_figures/data_overview_2024_overfit_ppt/magnitude_coverage_ppt.png'
IMG_GEO = ROOT / 'ppt_figures/data_overview_2024_overfit_ppt/geographic_coverage_ppt.png'
IMG_SPLIT = ROOT / 'ppt_figures/data_overview_2024_overfit_ppt/split_and_overfit_distribution_ppt.png'

# palette
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(247, 249, 252)
NAVY = RGBColor(22, 47, 78)
NAVY2 = RGBColor(31, 64, 104)
TEXT = RGBColor(30, 41, 59)
MUTED = RGBColor(100, 116, 139)
LINE = RGBColor(226, 232, 240)
BLUE = RGBColor(37, 99, 235)
SKY = RGBColor(14, 165, 233)
TEAL = RGBColor(13, 148, 136)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(245, 158, 11)
ORANGE = RGBColor(234, 88, 12)
RED = RGBColor(220, 38, 38)
PURPLE = RGBColor(124, 58, 237)
PALE_BLUE = RGBColor(239, 246, 255)
PALE_SKY = RGBColor(240, 249, 255)
PALE_TEAL = RGBColor(240, 253, 250)
PALE_GREEN = RGBColor(240, 253, 244)
PALE_ORANGE = RGBColor(255, 247, 237)
PALE_RED = RGBColor(254, 242, 242)
PALE_PURPLE = RGBColor(245, 243, 255)
PALE_GRAY = RGBColor(248, 250, 252)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

station_cache_df = pd.read_csv(STATION_CACHE_CSV)
split_event_all_df = pd.read_csv(SPLIT_EVENT_CSV)
split_station_all_df = pd.read_csv(SPLIT_STATION_CSV)

for df in [station_cache_df, split_event_all_df, split_station_all_df]:
    if 'EVENT' in df.columns:
        df['EVENT'] = df['EVENT'].astype(str)

full_events_df = station_cache_df.drop_duplicates('EVENT').copy()
overfit_events_df = split_event_all_df[split_event_all_df['is_overfit_selected']].copy()
overfit_stations_df = split_station_all_df[split_station_all_df['is_overfit_selected']].copy()

split_station_df = overfit_stations_df.groupby('split').size().rename('n_station_rows').reset_index()

notes = {
    'full_available_events': int(len(full_events_df)),
    'full_available_station_rows': int(len(station_cache_df)),
    'full_magnitude_range': [float(full_events_df['Magnitude'].min()), float(full_events_df['Magnitude'].max())],
    'full_network_counts': {k: int(v) for k, v in station_cache_df['source_network'].value_counts().sort_index().to_dict().items()},
    'overfit_selected_event_count': int(len(overfit_events_df)),
    'overfit_selected_mag_range': [float(overfit_events_df['Magnitude'].min()), float(overfit_events_df['Magnitude'].max())],
    'overfit_events_by_split': {k: int(v) for k, v in overfit_events_df['split'].value_counts().to_dict().items()},
}


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, x, y, w, h, text='', size=14, color=TEXT, bold=False,
                align=PP_ALIGN.LEFT, valign=MSO_VERTICAL_ANCHOR.TOP, font='Microsoft YaHei'):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    if text:
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tx


def add_shape_box(slide, x, y, w, h, fill=WHITE, line=LINE, radius=True):
    t = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(t, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1)
    return shp


def write_in_shape(shape, lines, sizes=None, colors=None, bolds=None, aligns=None, margins=(10, 10, 8, 6)):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(margins[0])
    tf.margin_right = Pt(margins[1])
    tf.margin_top = Pt(margins[2])
    tf.margin_bottom = Pt(margins[3])
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = aligns[i] if aligns else PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.name = 'Microsoft YaHei'
        r.font.size = Pt(sizes[i] if sizes else 14)
        r.font.color.rgb = colors[i] if colors else TEXT
        r.font.bold = bolds[i] if bolds else False


def add_header(slide, sec_no, title, subtitle=''):
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.16))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    badge = add_shape_box(slide, Inches(0.55), Inches(0.34), Inches(0.68), Inches(0.44), fill=NAVY, line=NAVY)
    write_in_shape(badge, [sec_no], sizes=[16], colors=[WHITE], bolds=[True], aligns=[PP_ALIGN.CENTER], margins=(0, 0, 6, 0))

    add_textbox(slide, Inches(1.38), Inches(0.28), Inches(9.9), Inches(0.42), title, size=24, color=NAVY, bold=True)
    if subtitle:
        add_textbox(slide, Inches(1.40), Inches(0.73), Inches(10.4), Inches(0.22), subtitle, size=10.5, color=MUTED)

    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.02), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def add_footer(slide, page_no, note=''):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(7.08), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()

    circle = add_shape_box(slide, Inches(12.18), Inches(6.88), Inches(0.55), Inches(0.34), fill=NAVY, line=NAVY)
    write_in_shape(circle, [str(page_no)], sizes=[12], colors=[WHITE], bolds=[True], aligns=[PP_ALIGN.CENTER], margins=(0, 0, 5, 0))
    if note:
        add_textbox(slide, Inches(0.58), Inches(7.0), Inches(11.3), Inches(0.14), note, size=8.5, color=MUTED)


def add_bullets(slide, x, y, w, h, lines, title=None, title_color=NAVY, font_size=15, bullet_color=TEXT):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    if title:
        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = title
        r0.font.name = 'Microsoft YaHei'
        r0.font.size = Pt(font_size + 1)
        r0.font.bold = True
        r0.font.color.rgb = title_color
    else:
        tf.clear()
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if (i > 0 or title) else tf.paragraphs[0]
        p.bullet = True
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = line
        r.font.name = 'Microsoft YaHei'
        r.font.size = Pt(font_size)
        r.font.color.rgb = bullet_color
    return tx


def add_chip(slide, x, y, w, h, text, fill, color):
    shp = add_shape_box(slide, x, y, w, h, fill=fill, line=fill)
    write_in_shape(shp, [text], sizes=[10.5], colors=[color], bolds=[True], aligns=[PP_ALIGN.CENTER], margins=(0, 0, 6, 0))
    return shp


def add_metric_card(slide, x, y, w, h, value, label, accent, fill=WHITE):
    shadow = add_shape_box(slide, x + Inches(0.03), y + Inches(0.04), w, h, fill=PALE_GRAY, line=PALE_GRAY)
    shadow.line.fill.background()
    box = add_shape_box(slide, x, y, w, h, fill=fill, line=LINE)
    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(0.08), h)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    write_in_shape(box, [str(value), label], sizes=[22, 10.5], colors=[accent, MUTED], bolds=[True, False])
    return box


def add_arrow(slide, x1, y1, x2, y2, color=MUTED, width=1.5):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_stage(slide, x, y, w, h, title, subtitle, fill, accent):
    box = add_shape_box(slide, x, y, w, h, fill=fill, line=LINE)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(10)
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = 'Microsoft YaHei'
    r1.font.size = Pt(17)
    r1.font.bold = True
    r1.font.color.rgb = accent
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.name = 'Microsoft YaHei'
    r2.font.size = Pt(11.5)
    r2.font.color.rgb = TEXT
    return box


def add_picture_framed(slide, path, x, y, w, h=None, caption=None):
    frame = add_shape_box(slide, x, y, w, h if h else Inches(2), fill=WHITE, line=LINE)
    pic = slide.shapes.add_picture(str(path), x + Inches(0.04), y + Inches(0.04), width=w - Inches(0.08), height=(h - Inches(0.08) if h else None))
    if caption:
        add_textbox(slide, x + Inches(0.06), y + h + Inches(0.05), w - Inches(0.12), Inches(0.18), caption, size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
    return frame, pic


# Slide 1: cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG)

left_bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(7.2), Inches(7.5))
left_bg.fill.solid()
left_bg.fill.fore_color.rgb = NAVY
left_bg.line.fill.background()

accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(7.2), Inches(0.18))
accent.fill.solid()
accent.fill.fore_color.rgb = SKY
accent.line.fill.background()

add_chip(slide, Inches(0.72), Inches(0.62), Inches(1.45), Inches(0.35), '阶段性进展汇报', PALE_SKY, SKY)
add_textbox(slide, Inches(0.72), Inches(1.22), Inches(5.8), Inches(1.35), '基于 TEAM + DiTing 的\n日本强震数据 PGA 估计模型', size=28, color=WHITE, bold=True)
add_textbox(slide, Inches(0.75), Inches(2.75), Inches(5.9), Inches(0.45), '模型接入 · 数据准备 · small-sample 调试 · 输入排查', size=13.5, color=RGBColor(226, 232, 240))

summary_panel = add_shape_box(slide, Inches(0.72), Inches(3.45), Inches(5.95), Inches(2.55), fill=NAVY2, line=NAVY2)
write_in_shape(
    summary_panel,
    [
        '当前阶段总结',
        '• 保留 TEAM 多台站框架，将原 waveform model 替换为 DiTing encoder。',
        '• 训练数据切换为自建日本强震数据，已完成基础清洗与概览统计。',
        '• 当前重点是 small-sample / overfit-style 调试实验。',
        '• PGA 训练暂未收敛，已开始基于模型输入 dump 的系统排查。',
    ],
    sizes=[18, 13.2, 13.2, 13.2, 13.2],
    colors=[WHITE, WHITE, WHITE, WHITE, WHITE],
    bolds=[True, False, False, False, False],
)

add_textbox(slide, Inches(0.75), Inches(6.5), Inches(5.7), Inches(0.22), 'team_pytorch/train_light.py · gemini_models.py', size=10.5, color=RGBColor(203, 213, 225))

right_title = add_textbox(slide, Inches(7.65), Inches(0.72), Inches(4.8), Inches(0.28), '进展快照', size=18, color=NAVY, bold=True)
add_metric_card(slide, Inches(7.65), Inches(1.25), Inches(2.1), Inches(1.08), '模型接入', 'TEAM + DiTing 已完成', BLUE, fill=WHITE)
add_metric_card(slide, Inches(9.95), Inches(1.25), Inches(2.35), Inches(1.08), '数据准备', '日本强震数据已整理', GREEN, fill=WHITE)
add_metric_card(slide, Inches(7.65), Inches(2.62), Inches(4.65), Inches(1.08), '当前问题', 'PGA 训练结果暂未收敛', RED, fill=WHITE)
add_metric_card(slide, Inches(7.65), Inches(4.0), Inches(4.65), Inches(1.08), '本次重点', '架构 / 数据 / 排查进展', ORANGE, fill=WHITE)

agenda = add_shape_box(slide, Inches(7.65), Inches(5.6), Inches(4.65), Inches(1.1), fill=WHITE, line=LINE)
write_in_shape(
    agenda,
    ['汇报结构', '1) 模型架构  2) 训练实现  3) 数据概览  4) 当前问题与排查  5) 下一步'],
    sizes=[16, 11.6], colors=[NAVY, TEXT], bolds=[True, False]
)

# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_header(slide, '01', '本次汇报内容与当前阶段定位')

stage_y = Inches(1.45)
boxes = [
    ('任务目标', '面向日本强震数据做 PGA 估计；以 TEAM 作为总体多台站框架。', PALE_BLUE, BLUE),
    ('模型改造', '保留 TEAM 的多台站融合逻辑，用 DiTing 替换 waveform encoder。', PALE_TEAL, TEAL),
    ('数据实验', '先做 small-sample / overfit-style 调试，验证训练链路是否合理。', PALE_GREEN, GREEN),
    ('当前状态', '训练已跑通，但 PGA 暂未收敛；当前重点是定位收敛瓶颈。', PALE_ORANGE, ORANGE),
]
xs = [Inches(0.85), Inches(3.95), Inches(7.05), Inches(10.15)]
for i, (title, subtitle, fill, accent) in enumerate(boxes):
    add_stage(slide, xs[i], stage_y, Inches(2.45), Inches(1.7), title, subtitle, fill, accent)
    if i < len(boxes) - 1:
        add_arrow(slide, xs[i] + Inches(2.45), stage_y + Inches(0.85), xs[i + 1], stage_y + Inches(0.85), color=MUTED, width=1.8)

left = add_shape_box(slide, Inches(0.85), Inches(3.0), Inches(5.85), Inches(3.35), fill=WHITE, line=LINE)
add_bullets(
    slide, Inches(1.08), Inches(3.28), Inches(5.4), Inches(2.8),
    [
        '这次汇报不是报告最终模型性能，而是汇报“从模型到数据到排查”的阶段性进展。',
        '当前最重要的是把 TEAM + DiTing 的训练链路打通，并确认进入模型的数据是可信的。',
        '因此本阶段优先看：架构改造是否合理、数据准备是否到位、小样本实验是否可解释。',
    ],
    title='当前阶段工作的核心目标', title_color=NAVY, font_size=14.2
)

right = add_shape_box(slide, Inches(7.0), Inches(3.0), Inches(5.35), Inches(3.35), fill=WHITE, line=LINE)
add_chip(slide, Inches(7.28), Inches(3.28), Inches(1.28), Inches(0.3), '为什么先做小样本', PALE_ORANGE, ORANGE)
add_bullets(
    slide, Inches(7.22), Inches(3.72), Inches(4.9), Inches(2.4),
    [
        '如果连小样本都难以拟合，通常优先说明训练链路仍存在问题。',
        '小样本更适合快速定位标签定义、冻结策略、学习率或输入对齐异常。',
        '因此当前把它作为 sanity check，而不是最终实验结论。',
    ],
    font_size=13.6
)
add_footer(slide, 2, '本次不展开 DiTing 内部网络细节，重点说明其在 TEAM 中的接入方式')

# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_header(slide, '02', '模型架构：从 TEAM 到 TEAM + DiTing', '保留 TEAM 多台站主干，替换 waveform encoder')

left_panel = add_shape_box(slide, Inches(0.7), Inches(1.35), Inches(5.95), Inches(5.3), fill=WHITE, line=LINE)
right_panel = add_shape_box(slide, Inches(6.95), Inches(1.35), Inches(5.7), Inches(5.3), fill=WHITE, line=LINE)
add_chip(slide, Inches(0.95), Inches(1.58), Inches(1.35), Inches(0.3), '原始 TEAM', PALE_BLUE, BLUE)
add_chip(slide, Inches(7.18), Inches(1.58), Inches(1.6), Inches(0.3), '当前改造版本', PALE_GREEN, GREEN)

# architecture blocks
left_xs = [0.95, 2.10, 3.28, 4.47, 5.58]
right_xs = [7.18, 8.33, 9.48, 10.63, 11.75]
arch_y = Inches(2.45)
labels_left = ['多台站\nwaveforms', 'TEAM\nwaveform model', 'station\nembedding', '位置编码 +\nTransformer', 'PGA / Mag / Loc\nheads']
labels_right = ['多台站\nwaveforms', 'DiTing\nencoder', 'station\nadapter', 'TEAM 位置编码\n+ Transformer', 'PGA\nhead']
left_fills = [PALE_GRAY, PALE_BLUE, PALE_GREEN, PALE_TEAL, PALE_ORANGE]
right_fills = [PALE_GRAY, PALE_BLUE, PALE_GREEN, PALE_TEAL, PALE_ORANGE]
for i in range(5):
    add_shape_box(slide, Inches(left_xs[i]), arch_y, Inches(0.9), Inches(1.0), fill=left_fills[i], line=LINE)
    add_shape_box(slide, Inches(right_xs[i]), arch_y, Inches(0.9), Inches(1.0), fill=right_fills[i], line=LINE)
    write_in_shape(slide.shapes[-2], [labels_left[i]], sizes=[12.2], colors=[TEXT], bolds=[True], aligns=[PP_ALIGN.CENTER], margins=(2, 2, 20, 0))
    write_in_shape(slide.shapes[-1], [labels_right[i]], sizes=[12.2], colors=[TEXT], bolds=[True], aligns=[PP_ALIGN.CENTER], margins=(2, 2, 20, 0))
    if i < 4:
        add_arrow(slide, Inches(left_xs[i] + 0.9), arch_y + Inches(0.5), Inches(left_xs[i + 1]), arch_y + Inches(0.5), width=1.6)
        add_arrow(slide, Inches(right_xs[i] + 0.9), arch_y + Inches(0.5), Inches(right_xs[i + 1]), arch_y + Inches(0.5), width=1.6)

add_chip(slide, Inches(8.28), Inches(2.06), Inches(1.0), Inches(0.24), '默认冻结', PALE_RED, RED)
add_chip(slide, Inches(9.43), Inches(2.06), Inches(1.0), Inches(0.24), '当前训练', PALE_GREEN, GREEN)

add_bullets(
    slide, Inches(1.0), Inches(4.05), Inches(5.25), Inches(2.1),
    [
        '原始 TEAM 先对每个台站的波形做编码，再通过位置编码和 Transformer 聚合多台站信息。',
        '这一设计保留了“单台站表征 + 多台站融合”的清晰层次。',
        '因此适合作为接入新 waveform encoder 的基础框架。',
    ],
    title='为什么选择 TEAM 作为外层框架', title_color=NAVY, font_size=13.4
)

add_bullets(
    slide, Inches(7.2), Inches(4.05), Inches(5.1), Inches(2.3),
    [
        'gemini_models.py 中，waveform_model 现在由 [DiTing encoder, station adapter] 组成。',
        'TEAM 的位置编码、Transformer 与 PGA head 保持不变。',
        'train_light.py 中当前默认冻结 DiTing encoder，先训练 station adapter + TEAM 上层。',
    ],
    title='当前实现要点', title_color=NAVY, font_size=13.2
)
add_footer(slide, 3, '代码参考：gemini_models.py:get_diting_model / build_full_model；train_light.py: 冻结 waveform_model[0]')

# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_header(slide, '03', '训练实现与当前训练策略')

left = add_shape_box(slide, Inches(0.72), Inches(1.35), Inches(6.1), Inches(5.55), fill=WHITE, line=LINE)
add_chip(slide, Inches(0.98), Inches(1.62), Inches(1.55), Inches(0.3), '训练实现', PALE_BLUE, BLUE)
add_bullets(
    slide, Inches(0.98), Inches(2.02), Inches(5.55), Inches(4.2),
    [
        '训练入口为 team_pytorch/train_light.py，支持 torchrun / DDP。',
        'DiTing 的配置由 YAML 读入，再构建 waveform encoder。',
        '优化器参数组区分 adapter / team / encoder，可分别设置学习率。',
        '支持 single-station pretrain 与 station pretrain 权重加载。',
        '训练时还能导出 model_input_dumps，便于对 batch 级输入做可视化排查。',
    ],
    font_size=14.1
)

right = add_shape_box(slide, Inches(7.02), Inches(1.35), Inches(5.58), Inches(5.55), fill=WHITE, line=LINE)
add_chip(slide, Inches(7.28), Inches(1.62), Inches(1.7), Inches(0.3), '当前训练对象', PALE_GREEN, GREEN)
for (x, title, sub, fill, accent) in [
    (7.32, 'DiTing encoder', '冻结', PALE_RED, RED),
    (9.13, 'station adapter', '训练', PALE_GREEN, GREEN),
    (10.94, 'TEAM 上层', '训练', PALE_BLUE, BLUE),
]:
    card = add_shape_box(slide, Inches(x), Inches(2.1), Inches(1.46), Inches(1.18), fill=fill, line=LINE)
    write_in_shape(card, [title, sub], sizes=[13, 14], colors=[TEXT, accent], bolds=[True, True], aligns=[PP_ALIGN.CENTER, PP_ALIGN.CENTER], margins=(4, 4, 12, 0))
add_arrow(slide, Inches(8.78), Inches(2.69), Inches(9.13), Inches(2.69), width=1.8)
add_arrow(slide, Inches(10.59), Inches(2.69), Inches(10.94), Inches(2.69), width=1.8)

add_chip(slide, Inches(7.28), Inches(3.65), Inches(1.8), Inches(0.3), '当前实验目的', PALE_ORANGE, ORANGE)
add_bullets(
    slide, Inches(7.3), Inches(4.05), Inches(5.0), Inches(2.2),
    [
        '先确认前向、反向与 loss 计算链路都正常。',
        '确认输入波形、metadata、P pick、PGA 标签对齐没有明显错误。',
        '观察小样本下 loss 是否具备显著下降趋势。',
    ],
    font_size=13.5
)
add_footer(slide, 4, '代码参考：build_diting_args / build_optimizer_with_groups / maybe_dump_model_batch')

# Slide 5
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_header(slide, '04', '训练数据概览：日本强震数据')

mag_rng = notes['full_magnitude_range']
add_metric_card(slide, Inches(0.72), Inches(1.25), Inches(2.0), Inches(0.98), notes['full_available_events'], '全量可用事件数', BLUE)
add_metric_card(slide, Inches(2.92), Inches(1.25), Inches(2.2), Inches(0.98), f"{notes['full_available_station_rows']:,}", '全量台站记录数', GREEN)
add_metric_card(slide, Inches(5.32), Inches(1.25), Inches(1.9), Inches(0.98), f"M {mag_rng[0]}–{mag_rng[1]}", '震级范围', ORANGE)
add_metric_card(slide, Inches(7.42), Inches(1.25), Inches(2.35), Inches(0.98), f"KiK-net {notes['full_network_counts']['kik']}", '网络 1', TEAL)
add_metric_card(slide, Inches(9.97), Inches(1.25), Inches(2.35), Inches(0.98), f"K-NET {notes['full_network_counts']['knt']}", '网络 2', SKY)

add_picture_framed(slide, IMG_FULL_SCALE, Inches(0.72), Inches(2.55), Inches(5.85), Inches(3.35), caption='全量数据规模与网络/传感器组成（已按 PPT 场景重绘大字体版本）')
add_picture_framed(slide, IMG_MAG, Inches(6.77), Inches(2.55), Inches(5.55), Inches(3.35), caption='震级分布与按 split 的震级覆盖情况（已按 PPT 场景重绘大字体版本）')

note_box = add_shape_box(slide, Inches(0.9), Inches(6.35), Inches(11.45), Inches(0.42), fill=WHITE, line=LINE)
write_in_shape(note_box, ['当前数据具备一定规模与覆盖范围，适合作为后续完整训练的基础；本阶段先用其中的 small-sample 子集做调试实验。'], sizes=[11.5], colors=[NAVY], bolds=[True], aligns=[PP_ALIGN.CENTER], margins=(10, 10, 9, 0))
add_footer(slide, 5, '图像已根据 generate_data_overview_ppt_assets.ipynb 重绘为大字体 PPT 版本')

# Slide 6
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_header(slide, '05', 'small-sample / overfit-style 调试实验设置')

add_picture_framed(slide, IMG_GEO, Inches(0.72), Inches(1.42), Inches(5.7), Inches(3.7), caption='事件空间覆盖与台站空间分布（已按 PPT 场景重绘大字体版本）')
add_picture_framed(slide, IMG_SPLIT, Inches(6.65), Inches(1.42), Inches(5.68), Inches(3.7), caption='split 与当前 small-sample 事件分布（已按 PPT 场景重绘大字体版本）')

warn = add_shape_box(slide, Inches(0.95), Inches(6.28), Inches(11.15), Inches(0.32), fill=PALE_ORANGE, line=PALE_ORANGE)
write_in_shape(warn, ['更准确地说，这是一组 small-sample / overfit-style 调试实验，而不是严格意义上 train==dev 的经典 overfit 设置。'], sizes=[10.5], colors=[ORANGE], bolds=[True], aligns=[PP_ALIGN.CENTER], margins=(4, 4, 7, 0))
add_footer(slide, 6, '图像已根据 generate_data_overview_ppt_assets.ipynb 重绘为大字体 PPT 版本')

# Slide 7
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_header(slide, '06', '当前训练结果与已做排查')

status = add_shape_box(slide, Inches(0.72), Inches(1.35), Inches(4.75), Inches(5.6), fill=WHITE, line=LINE)
add_chip(slide, Inches(0.98), Inches(1.62), Inches(1.2), Inches(0.3), '当前状态', PALE_RED, RED)
add_bullets(
    slide, Inches(0.98), Inches(2.02), Inches(4.2), Inches(2.0),
    [
        'PGA 训练已经能正常启动并运行。',
        '当前 small-sample 实验中，loss 尚未体现出理想的收敛趋势。',
        '因此现阶段重点不是性能汇报，而是定位训练瓶颈。',
    ],
    font_size=14.1
)
add_chip(slide, Inches(0.98), Inches(4.38), Inches(1.2), Inches(0.3), '当前判断', PALE_BLUE, BLUE)
add_bullets(
    slide, Inches(0.98), Inches(4.78), Inches(4.15), Inches(1.4),
    [
        '暂未发现明显的数据读取维度错误。',
        '更可能的问题集中在标签定义、尺度变换、损失设计、冻结策略或学习率设置。',
    ],
    font_size=13.2
)

p1 = add_shape_box(slide, Inches(5.75), Inches(1.35), Inches(3.1), Inches(2.25), fill=PALE_GRAY, line=LINE)
write_in_shape(p1, ['Loss 曲线', '稍后补充'], sizes=[22, 17], colors=[MUTED, MUTED], bolds=[True, True], aligns=[PP_ALIGN.CENTER, PP_ALIGN.CENTER], margins=(10, 10, 28, 0))
p2 = add_shape_box(slide, Inches(9.15), Inches(1.35), Inches(3.15), Inches(2.25), fill=PALE_GRAY, line=LINE)
write_in_shape(p2, ['测试集推理样例', '稍后补充'], sizes=[20, 17], colors=[MUTED, MUTED], bolds=[True, True], aligns=[PP_ALIGN.CENTER, PP_ALIGN.CENTER], margins=(10, 10, 28, 0))

add_chip(slide, Inches(5.98), Inches(3.95), Inches(1.55), Inches(0.3), '已做输入检查', PALE_GREEN, GREEN)
check_data = [
    ('Waveform', '检查 batch 内波形形状、振幅范围与样例可视化。', PALE_BLUE, BLUE),
    ('Metadata', '检查台站坐标、有效台站 mask 与 selected station 索引。', PALE_TEAL, TEAL),
    ('P pick', '检查 raw / shifted pick 与 cutout 后对齐关系。', PALE_GREEN, GREEN),
    ('PGA label', '检查输入 PGA 值、分布范围以及异常值风险。', PALE_ORANGE, ORANGE),
]
coords = [(5.82, 4.35), (9.1, 4.35), (5.82, 5.55), (9.1, 5.55)]
for (title, desc, fill, accent), (x, y) in zip(check_data, coords):
    b = add_shape_box(slide, Inches(x), Inches(y), Inches(2.95), Inches(0.95), fill=fill, line=LINE)
    write_in_shape(b, [title, desc], sizes=[12.5, 10.1], colors=[accent, TEXT], bolds=[True, False])
add_footer(slide, 7, '输入检查参考：analyze_model_input_dumps.ipynb；训练时可导出 model_input_dumps 做 batch 级排查')

# Slide 8
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_header(slide, '07', '下一步工作')

plans = [
    ('数据与标签链路', [
        '继续核查 PGA target 的定义、尺度与变换。',
        '检查事件-台站配对、有效台站筛选与 station selection 逻辑。',
        '确认当前 small-sample 设置是否满足“应当能拟合”的条件。',
    ], PALE_BLUE, BLUE),
    ('优化与模型策略', [
        '系统比较冻结策略、学习率组和 station adapter 初始化方式。',
        '必要时先用更简单 baseline 做 sanity check。',
        '评估 single-station pretrain / station pretrain 对收敛性的帮助。',
    ], PALE_GREEN, GREEN),
    ('结果补充与汇报', [
        '补充训练 loss 曲线与关键中间指标。',
        '补充测试集推理样例与误差分析。',
        '形成 TEAM 原版 vs 当前改造版 的阶段性对比结论。',
    ], PALE_ORANGE, ORANGE),
]
for i, (title, bullets, fill, accent) in enumerate(plans):
    x = Inches(0.78 + i * 4.15)
    box = add_shape_box(slide, x, Inches(1.55), Inches(3.7), Inches(4.85), fill=fill, line=LINE)
    write_in_shape(box, [title], sizes=[17], colors=[accent], bolds=[True], aligns=[PP_ALIGN.CENTER], margins=(8, 8, 14, 0))
    add_bullets(slide, x + Inches(0.18), Inches(2.12), Inches(3.35), Inches(3.85), bullets, font_size=12.6)

close = add_shape_box(slide, Inches(0.95), Inches(6.48), Inches(11.15), Inches(0.38), fill=WHITE, line=LINE)
write_in_shape(close, ['结论：当前已完成 “架构接入 + 数据准备 + 初步排查”，下一阶段聚焦收敛瓶颈定位与可解释结果补充。'], sizes=[12.4], colors=[NAVY], bolds=[True], aligns=[PP_ALIGN.CENTER], margins=(8, 8, 8, 0))
add_footer(slide, 8, '后续可在此基础上继续补充 loss 曲线、推理样例和对比实验结果')

prs.save(str(OUT))
print(f'Wrote {OUT}')
