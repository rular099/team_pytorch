#!/usr/bin/env python3
"""Build a draft academic PPT for the PGA project from generated report assets."""

from __future__ import annotations

import argparse
import math
import os
import re
from pathlib import Path

os.environ.setdefault("MPLCONFIGDIR", "/tmp/matplotlib-cache")

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


FONT = "Microsoft YaHei"
TITLE = RGBColor(28, 47, 74)
TEXT = RGBColor(35, 40, 48)
MUTED = RGBColor(105, 115, 128)
BLUE = RGBColor(52, 101, 164)
GREEN = RGBColor(70, 142, 92)
RED = RGBColor(190, 58, 52)
LIGHT_BLUE = RGBColor(227, 237, 248)
LIGHT_GRAY = RGBColor(244, 246, 248)
WHITE = RGBColor(255, 255, 255)


def emu(value):
    return value


def add_textbox(slide, left, top, width, height, text, font_size=20, bold=False,
                color=TEXT, align=PP_ALIGN.LEFT, line_spacing=1.05):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.55), Inches(0.25), Inches(12.25), Inches(0.55),
                title, 28, True, TITLE)
    if subtitle:
        add_textbox(slide, Inches(0.58), Inches(0.78), Inches(12.0), Inches(0.3),
                    subtitle, 12, False, MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.03), Inches(12.25), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_BLUE
    line.line.fill.background()


def add_footer(slide, text="Draft figures use current eval; replace with epidist-selection eval when available."):
    add_textbox(slide, Inches(0.55), Inches(7.12), Inches(12.25), Inches(0.22),
                text, 8, False, MUTED, PP_ALIGN.RIGHT)


def add_bullets(slide, items, left, top, width, height, font_size=18, bullet=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.level = 0
        if bullet:
            p.text = f"• {item}"
    return box


def add_table(slide, df, left, top, width, height, font_size=14):
    rows, cols = df.shape[0] + 1, df.shape[1]
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_BLUE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = FONT
                r.font.size = Pt(font_size)
                r.font.bold = True
                r.font.color.rgb = TEXT
    for i in range(df.shape[0]):
        for j in range(cols):
            value = df.iloc[i, j]
            if isinstance(value, float):
                text = f"{value:.3f}"
            else:
                text = str(value)
            cell = table.cell(i + 1, j)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_GRAY
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.name = FONT
                    r.font.size = Pt(font_size)
                    r.font.color.rgb = TEXT
    return table


def fit_image(slide, image_path: Path, left, top, width, height):
    if not image_path.exists():
        add_textbox(slide, left, top, width, height, f"Missing image:\n{image_path.name}", 18, True, RED, PP_ALIGN.CENTER)
        return None
    return slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def add_box(slide, left, top, width, height, text, fill, line=BLUE, font_size=16):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = TEXT
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=BLUE):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(2)
    line.line.end_arrowhead = True
    return line


def metric_summary(y, p):
    y = np.asarray(y, dtype=float).reshape(-1)
    p = np.asarray(p, dtype=float).reshape(-1)
    mask = np.isfinite(y) & np.isfinite(p)
    y, p = y[mask], p[mask]
    residual = p - y
    mae = float(np.mean(np.abs(residual)))
    rmse = float(np.sqrt(np.mean(residual ** 2)))
    corr = float(np.corrcoef(y, p)[0, 1]) if len(y) > 1 else np.nan
    denom = max(float(np.sum((y - y.mean()) ** 2)), 1e-12)
    r2 = float(1.0 - np.sum(residual ** 2) / denom)
    return {"MAE": mae, "RMSE": rmse, "Corr": corr, "R2": r2, "N": len(y)}


def npz_array(data, key):
    return np.asarray(data[key], dtype=float)


def build_single_vs_multi_asset(npz_path: Path, metrics_csv: Path, out_dir: Path) -> tuple[Path, pd.DataFrame]:
    data = np.load(npz_path, allow_pickle=True)
    metrics = pd.read_csv(metrics_csv)
    rows = []
    for split in ("train", "val"):
        full_row = metrics[(metrics["checkpoint"] == "best") & (metrics["split"] == split)].iloc[0]
        rows.append({
            "Model": "Multi-station",
            "Split": split,
            "MAE": float(full_row["mae"]),
            "RMSE": float(full_row["rmse"]),
            "Corr": float(full_row["corr"]),
            "R2": float(full_row["r2"]),
            "N": int(np.asarray(data[f"{split}_pga_target_valid"]).astype(bool).sum()),
        })
        single = metric_summary(npz_array(data, f"single_{split}_pga_label"), npz_array(data, f"single_{split}_pga_mu"))
        rows.append({
            "Model": "Single-station",
            "Split": split,
            **single,
        })
    df = pd.DataFrame(rows)
    out_dir.mkdir(parents=True, exist_ok=True)
    out = out_dir / "single_vs_multi_pga_metrics.png"
    fig, axes = plt.subplots(1, 2, figsize=(13.5, 5.4))
    colors = {"Multi-station": "#377eb8", "Single-station": "#4daf4a"}
    for ax, metric, title in zip(axes, ["MAE", "Corr"], ["PGA Error", "PGA Correlation"]):
        x = np.arange(2)
        width = 0.34
        for offset, model in [(-width / 2, "Multi-station"), (width / 2, "Single-station")]:
            vals = [
                float(df[(df.Model == model) & (df.Split == split)][metric].iloc[0])
                for split in ("train", "val")
            ]
            ax.bar(x + offset, vals, width, label=model, color=colors[model])
        ax.set_xticks(x)
        ax.set_xticklabels(["Train", "Validation"], fontsize=15)
        ax.set_title(title, fontsize=20, weight="bold")
        ax.set_ylabel(metric, fontsize=17)
        ax.tick_params(axis="y", labelsize=14)
        ax.grid(axis="y", alpha=0.25)
        ax.legend(fontsize=13)
    fig.tight_layout()
    fig.savefig(out, dpi=220)
    plt.close(fig)
    return out, df


def first_existing(asset_dir: Path, names: list[str]) -> Path:
    for name in names:
        p = asset_dir / name
        if p.exists():
            return p
    return asset_dir / names[0]


def choose_diagnostic_images(diagnostics_dir: Path, limit: int = 4) -> list[Path]:
    if not diagnostics_dir.exists():
        return []
    waveform = sorted(diagnostics_dir.glob("waveform_pick_check_*.png"))
    if waveform:
        rng = np.random.default_rng(2026)
        pick = rng.choice(len(waveform), size=min(limit, len(waveform)), replace=False)
        return [waveform[int(i)] for i in pick]
    candidates = sorted(diagnostics_dir.rglob("*.png"))
    return candidates[:limit]


def read_data_summary(asset_dir: Path) -> pd.DataFrame:
    path = asset_dir / "data_split_summary.csv"
    if path.exists():
        return pd.read_csv(path)
    return pd.DataFrame()


def add_data_stat_cards(slide, summary: pd.DataFrame, left, top, width, height):
    if summary.empty:
        return
    total_events = int(summary["Events"].sum()) if "Events" in summary else 0
    total_records = int(summary["Station records"].sum()) if "Station records" in summary else 0
    train = summary[summary["Split"].astype(str).str.lower() == "train"]
    val = summary[summary["Split"].astype(str).str.lower().isin(["validation", "val", "dev"])]
    test = summary[summary["Split"].astype(str).str.lower() == "test"]
    lines = [
        f"Events: {total_events:,}",
        f"Station records: {total_records:,}",
    ]
    if not train.empty:
        lines.append(f"Train: {int(train['Events'].iloc[0]):,} events")
    if not val.empty:
        lines.append(f"Validation: {int(val['Events'].iloc[0]):,} events")
    if not test.empty:
        lines.append(f"Test: {int(test['Events'].iloc[0]):,} events")
    if "Magnitude max" in summary:
        lines.append(f"Max M: {summary['Magnitude max'].max():.1f}")
    if "PGA max m/s2" in summary:
        lines.append(f"Max PGA: {summary['PGA max m/s2'].max():.1f} m/s2")
    add_bullets(slide, lines, left, top, width, height, 15)


def sorted_numbered_images(asset_dir: Path, prefix: str) -> list[Path]:
    images = list(asset_dir.glob(f"{prefix}_*.png"))
    def key(path: Path):
        m = re.search(rf"{re.escape(prefix)}_(\d+)_", path.name)
        return int(m.group(1)) if m else 999
    return sorted(images, key=key)


def build_ppt(
    asset_dir: Path,
    report_input_dir: Path,
    output: Path,
    diagnostics_dir: Path | None = None,
    team_diagram: Path | None = None,
    quakeformer_diagram: Path | None = None,
):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    primary = "japan_overfit_pga15_cross_attention_overfitall_fixed_inputs_random_targets"
    metrics_csv = asset_dir / "main_pga_metrics.csv"
    npz_path = report_input_dir / "eval_results_best.npz"
    single_chart, single_df = build_single_vs_multi_asset(npz_path, metrics_csv, asset_dir)
    data_summary = read_data_summary(asset_dir)
    diagnostics_dir = diagnostics_dir or Path(os.environ.get("PGA_DIAGNOSTICS_DIR", "/opt/zb/data/japan/diagnostics_2024"))
    diagnostic_images = choose_diagnostic_images(diagnostics_dir, limit=10)
    team_diagram = team_diagram or Path("../team_diagram.png")
    quakeformer_diagram = quakeformer_diagram or Path("../quakeformer_diagram.png")

    # 1
    slide = prs.slides.add_slide(blank)
    bg = first_existing(asset_dir, [f"case_station_maps_1_{primary}_best_val.png"])
    fit_image(slide, bg, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(255, 255, 255)
    overlay.fill.transparency = 22
    overlay.line.fill.background()
    add_textbox(slide, Inches(0.7), Inches(0.85), Inches(8.4), Inches(1.0),
                "基于多台站波形表征的目标台站 PGA 估计", 30, True, TITLE)
    add_textbox(slide, Inches(0.75), Inches(2.0), Inches(7.0), Inches(0.4),
                "日本强震数据上的 log PGA 预测与空间诊断", 18, False, TEXT)
    add_textbox(slide, Inches(0.75), Inches(6.55), Inches(7.0), Inches(0.35),
                "Draft version | current figures use existing eval outputs", 12, False, MUTED)

    # 2
    slide = prs.slides.add_slide(blank)
    add_title(slide, "报告主线")
    add_bullets(slide, [
        "PGA 是地震动快速评估和工程应用中的核心指标",
        "目标：用有限输入台站波形，预测多个目标台站的 log PGA",
        "方法：单台波形表征 + 多台站 Transformer + target cross-attention",
        "结果：从总体精度、动态范围、输入台站数量、震中距和空间案例五个角度评估",
    ], Inches(0.85), Inches(1.45), Inches(11.7), Inches(4.4), 22)
    add_footer(slide, "30-minute talk structure: background, method, result, diagnosis, case studies.")

    # 3
    slide = prs.slides.add_slide(blank)
    add_title(slide, "研究问题：目标台站 PGA 估计")
    add_box(slide, Inches(0.8), Inches(1.6), Inches(2.7), Inches(1.0), "输入台站波形\n3C waveform", LIGHT_BLUE)
    add_box(slide, Inches(0.8), Inches(3.2), Inches(2.7), Inches(0.8), "输入台站坐标\nstation coords", LIGHT_BLUE)
    add_arrow(slide, Inches(3.55), Inches(2.1), Inches(4.75), Inches(2.35))
    add_arrow(slide, Inches(3.55), Inches(3.6), Inches(4.75), Inches(2.65))
    add_box(slide, Inches(4.85), Inches(1.85), Inches(3.1), Inches(1.3), "多台站模型\nvariable station set", RGBColor(238, 245, 232), GREEN)
    add_arrow(slide, Inches(8.05), Inches(2.35), Inches(9.1), Inches(2.35))
    add_box(slide, Inches(9.2), Inches(1.55), Inches(3.1), Inches(1.6), "目标台站 log PGA\nsite-specific output", RGBColor(249, 232, 230), RED)
    add_bullets(slide, [
        "输入台站数量可变，且空间分布每个事件不同",
        "目标台站可以不是输入台站",
        "报告中 PGA 结果均在 log PGA 空间评估",
    ], Inches(0.95), Inches(4.75), Inches(11.5), Inches(1.4), 19)

    # 4
    slide = prs.slides.add_slide(blank)
    add_title(slide, "为什么关注 PGA")
    add_bullets(slide, [
        "PGA = max |a(t)|，描述特定地点峰值地面加速度",
        "PGA 是烈度快速估计、工程地震动评价和应急决策中的基础指标",
        "PGA 受震级、震中距、传播路径、场地条件和局部放大共同影响",
        "本项目关注：不显式使用绝对振幅信息时，多台站模型能否恢复空间 PGA 分布",
    ], Inches(0.85), Inches(1.35), Inches(7.3), Inches(4.8), 20)
    refs = pd.DataFrame({
        "Factor": ["Source", "Path", "Site", "Observation"],
        "Examples": ["Magnitude, rupture", "Distance, attenuation", "VS30, basin", "Waveform shape"],
    })
    add_table(slide, refs, Inches(8.35), Inches(1.55), Inches(4.25), Inches(2.35), 14)

    # 5
    slide = prs.slides.add_slide(blank)
    add_title(slide, "传统 GMPE 方法：经验方程预测地震动")
    add_box(slide, Inches(0.75), Inches(1.45), Inches(1.8), Inches(0.75), "震源参数\nM, fault type", LIGHT_BLUE, BLUE, 13)
    add_box(slide, Inches(0.75), Inches(2.65), Inches(1.8), Inches(0.75), "传播距离\nRrup / RJB", LIGHT_BLUE, BLUE, 13)
    add_box(slide, Inches(0.75), Inches(3.85), Inches(1.8), Inches(0.75), "场地参数\nVS30 / basin", LIGHT_BLUE, BLUE, 13)
    add_arrow(slide, Inches(2.65), Inches(1.82), Inches(4.05), Inches(2.45))
    add_arrow(slide, Inches(2.65), Inches(3.02), Inches(4.05), Inches(2.88))
    add_arrow(slide, Inches(2.65), Inches(4.22), Inches(4.05), Inches(3.3))
    add_box(slide, Inches(4.15), Inches(2.1), Inches(3.0), Inches(1.55),
            "经验回归方程\nln PGA = f(M, R, site, fault) + ε", RGBColor(238, 245, 232), GREEN, 14)
    add_arrow(slide, Inches(7.25), Inches(2.95), Inches(8.55), Inches(2.95))
    add_box(slide, Inches(8.65), Inches(2.1), Inches(2.2), Inches(1.55),
            "中位数 PGA\n+ 不确定性 σ", RGBColor(249, 232, 230), RED, 14)
    add_bullets(slide, [
        "优势：物理量含义清楚，计算快，工程应用成熟",
        "通常需要先由定位/震级估计模块得到震源参数，再结合距离和场地项计算地震动",
        "局限 1：输入变量低维，难以直接利用实时多台站波形形状",
        "局限 2：区域、路径和局部场地效应常被压缩进经验项和随机残差",
        "局限 3：对复杂空间残差和事件内台站间相关性表达有限",
        "本项目补充方向：让模型从多台站波形和几何关系中直接学习目标台站 PGA",
    ], Inches(0.95), Inches(5.0), Inches(11.5), Inches(1.45), 14)

    # 6
    slide = prs.slides.add_slide(blank)
    add_title(slide, "文献：TEAM")
    if team_diagram.exists():
        fit_image(slide, team_diagram, Inches(0.55), Inches(1.15), Inches(7.4), Inches(5.85))
        add_bullets(slide, [
            "图中核心是任意数量台站输入：每个台站波形与坐标形成 station token",
            "Transformer 在台站集合上建模，再面向目标位置输出地震动预测",
            "TEAM 避免固定台站输入限制，适合事件附近可变台站组合",
            "与 QuakeFormer 相比，TEAM 更偏早期预警任务；本项目参考其可变台站集合建模思想",
        ], Inches(8.15), Inches(1.3), Inches(4.55), Inches(5.25), 15)
    else:
        add_bullets(slide, [
            "Münchmeyer et al. (2021) 指出传统早期预警可分为 source-based 与 propagation-based；前者常要先估计震源再接 GMPE",
            "TEAM 直接分析任意数量、任意位置台站的强震波形，目标是同时兼顾准确性、及时性和可变台网适应性",
            "论文强调固定台站深度模型难以适应大规模密集台网，因为离事件最近的台站会随事件变化",
            "与 QuakeFormer 相比，TEAM 更偏早期预警任务；本项目参考其任意台站集合建模思想",
        ], Inches(0.85), Inches(1.3), Inches(11.5), Inches(4.7), 19)

    # 7
    slide = prs.slides.add_slide(blank)
    add_title(slide, "文献：QuakeFormer")
    if quakeformer_diagram.exists():
        fit_image(slide, quakeformer_diagram, Inches(0.55), Inches(1.15), Inches(7.4), Inches(5.85))
        add_bullets(slide, [
            "图中 masked Transformer 同时服务 forecasting、early warning 和 interpolation",
            "mask 表示把部分台站/位置的地震动观测隐藏，让模型用可见台站去预测这些被隐藏位置",
            "因此被 mask 的位置在当前任务中就是 target；区别只在于哪些台站可见、哪些位置待预测",
            "与 TEAM 相比，QuakeFormer 更强调 masked 多任务框架和显式绝对/相对空间坐标嵌入",
            "本项目当前未加入相对几何 bias，后续可借鉴其空间建模设计",
        ], Inches(8.15), Inches(1.25), Inches(4.55), Inches(5.45), 14)
    else:
        add_bullets(slide, [
            "Feng et al. (2024) 将 forecasting、early warning 和 interpolation 统一到 masked Transformer 框架",
            "mask 表示把部分台站/位置的地震动观测隐藏，让模型用可见台站去预测这些被隐藏位置",
            "因此被 mask 的位置在当前任务中就是 target；区别只在于哪些台站可见、哪些位置待预测",
            "与 TEAM 相比，QuakeFormer 更强调 masked 多任务框架和显式绝对/相对空间坐标嵌入",
            "本项目当前未加入相对几何 bias，后续可借鉴其空间建模设计",
        ], Inches(0.85), Inches(1.3), Inches(11.5), Inches(4.7), 19)

    # 8
    slide = prs.slides.add_slide(blank)
    add_title(slide, "方案、数据与评估口径")
    data = np.load(npz_path, allow_pickle=True)
    data_rows = []
    for split in ("train", "val"):
        data_rows.append({
            "Split": "Train" if split == "train" else "Validation",
            "Events": len(data[f"{split}_event_index"]),
            "PGA targets": int(np.asarray(data[f"{split}_pga_target_valid"]).astype(bool).sum()),
            "Single-station samples": len(data[f"single_{split}_pga_label"]),
        })
    add_table(slide, pd.DataFrame(data_rows), Inches(0.75), Inches(1.25), Inches(5.75), Inches(1.75), 13)
    add_bullets(slide, [
        "可变数量输入台站波形经过单台表征后形成 station tokens",
        "PGA 采用 target-wise cross-attention readout，每个目标台站独立查询输入台站 tokens",
        "当前模型使用台站/目标位置编码，但未加入显式相对几何 bias",
        "主结果使用按震中距最近台站选择的 epidist eval",
        "评估指标：MAE、RMSE、Corr、R2、linear-fit slope、Bias",
        "single-station 结果用于说明单台波形表征能力；multi-station 结果用于目标台站 PGA 估计",
    ], Inches(6.85), Inches(1.22), Inches(5.85), Inches(4.75), 15)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "训练数据分布：事件与台站记录")
    add_data_stat_cards(slide, data_summary, Inches(0.65), Inches(1.25), Inches(2.35), Inches(4.8))
    fit_image(slide, asset_dir / "data_distribution_overview.png", Inches(3.15), Inches(1.08), Inches(9.7), Inches(5.95))
    add_footer(slide, "Split CSVs: split_events.csv and split_stations.csv.")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "训练数据分布：统计汇总")
    fit_image(slide, asset_dir / "data_split_summary.png", Inches(0.65), Inches(1.15), Inches(12.0), Inches(3.15))
    fit_image(slide, asset_dir / "data_event_station_depth.png", Inches(0.75), Inches(4.25), Inches(11.85), Inches(2.55))
    add_footer(slide, "PGA distribution uses pga_norm_resampled_mps2 when available.")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "训练数据分布：事件与台站空间覆盖")
    fit_image(slide, asset_dir / "data_event_station_map.png", Inches(0.55), Inches(1.05), Inches(12.25), Inches(5.95))
    add_footer(slide, "Events colored by magnitude; station locations are deduplicated from station records.")

    # 13
    slide = prs.slides.add_slide(blank)
    add_title(slide, "模型结构")
    labels = [
        ("输入台站波形", 0.75), ("DiTing encoder", 2.55), ("Station adapter", 4.35),
        ("Station tokens\n+ coords", 6.15), ("Multistation\nTransformer", 7.95),
        ("PGA target query\n+ location encoding", 9.75), ("log PGA", 11.6)
    ]
    for text, x in labels:
        add_box(slide, Inches(x), Inches(2.05), Inches(1.45), Inches(1.05), text, LIGHT_BLUE if x < 6 else RGBColor(238, 245, 232), font_size=13)
    for i in range(len(labels) - 1):
        add_arrow(slide, Inches(labels[i][1] + 1.45), Inches(2.58), Inches(labels[i + 1][1]), Inches(2.58))
    add_bullets(slide, [
        "station tokens 先经过多台站 Transformer 混合",
        "PGA target 不直接作为普通 token 参与 self-attention，而是 cross-attend station tokens",
        "target query 包含目标位置编码和可学习 query token",
    ], Inches(0.95), Inches(4.25), Inches(11.5), Inches(1.65), 18)

    # 14
    slide = prs.slides.add_slide(blank)
    add_title(slide, "训练策略")
    add_bullets(slide, [
        "Single-station pretrain: mag / epidist / PGA 多任务，强化 PGA 权重",
        "Full model: target-station PGA objective，使用 Huber loss",
        "PGA target normalization: 使用训练集统计量，在 eval 时反归一化后计算指标",
        "不显式使用绝对振幅信息：当前关注波形形状表征与台站位置的贡献",
        "Checkpoint: 同时评估 model_best 和 model_last，主报告优先展示 best",
    ], Inches(0.9), Inches(1.25), Inches(11.4), Inches(4.9), 20)

    # 15
    slide = prs.slides.add_slide(blank)
    add_title(slide, "主结果：Train 与 Validation")
    fit_image(slide, asset_dir / "train_val_metric_bars.png", Inches(0.65), Inches(1.18), Inches(12.0), Inches(5.65))
    add_footer(slide)

    # 16
    slide = prs.slides.add_slide(blank)
    add_title(slide, "单台模型与多台模型对比")
    fit_image(slide, single_chart, Inches(0.65), Inches(1.15), Inches(7.0), Inches(3.1))
    summary = single_df[["Model", "Split", "MAE", "RMSE", "Corr", "R2", "N"]].copy()
    summary["Split"] = summary["Split"].map({"train": "Train", "val": "Val"})
    add_table(slide, summary, Inches(0.85), Inches(4.55), Inches(11.6), Inches(1.8), 11)
    add_bullets(slide, [
        "single-station 预测的是输入台站 PGA，multi-station 预测目标台站 PGA；二者样本集合不同，比较用于展示表征层面的参考基线",
        "正式版本建议用同一 epidist eval 重新计算该页",
    ], Inches(7.95), Inches(1.25), Inches(4.65), Inches(2.7), 14)

    # 17
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Predicted vs True PGA")
    fit_image(slide, asset_dir / f"scatter_{primary}_best_train.png", Inches(0.65), Inches(1.22), Inches(5.95), Inches(5.45))
    fit_image(slide, asset_dir / f"scatter_{primary}_best_val.png", Inches(6.75), Inches(1.22), Inches(5.95), Inches(5.45))
    add_footer(slide)

    # 18
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Residual 诊断：强弱 PGA 是否系统偏差")
    fit_image(slide, asset_dir / f"residual_vs_true_{primary}_best_val.png", Inches(0.7), Inches(1.18), Inches(6.1), Inches(5.55))
    add_bullets(slide, [
        "Residual = predicted log PGA - true log PGA",
        "MAE 表示误差绝对大小",
        "Bias 表示系统性高估或低估",
        "强 PGA 若长期负 Bias，说明动态范围仍被压缩",
    ], Inches(7.25), Inches(1.55), Inches(5.0), Inches(3.2), 18)
    add_footer(slide)

    # 19
    slide = prs.slides.add_slide(blank)
    add_title(slide, "按 True PGA 强度分桶")
    fit_image(slide, asset_dir / f"pga_strength_bins_{primary}_best_val.png", Inches(0.75), Inches(1.15), Inches(11.85), Inches(5.75))
    add_footer(slide)

    # 20
    slide = prs.slides.add_slide(blank)
    add_title(slide, "输入台站数量对误差的影响")
    fit_image(slide, asset_dir / "val_mae_by_station_count.png", Inches(1.15), Inches(1.12), Inches(11.0), Inches(5.85))
    add_footer(slide)

    # 21
    slide = prs.slides.add_slide(blank)
    add_title(slide, "震中距对误差的影响")
    fit_image(slide, asset_dir / "val_mae_by_epicentral_distance.png", Inches(1.15), Inches(1.12), Inches(11.0), Inches(5.85))
    add_footer(slide)

    # 22-27: each case is shown as a paired station-count map and target-wise diagnostic.
    for rank in range(1, 4):
        slide = prs.slides.add_slide(blank)
        add_title(slide, f"Case Study {rank}：输入台站数量变化下的空间残差")
        fit_image(slide, asset_dir / f"case_station_maps_{rank}_{primary}_best_val.png", Inches(0.55), Inches(1.05), Inches(12.35), Inches(5.95))
        add_footer(slide)
        slide = prs.slides.add_slide(blank)
        add_title(slide, f"Case Study {rank}：目标台站级 PGA 与空间残差")
        fit_image(slide, asset_dir / f"case_study_{rank}_{primary}_best_val.png", Inches(0.55), Inches(1.12), Inches(12.3), Inches(5.85))
        add_footer(slide)

    # 28
    attention_all = sorted_numbered_images(asset_dir, "attention_target")
    attention_all = [path for path in attention_all if not path.name.startswith("attention_target_1_")]
    if attention_all:
        rng = np.random.default_rng(2026)
        selected_idx = rng.choice(len(attention_all), size=min(4, len(attention_all)), replace=False)
        attention_examples = [attention_all[int(i)] for i in selected_idx]
    else:
        attention_examples = []

    if attention_examples:
        for idx, attention_image in enumerate(attention_examples, start=1):
            slide = prs.slides.add_slide(blank)
            add_title(slide, "Attention：PGA target 关注哪些输入台站")
            fit_image(slide, attention_image, Inches(0.55), Inches(1.06), Inches(7.35), Inches(5.95))
            if idx == 1:
                add_bullets(slide, [
                    "红圈是当前正在计算 PGA 的目标台站",
                    "三角形表示输入台站，大小/颜色表示该目标台站 query 对输入台站的 attention 权重",
                    "红线连接 attention 最高的输入台站，用来直观看模型依赖哪些观测",
                    f"当前 attention npz 去掉第 1 个后剩余 {len(attention_all)} 个事件示例；本稿随机展示其中 {len(attention_examples)} 个",
                ], Inches(8.15), Inches(1.35), Inches(4.45), Inches(4.8), 15)
            else:
                add_bullets(slide, [
                    "同一解释口径：一个目标台站 query，一组输入台站 key/value",
                    "目标台站在不同事件中覆盖不同震中距，用于观察 attention 是否更偏近场或局部台站",
                    "当前导出的 attention 结果只包含 requested station count = 3",
                ], Inches(8.15), Inches(1.35), Inches(4.45), Inches(4.8), 15)
            add_footer(slide, "Attention uses model_last because eval_attention_best.npz was not produced.")
    else:
        slide = prs.slides.add_slide(blank)
        add_title(slide, "Attention 可视化：PGA target 关注哪些输入台站")
        add_bullets(slide, [
            "attention target 图未找到，保留可视化设计：目标台站 residual map + 输入台站 marker 大小/颜色表示 attention",
            "需要 eval_attention.py 输出 eval_attention_last.npz 或 eval_attention_best.npz",
            "当前报告主结果不依赖 attention；attention 仅作为解释性补充",
        ], Inches(0.85), Inches(1.3), Inches(11.4), Inches(4.6), 19)

    slide = prs.slides.add_slide(blank)
    add_title(slide, "训练数据局限性：自动 P pick 标签")
    add_bullets(slide, [
        "当前数据中的 P pick 标签主要由自动算法拾取，不是逐条人工精修标签",
        "P pick 偏差会影响输入窗口对齐，进而影响单台波形表征和多台站时序一致性",
        "近场强震、低信噪比、截断记录和复杂震相更容易引入 pick 不确定性",
        "因此模型误差不完全等价于结构能力不足，也包含上游数据标签噪声",
    ], Inches(0.85), Inches(1.3), Inches(11.5), Inches(4.6), 20)
    add_footer(slide, "Diagnostics directory: /opt/zb/data/japan/diagnostics_2024 when available.")

    for idx, image_path in enumerate(diagnostic_images, start=1):
        slide = prs.slides.add_slide(blank)
        add_title(slide, "训练数据局限性：P pick 诊断")
        fit_image(slide, image_path, Inches(0.65), Inches(1.05), Inches(12.0), Inches(5.95))
        add_footer(slide, f"Randomly selected from {diagnostics_dir}.")

    # 29
    slide = prs.slides.add_slide(blank)
    add_title(slide, "讨论和结论")
    add_bullets(slide, [
        "多台站 Transformer 可以自然处理可变输入台站集合和目标台站 PGA 预测",
        "实验结果表明，不显式使用绝对振幅信息时，波形形状表征仍然可以支持 PGA 强度估计",
        "强 PGA 动态范围仍是关键问题：需要关注 slope、强 PGA bin 的 Bias",
        "文献对比目前只做方法定位；没有 TEAM/QuakeFormer 同数据集复现，不能做严格数值排名",
        "当前数据集 P pick 不够精准，需要更高质量的数据集来训练模型在小秒数窗口下的能力",
    ], Inches(0.9), Inches(1.25), Inches(11.4), Inches(4.9), 19)

    # 31
    slide = prs.slides.add_slide(blank)
    add_title(slide, "References")
    refs = [
        "Muenchmeyer et al. (2021), The Transformer Earthquake Alerting Model, GJI.",
        "Feng, Zhu & Lu (2024), QuakeFormer: A Uniform Approach to Earthquake Ground Motion Prediction Using Masked Transformers.",
        "Boore & Atkinson (2008), Ground-motion prediction equations for PGA, PGV and PSA, Earthquake Spectra.",
    ]
    add_bullets(slide, refs, Inches(0.85), Inches(1.25), Inches(11.6), Inches(4.8), 16)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    single_df.to_csv(output.with_name(output.stem + "_single_vs_multi_metrics.csv"), index=False)
    return output


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--asset-dir", default="reports/pga_academic_report_assets_current")
    parser.add_argument("--report-input-dir", default="reports/pga_report_inputs/weights_japan_overfit_pga15_cross_attention_overfitall_fixed_inputs_random_targets")
    parser.add_argument("--output", default="reports/pga_academic_report_draft_current.pptx")
    default_diag = "../diagnostics_2024" if Path("../diagnostics_2024").exists() else os.environ.get("PGA_DIAGNOSTICS_DIR", "/opt/zb/data/japan/diagnostics_2024")
    parser.add_argument("--diagnostics-dir", default=default_diag)
    parser.add_argument("--team-diagram", default="../team_diagram.png")
    parser.add_argument("--quakeformer-diagram", default="../quakeformer_diagram.png")
    args = parser.parse_args()
    out = build_ppt(
        Path(args.asset_dir),
        Path(args.report_input_dir),
        Path(args.output),
        Path(args.diagnostics_dir),
        Path(args.team_diagram),
        Path(args.quakeformer_diagram),
    )
    print(f"[INFO] wrote {out}")


if __name__ == "__main__":
    main()
