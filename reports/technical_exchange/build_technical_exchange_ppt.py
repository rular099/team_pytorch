from __future__ import annotations

import math
import re
from pathlib import Path

import h5py
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[3]
OUT_DIR = ROOT / "team_pytorch" / "reports" / "technical_exchange"
ASSET_DIR = OUT_DIR / "assets"
TABLE_DIR = OUT_DIR / "tables"
PPTX_PATH = OUT_DIR / "diting_team_graph_pga_technical_exchange.pptx"

TEAM_RES = ROOT / "chaosuan_res"
GRAPH_RES = ROOT / "team_pytorch2" / "chaosuan_res"
DATA32 = ROOT / "team_pytorch" / "data_example" / "32"
DATA128 = ROOT / "team_pytorch" / "data_example" / "128"
H5_PATH = ROOT / "team_pytorch" / "japan_overfit.hdf5"


BLUE = RGBColor(31, 78, 121)
TEAL = RGBColor(0, 128, 128)
ORANGE = RGBColor(198, 89, 17)
GRAY = RGBColor(89, 89, 89)
LIGHT_BLUE = RGBColor(221, 235, 247)
LIGHT_ORANGE = RGBColor(252, 228, 214)
LIGHT_GREEN = RGBColor(226, 239, 218)
WHITE = RGBColor(255, 255, 255)


def ensure_dirs() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    TABLE_DIR.mkdir(parents=True, exist_ok=True)


def savefig(path: Path) -> None:
    plt.savefig(path, dpi=220, bbox_inches="tight")
    plt.close()


def read_csv(path: Path) -> pd.DataFrame:
    return pd.read_csv(path)


def last_value(csv_path: Path) -> float | None:
    if not csv_path.exists():
        return None
    df = pd.read_csv(csv_path)
    if df.empty:
        return None
    return float(df["value"].iloc[-1])


def extract_eval_metrics(path: Path) -> dict:
    metrics: dict[str, dict[str, dict[str, float]]] = {}
    if not path.exists() or path.stat().st_size == 0:
        return metrics
    current_set = None
    current_task = None
    for raw in path.read_text(errors="ignore").splitlines():
        line = raw.strip()
        m = re.search(r"^(TRAIN|VAL) set:", line)
        if m:
            current_set = m.group(1).lower()
            metrics.setdefault(current_set, {})
            current_task = None
            continue
        if line.startswith("SINGLE-STATION TRAIN"):
            current_set = "single_train"
            metrics.setdefault(current_set, {})
            current_task = None
            continue
        if line.startswith("SINGLE-STATION VAL"):
            current_set = "single_val"
            metrics.setdefault(current_set, {})
            current_task = None
            continue
        m = re.match(r"---\s+(.+?)\s+---", line)
        if m and current_set:
            current_task = m.group(1).replace("/", "_").lower()
            metrics[current_set].setdefault(current_task, {})
            continue
        if not current_set or not current_task:
            continue
        m = re.search(r"MAE=([-+0-9.eE]+), RMSE=([-+0-9.eE]+)", line)
        if m:
            metrics[current_set][current_task]["mae"] = float(m.group(1))
            metrics[current_set][current_task]["rmse"] = float(m.group(2))
            continue
        m = re.search(r"Correlation:\s*([-+0-9.eE]+)", line)
        if m:
            metrics[current_set][current_task]["corr"] = float(m.group(1))
            continue
        m = re.search(r"R\^2:\s*([-+0-9.eE]+)", line)
        if m:
            metrics[current_set][current_task]["r2"] = float(m.group(1))
            continue
        m = re.search(r"Linear fit: pred =\s*([-+0-9.eE]+)\s*\* label \+\s*([-+0-9.eE]+)", line)
        if m:
            metrics[current_set][current_task]["slope"] = float(m.group(1))
            metrics[current_set][current_task]["intercept"] = float(m.group(2))
    return metrics


def metric(path: Path, split: str = "val", task: str = "pga", key: str = "corr") -> float | None:
    data = extract_eval_metrics(path)
    return data.get(split, {}).get(task, {}).get(key)


def build_result_tables() -> dict[str, pd.DataFrame]:
    rows_transformer = [
        {
            "method": "query_transformer",
            "path": TEAM_RES / "weights_japan_overfit_exp1_attn_diag" / "eval_results_attn_diag.txt",
            "note": "PGA query in full self-attention",
        },
        {
            "method": "mask_batch1",
            "path": TEAM_RES / "weights_japan_overfit_exp1_mask_batch1" / "eval_results_mask_batch1.txt",
            "note": "batch=1 mask sanity",
        },
        {
            "method": "query_no_transformer",
            "path": TEAM_RES / "weights_japan_overfit_exp1_query_no_transformer" / "eval_results_query_no_transformer.txt",
            "note": "target coord + learned query only",
        },
        {
            "method": "direct_station",
            "path": TEAM_RES / "weights_japan_overfit_exp1_direct_station" / "eval_results_direct_station.txt",
            "note": "station embedding direct readout",
        },
    ]
    transformer = []
    for row in rows_transformer:
        ev = extract_eval_metrics(row["path"])
        transformer.append(
            {
                "Method": row["method"],
                "Train MAE": ev.get("train", {}).get("pga", {}).get("mae"),
                "Train Corr": ev.get("train", {}).get("pga", {}).get("corr"),
                "Train R2": ev.get("train", {}).get("pga", {}).get("r2"),
                "Train slope": ev.get("train", {}).get("pga", {}).get("slope"),
                "Val MAE": ev.get("val", {}).get("pga", {}).get("mae"),
                "Val Corr": ev.get("val", {}).get("pga", {}).get("corr"),
                "Val R2": ev.get("val", {}).get("pga", {}).get("r2"),
                "Val slope": ev.get("val", {}).get("pga", {}).get("slope"),
                "Note": row["note"],
            }
        )
    transformer_df = pd.DataFrame(transformer)

    rows_cross = [
        ("pga15_cross_overfit32", TEAM_RES / "weights_japan_overfit_pga15_cross_attention_overfit32" / "eval_results.txt"),
        ("pga15_cross_overfit128", TEAM_RES / "weights_japan_overfit_pga15_cross_attention_overfit128" / "eval_results.txt"),
        ("fixed_inputs_targets", TEAM_RES / "weights_japan_overfit_pga15_cross_attention_overfit32_fixed_inputs_targets" / "eval_results.txt"),
        ("input_targets", TEAM_RES / "weights_japan_overfit_pga15_cross_attention_overfit32_input_targets" / "eval_results.txt"),
        ("event_pga_cross_first_inputs", TEAM_RES / "weights_japan_overfit_event_pga_cross_attention_first_inputs" / "eval_results.txt"),
    ]
    cross = []
    for name, path in rows_cross:
        ev = extract_eval_metrics(path)
        cross.append(
            {
                "Experiment": name,
                "Train MAE": ev.get("train", {}).get("pga", {}).get("mae"),
                "Train Corr": ev.get("train", {}).get("pga", {}).get("corr"),
                "Train R2": ev.get("train", {}).get("pga", {}).get("r2"),
                "Val MAE": ev.get("val", {}).get("pga", {}).get("mae"),
                "Val Corr": ev.get("val", {}).get("pga", {}).get("corr"),
                "Val R2": ev.get("val", {}).get("pga", {}).get("r2"),
                "Val slope": ev.get("val", {}).get("pga", {}).get("slope"),
                "Val mag MAE": ev.get("val", {}).get("mag", {}).get("mae"),
            }
        )
    cross_df = pd.DataFrame(cross)

    rows_graph = [
        ("old graph first-inputs", GRAPH_RES / "weights_japan_overfit_graph_pga_first_inputs" / "eval_results.txt"),
        ("prior-residual first-inputs", GRAPH_RES / "weights_japan_overfit_graph_pga_first_inputs_prior_residual" / "eval_results.txt"),
        ("prior-residual random-inputs", GRAPH_RES / "weights_japan_overfit_graph_pga_random_inputs_prior_residual" / "eval_results.txt"),
        ("graph exp1 same-station", GRAPH_RES / "weights_japan_overfit_exp1_graph_single_input_same_station_pga" / "eval_results.txt"),
        ("graph exp2 multi-target", GRAPH_RES / "weights_japan_overfit_exp2_graph_single_input_multi_target_pga" / "eval_results.txt"),
        ("graph exp3 holdout", GRAPH_RES / "weights_japan_overfit_exp3_graph_snr_filtered_holdout_pga" / "eval_results.txt"),
    ]
    graph = []
    for name, path in rows_graph:
        ev = extract_eval_metrics(path)
        graph.append(
            {
                "Experiment": name,
                "Train MAE": ev.get("train", {}).get("pga", {}).get("mae"),
                "Train Corr": ev.get("train", {}).get("pga", {}).get("corr"),
                "Train R2": ev.get("train", {}).get("pga", {}).get("r2"),
                "Val MAE": ev.get("val", {}).get("pga", {}).get("mae"),
                "Val Corr": ev.get("val", {}).get("pga", {}).get("corr"),
                "Val R2": ev.get("val", {}).get("pga", {}).get("r2"),
                "Val slope": ev.get("val", {}).get("pga", {}).get("slope"),
                "Single val PGA corr": ev.get("single_val", {}).get("single_pga", {}).get("corr"),
            }
        )
    graph_df = pd.DataFrame(graph)

    tables = {
        "transformer_ablation": transformer_df,
        "cross_attention": cross_df,
        "graph_results": graph_df,
    }
    for name, df in tables.items():
        df.to_csv(TABLE_DIR / f"{name}.csv", index=False)
        df.to_markdown(TABLE_DIR / f"{name}.md", index=False, floatfmt=".4f")
    return tables


def plot_data_overview() -> Path:
    events32 = read_csv(DATA32 / "split_events.csv")
    events128 = read_csv(DATA128 / "split_events.csv")
    stations128 = read_csv(DATA128 / "split_stations.csv")
    fig = plt.figure(figsize=(12, 7))
    gs = fig.add_gridspec(2, 2, width_ratios=[1.35, 1.0])

    ax0 = fig.add_subplot(gs[:, 0])
    sc = ax0.scatter(
        stations128["station_lon"],
        stations128["station_lat"],
        c=np.log10(stations128["pga_norm_resampled_mps2"].clip(lower=1e-4)),
        s=8,
        cmap="viridis",
        alpha=0.35,
        linewidths=0,
        label="Stations",
    )
    ax0.scatter(events128["Longitude"], events128["Latitude"], c="crimson", s=28, marker="*", label="Events")
    ax0.set_title("Japan overfit sample: events and stations")
    ax0.set_xlabel("Longitude")
    ax0.set_ylabel("Latitude")
    ax0.legend(loc="lower right", frameon=True)
    cb = fig.colorbar(sc, ax=ax0, fraction=0.035, pad=0.02)
    cb.set_label("log10(PGA m/s²)")

    ax1 = fig.add_subplot(gs[0, 1])
    bins = np.arange(2.5, 8.5, 0.5)
    ax1.hist(events128["Magnitude"], bins=bins, alpha=0.75, color="#4C78A8", label="128 events")
    ax1.hist(events32["Magnitude"], bins=bins, alpha=0.65, color="#F58518", label="32 events")
    ax1.set_title("Magnitude distribution")
    ax1.set_xlabel("Magnitude")
    ax1.set_ylabel("Event count")
    ax1.legend(frameon=False)

    ax2 = fig.add_subplot(gs[1, 1])
    split_counts = events128["split"].value_counts().reindex(["train", "dev", "test"]).fillna(0)
    ax2.bar(split_counts.index.astype(str), split_counts.values, color=["#54A24B", "#ECA82C", "#B279A2"])
    ax2.set_title("128-event split")
    ax2.set_ylabel("Event count")
    for i, v in enumerate(split_counts.values):
        ax2.text(i, v + 0.5, f"{int(v)}", ha="center", va="bottom")

    fig.suptitle("Data diagnostics used in current overfit experiments", fontsize=15, fontweight="bold")
    path = ASSET_DIR / "data_overview.png"
    savefig(path)
    return path


def plot_sample_event() -> Path:
    stations = read_csv(DATA32 / "split_stations.csv")
    event_id = stations.groupby("EVENT").size().sort_values(ascending=False).index[0]
    df = stations[stations["EVENT"] == event_id].copy()
    df["log_pga"] = np.log10(df["pga_norm_resampled_mps2"].clip(lower=1e-4))
    df = df.sort_values("p_pick_aligned")
    input_df = df.head(25)

    fig = plt.figure(figsize=(12, 6.6))
    gs = fig.add_gridspec(1, 2, width_ratios=[1.1, 1.0])
    ax0 = fig.add_subplot(gs[0, 0])
    ax0.scatter(df["station_lon"], df["station_lat"], c=df["log_pga"], cmap="magma", s=20, alpha=0.55, label="Target stations")
    ax0.scatter(input_df["station_lon"], input_df["station_lat"], facecolors="none", edgecolors="cyan", s=70, linewidths=1.4, label="First 25 input stations")
    ax0.scatter(df["Longitude"].iloc[0], df["Latitude"].iloc[0], c="white", edgecolors="black", marker="*", s=180, label="Event")
    ax0.set_title(f"Sample event {event_id}: station geometry")
    ax0.set_xlabel("Longitude")
    ax0.set_ylabel("Latitude")
    ax0.legend(loc="best", frameon=True, fontsize=9)

    ax1 = fig.add_subplot(gs[0, 1])
    ax1.scatter(df["hypocentral_distance_km"], df["pga_norm_resampled_mps2"], s=18, alpha=0.55, color="#4C78A8", label="All target stations")
    ax1.scatter(input_df["hypocentral_distance_km"], input_df["pga_norm_resampled_mps2"], s=50, facecolors="none", edgecolors="#F58518", linewidths=1.4, label="Input stations")
    ax1.set_yscale("log")
    ax1.set_xlabel("Hypocentral distance (km)")
    ax1.set_ylabel("PGA (m/s², log scale)")
    ax1.set_title("Target PGA vs distance")
    ax1.legend(frameon=False)
    fig.suptitle("One event illustrates same-station / multi-target / holdout PGA tasks", fontsize=14, fontweight="bold")
    path = ASSET_DIR / "sample_event_geometry.png"
    savefig(path)
    return path


def plot_waveform_example() -> Path:
    with h5py.File(H5_PATH, "r") as f:
        event_key = sorted(f["data"].keys())[0]
        wave = np.asarray(f["data"][event_key]["waveforms"][:3, :, :])
        picks = np.asarray(f["data"][event_key]["p_picks"][:3])
        pga = np.asarray(f["data"][event_key]["pga"][:3])
    t = np.arange(wave.shape[1]) / 100.0
    fig, axes = plt.subplots(3, 1, figsize=(11, 5.4), sharex=True)
    comp_names = ["EW", "NS", "UD"]
    for i, ax in enumerate(axes):
        for c in range(3):
            sig = wave[i, :, c]
            sig = sig / (np.nanmax(np.abs(sig)) + 1e-8)
            ax.plot(t, sig + c * 2.2, lw=0.8, label=comp_names[c] if i == 0 else None)
        ax.axvline(picks[i] / 100.0, color="crimson", ls="--", lw=1.2)
        ax.text(0.01, 0.87, f"station {i}, PGA={pga[i]:.3f}", transform=ax.transAxes, fontsize=9)
        ax.set_yticks([])
        ax.set_ylabel(f"S{i}")
    axes[0].legend(loc="upper right", ncol=3, frameon=False)
    axes[-1].set_xlabel("Time in 100 s window (s)")
    fig.suptitle("Example waveforms from japan_overfit.hdf5 with P-pick markers", fontsize=14, fontweight="bold")
    path = ASSET_DIR / "waveform_example.png"
    savefig(path)
    return path


def draw_box(ax, xy, wh, text, fc="#DDEBF7", ec="#1F4E79", fontsize=10):
    x, y = xy
    w, h = wh
    rect = plt.Rectangle((x, y), w, h, facecolor=fc, edgecolor=ec, linewidth=1.5, joinstyle="round")
    ax.add_patch(rect)
    ax.text(x + w / 2, y + h / 2, text, ha="center", va="center", fontsize=fontsize, wrap=True)


def arrow(ax, start, end, color="#555555"):
    ax.annotate("", xy=end, xytext=start, arrowprops=dict(arrowstyle="->", lw=1.5, color=color))


def plot_architecture_team() -> Path:
    fig, ax = plt.subplots(figsize=(12, 5.7))
    ax.axis("off")
    boxes = [
        ((0.5, 3.6), (1.55, 0.65), "Waveforms\n(B,S,C,T)", "#E2F0D9"),
        ((2.45, 3.6), (1.55, 0.65), "DiTing\nencoder", "#E7E6E6"),
        ((4.35, 3.6), (1.8, 0.65), "Station adapter\nstation_feature_emb", "#DDEBF7"),
        ((6.55, 3.6), (1.65, 0.65), "LN + scale\n+ coord fusion", "#DDEBF7"),
        ((8.65, 3.6), (1.55, 0.65), "TEAM\nTransformer", "#FCE4D6"),
        ((10.6, 4.35), (1.35, 0.55), "Event token\nmag / loc", "#FFF2CC"),
        ((10.6, 2.85), (1.35, 0.55), "PGA target\nquery token", "#FFF2CC"),
    ]
    for b in boxes:
        draw_box(ax, *b)
    arrow(ax, (2.05, 3.93), (2.45, 3.93))
    arrow(ax, (4.0, 3.93), (4.35, 3.93))
    arrow(ax, (6.15, 3.93), (6.55, 3.93))
    arrow(ax, (8.2, 3.93), (8.65, 3.93))
    arrow(ax, (10.2, 4.0), (10.6, 4.55))
    arrow(ax, (10.2, 3.75), (10.6, 3.05))
    ax.text(0.5, 1.25, "Original readout risk: event/PGA readout tokens are expected to discover the routing through full self-attention.", fontsize=12, color="#7F6000")
    ax.set_xlim(0, 12.3)
    ax.set_ylim(1.0, 5.2)
    path = ASSET_DIR / "architecture_team.png"
    savefig(path)
    return path


def plot_architecture_cross() -> Path:
    fig, ax = plt.subplots(figsize=(12, 5.7))
    ax.axis("off")
    draw_box(ax, (0.7, 3.4), (2.0, 0.75), "station_feature_emb\n(B,S,D)", "#DDEBF7")
    draw_box(ax, (3.35, 4.25), (1.8, 0.6), "learned\nevent query", "#FFF2CC")
    draw_box(ax, (3.35, 2.55), (1.8, 0.7), "target coord emb\n+ learned PGA token", "#FFF2CC")
    draw_box(ax, (5.85, 4.1), (2.15, 0.75), "CrossAttention\nQ=event, K/V=stations", "#E2F0D9")
    draw_box(ax, (5.85, 2.55), (2.15, 0.75), "CrossAttention\nQ=target, K/V=stations", "#E2F0D9")
    draw_box(ax, (8.75, 4.15), (1.7, 0.6), "mag / loc\nheads", "#FCE4D6")
    draw_box(ax, (8.75, 2.65), (1.7, 0.6), "PGA head", "#FCE4D6")
    arrow(ax, (2.7, 3.8), (5.85, 4.45))
    arrow(ax, (2.7, 3.75), (5.85, 2.95))
    arrow(ax, (5.15, 4.55), (5.85, 4.55))
    arrow(ax, (5.15, 2.9), (5.85, 2.9))
    arrow(ax, (8.0, 4.48), (8.75, 4.48))
    arrow(ax, (8.0, 2.92), (8.75, 2.92))
    ax.text(0.7, 1.35, "Key diagnostic change: readout values come only from station tokens; PGA query is target-specific and task-specific.", fontsize=12, color="#1F4E79")
    ax.set_xlim(0, 11.2)
    ax.set_ylim(1.0, 5.2)
    path = ASSET_DIR / "architecture_cross_attention.png"
    savefig(path)
    return path


def plot_architecture_graph() -> Path:
    fig, ax = plt.subplots(figsize=(12, 5.7))
    ax.axis("off")
    draw_box(ax, (0.55, 3.75), (1.9, 0.7), "station_feature_emb\n(B,S,D)", "#DDEBF7")
    draw_box(ax, (3.0, 4.1), (1.75, 0.62), "station PGA\nprior head", "#FFF2CC")
    draw_box(ax, (5.25, 4.1), (1.9, 0.62), "distance weighted\nbaseline", "#E2F0D9")
    draw_box(ax, (3.0, 2.45), (1.75, 0.72), "target-station\nedge features", "#FCE4D6")
    draw_box(ax, (5.25, 2.45), (1.9, 0.72), "GraphPGAReadout\nmessage passing", "#D9EAD3")
    draw_box(ax, (7.75, 2.65), (1.5, 0.62), "residual\nhead", "#FCE4D6")
    draw_box(ax, (9.75, 3.35), (1.6, 0.75), "final PGA\nbaseline + residual", "#E2F0D9")
    arrow(ax, (2.45, 4.1), (3.0, 4.4))
    arrow(ax, (4.75, 4.4), (5.25, 4.4))
    arrow(ax, (2.45, 3.85), (5.25, 2.85))
    arrow(ax, (4.75, 2.8), (5.25, 2.8))
    arrow(ax, (7.15, 2.8), (7.75, 2.95))
    arrow(ax, (7.15, 4.4), (9.75, 3.85))
    arrow(ax, (9.25, 2.95), (9.75, 3.55))
    ax.text(0.55, 1.25, "Graph route: do not predict PGA from scratch; use single-station PGA prior and learn a spatial residual.", fontsize=12, color="#548235")
    ax.set_xlim(0, 11.8)
    ax.set_ylim(1.0, 5.2)
    path = ASSET_DIR / "architecture_graph_prior_residual.png"
    savefig(path)
    return path


def plot_results(tables: dict[str, pd.DataFrame]) -> dict[str, Path]:
    paths = {}
    df = tables["transformer_ablation"].copy()
    fig, ax = plt.subplots(figsize=(10.5, 5.2))
    x = np.arange(len(df))
    ax.bar(x - 0.18, df["Val Corr"], width=0.36, label="Val Corr", color="#4C78A8")
    ax.bar(x + 0.18, df["Val slope"], width=0.36, label="Val slope", color="#F58518")
    ax.axhline(0, color="black", lw=0.8)
    ax.set_xticks(x)
    ax.set_xticklabels(df["Method"], rotation=18, ha="right")
    ax.set_title("PGA readout ablations: original query_transformer collapses")
    ax.set_ylabel("Metric value")
    ax.legend(frameon=False)
    paths["transformer_ablation"] = ASSET_DIR / "result_transformer_ablation.png"
    savefig(paths["transformer_ablation"])

    df = tables["cross_attention"].copy().dropna(subset=["Val Corr"])
    fig, ax = plt.subplots(figsize=(10.5, 5.2))
    x = np.arange(len(df))
    ax.bar(x - 0.22, df["Val Corr"], width=0.22, label="Val Corr", color="#4C78A8")
    ax.bar(x, df["Val R2"], width=0.22, label="Val R²", color="#54A24B")
    ax.bar(x + 0.22, df["Val MAE"], width=0.22, label="Val MAE", color="#ECA82C")
    ax.axhline(0, color="black", lw=0.8)
    ax.set_xticks(x)
    ax.set_xticklabels(df["Experiment"], rotation=22, ha="right")
    ax.set_title("Cross-attention PGA results")
    ax.legend(frameon=False, ncol=3)
    paths["cross_attention"] = ASSET_DIR / "result_cross_attention.png"
    savefig(paths["cross_attention"])

    df = tables["graph_results"].iloc[:3].copy()
    fig, ax = plt.subplots(figsize=(10.5, 5.2))
    x = np.arange(len(df))
    ax.bar(x - 0.22, df["Val Corr"], width=0.22, label="Val Corr", color="#4C78A8")
    ax.bar(x, df["Val R2"], width=0.22, label="Val R²", color="#54A24B")
    ax.bar(x + 0.22, df["Val MAE"], width=0.22, label="Val MAE", color="#ECA82C")
    ax.axhline(0, color="black", lw=0.8)
    ax.set_xticks(x)
    ax.set_xticklabels(df["Experiment"], rotation=16, ha="right")
    ax.set_title("Graph prior-residual improves over old graph readout")
    ax.legend(frameon=False, ncol=3)
    paths["graph_results"] = ASSET_DIR / "result_graph_prior_residual.png"
    savefig(paths["graph_results"])
    return paths


def plot_variance() -> Path:
    rows = [
        ("query_transformer", TEAM_RES / "weights_japan_overfit_exp1_attn_diag"),
        ("direct_station", TEAM_RES / "weights_japan_overfit_exp1_direct_station"),
        ("cross fixed targets", TEAM_RES / "weights_japan_overfit_pga15_cross_attention_overfit32_fixed_inputs_targets"),
        ("graph old", GRAPH_RES / "weights_japan_overfit_graph_pga_first_inputs"),
        ("graph prior", GRAPH_RES / "weights_japan_overfit_graph_pga_first_inputs_prior_residual"),
        ("graph random prior", GRAPH_RES / "weights_japan_overfit_graph_pga_random_inputs_prior_residual"),
    ]
    data = []
    for name, d in rows:
        pred = last_value(d / "diag_pga_mu_best_valid_std.csv") or last_value(d / "diag_pga_point_mu_std.csv")
        target = last_value(d / "diag_pga_target_std.csv")
        data.append((name, pred, target))
    df = pd.DataFrame(data, columns=["Experiment", "Pred std", "Target std"])
    df.to_csv(TABLE_DIR / "variance_summary.csv", index=False)
    fig, ax = plt.subplots(figsize=(10.8, 5.4))
    x = np.arange(len(df))
    ax.bar(x - 0.18, df["Pred std"], width=0.36, label="Pred std", color="#4C78A8")
    ax.bar(x + 0.18, df["Target std"], width=0.36, label="Target std", color="#F58518")
    ax.set_xticks(x)
    ax.set_xticklabels(df["Experiment"], rotation=20, ha="right")
    ax.set_ylabel("Std of PGA log/value")
    ax.set_title("Prediction variance: collapse and conservative outputs")
    ax.legend(frameon=False)
    path = ASSET_DIR / "variance_compression.png"
    savefig(path)
    return path


def generate_assets() -> tuple[dict[str, Path], dict[str, pd.DataFrame]]:
    ensure_dirs()
    tables = build_result_tables()
    assets = {
        "data_overview": plot_data_overview(),
        "sample_event": plot_sample_event(),
        "waveform": plot_waveform_example(),
        "architecture_team": plot_architecture_team(),
        "architecture_cross": plot_architecture_cross(),
        "architecture_graph": plot_architecture_graph(),
        "variance": plot_variance(),
    }
    assets.update(plot_results(tables))
    return assets, tables


def set_text(run, size=18, bold=False, color=RGBColor(0, 0, 0)):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"


def add_title(slide, title: str, subtitle: str | None = None):
    box = slide.shapes.add_textbox(Inches(0.35), Inches(0.2), Inches(12.6), Inches(0.52))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLUE
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.38), Inches(0.76), Inches(12.2), Inches(0.32))
        p2 = box2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY


def add_bullets(slide, bullets, x, y, w, h, font_size=16):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(35, 35, 35)
        p.space_after = Pt(4)
    return box


def add_image(slide, path: Path, x, y, w=None, h=None):
    kwargs = {}
    if w is not None:
        kwargs["width"] = Inches(w)
    if h is not None:
        kwargs["height"] = Inches(h)
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), **kwargs)


def add_table(slide, df: pd.DataFrame, x, y, w, h, font_size=9, cols=None):
    if cols is not None:
        df = df[cols]
    df = df.copy()
    for c in df.columns:
        if pd.api.types.is_float_dtype(df[c]):
            df[c] = df[c].map(lambda v: "" if pd.isna(v) else f"{v:.3f}")
    rows, cols_n = df.shape[0] + 1, df.shape[1]
    table = slide.shapes.add_table(rows, cols_n, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
    for display_i, (_, row) in enumerate(df.iterrows()):
        for j, val in enumerate(row):
            cell = table.cell(display_i + 1, j)
            cell.text = str(val)
            if display_i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(242, 246, 250)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = RGBColor(30, 30, 30)
    return table


def add_section_tag(slide, text: str, color=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(6.78), Inches(2.4), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE


def build_ppt(assets: dict[str, Path], tables: dict[str, pd.DataFrame]) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
    title = s.shapes.add_textbox(Inches(0.65), Inches(1.45), Inches(12), Inches(1.2))
    p = title.text_frame.paragraphs[0]
    p.text = "DiTing 表征接入多台站 PGA / Event 预测"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLUE
    sub = s.shapes.add_textbox(Inches(0.7), Inches(2.8), Inches(11.5), Inches(0.8))
    p = sub.text_frame.paragraphs[0]
    p.text = "Transformer Readout 与 Graph Readout 的实验诊断"
    p.font.size = Pt(22)
    p.font.color.rgb = GRAY
    add_bullets(s, ["技术交流：当前进展、问题定位、下一步请专家建议", "2026-05-01"], 0.75, 4.3, 10.5, 1.0, 16)

    # 2
    s = prs.slides.add_slide(blank)
    add_title(s, "这次交流想回答的问题")
    add_bullets(
        s,
        [
            "DiTing station feature 是否包含 PGA / event 信息？",
            "为什么 full model 容易输出常数？问题在 backbone、head，还是 readout 路由？",
            "Cross-attention 和 graph prior-residual 两条路线分别解决了什么？",
            "下一步模型结构、物理先验和验证实验该如何设计？",
        ],
        0.8,
        1.35,
        11.6,
        4.8,
        22,
    )
    add_section_tag(s, "Motivation", BLUE)

    # 3
    s = prs.slides.add_slide(blank)
    add_title(s, "任务与数据设置", "输入多台站波形与坐标；输出 event-level mag/loc 和 target-level PGA")
    add_bullets(
        s,
        [
            "输入：N 个 station 的三分量波形、station 坐标、P pick、target station 坐标。",
            "输出：event-level magnitude / location；target-level PGA。",
            "当前 Japan overfit 小样本用于结构诊断，验证集很小，不能直接作为泛化性能结论。",
            "实验对照包括 same-station、multi-target、holdout、fixed input/target 等设置。",
        ],
        0.7,
        1.35,
        5.2,
        5.5,
        15,
    )
    add_image(s, assets["data_overview"], 6.0, 1.15, w=6.8)
    add_section_tag(s, "Data", TEAL)

    # 4
    s = prs.slides.add_slide(blank)
    add_title(s, "数据覆盖与抽样样本", "地理覆盖、震级分布和 overfit split 是当前结果解释的边界条件")
    add_image(s, assets["data_overview"], 0.45, 1.0, w=7.0)
    add_image(s, assets["sample_event"], 7.0, 1.0, w=5.9)
    add_section_tag(s, "Data", TEAL)

    # 5
    s = prs.slides.add_slide(blank)
    add_title(s, "一个事件的波形与 PGA 样本", "P pick、PGA label 和台站距离共同定义了 station-level 与 target-level 学习信号")
    add_image(s, assets["waveform"], 0.6, 1.05, w=6.2)
    add_image(s, assets["sample_event"], 6.8, 1.05, w=6.0)
    add_section_tag(s, "Data", TEAL)

    # 6
    s = prs.slides.add_slide(blank)
    add_title(s, "基础模型结构：DiTing + TEAM full model", "DiTing 作为 station feature extractor；问题主要集中在 full model readout")
    add_image(s, assets["architecture_team"], 0.45, 1.0, w=12.4)
    add_bullets(s, ["关键中间量：station_feature_emb (B,S,D)", "原始 PGA/event readout 依赖 full self-attention 中的 token 路由"], 0.85, 6.15, 11.5, 0.55, 12)
    add_section_tag(s, "Model", BLUE)

    # 7
    s = prs.slides.add_slide(blank)
    add_title(s, "Single-station 预训练：station feature 有效", "单台站任务能学到 mag、epidist、PGA 信号")
    add_bullets(
        s,
        [
            "Single-station PGA 在多个实验中保持较高相关性。",
            "same-station 场景下，单台站 head 能直接从 DiTing station embedding 读出 PGA。",
            "这说明 full model 常数化不能简单归因于 DiTing feature 无效。",
        ],
        0.7,
        1.25,
        5.3,
        2.8,
        15,
    )
    single_rows = pd.DataFrame(
        [
            ["TEAM same-station", "single/PGA", 0.3009, 0.6674],
            ["Graph first prior", "single/PGA", 0.1977, 0.8721],
            ["Cross overfit128", "single/PGA", 0.1654, 0.9313],
        ],
        columns=["Experiment", "Task", "Val MAE", "Val Corr"],
    )
    add_table(s, single_rows, 6.2, 1.35, 6.4, 1.6, 10)
    add_image(s, assets["variance"], 6.2, 3.25, w=6.3)
    add_section_tag(s, "Evidence", TEAL)

    # 8
    s = prs.slides.add_slide(blank)
    add_title(s, "原始 full self-attention readout 的 PGA 坍塌")
    add_image(s, assets["transformer_ablation"], 0.55, 1.1, w=6.3)
    df = tables["transformer_ablation"][["Method", "Train Corr", "Train slope", "Val Corr", "Val slope"]]
    add_table(s, df, 7.05, 1.2, 5.8, 2.45, 9)
    add_bullets(
        s,
        [
            "query_transformer：PGA slope≈0，输出近似常数。",
            "mask_batch1：batch=1 仍常数，跨 batch 泄露不是主因。",
            "direct_station：train corr≈0.84，证明 station embedding + PGA head 可学。",
        ],
        7.1,
        4.25,
        5.6,
        1.7,
        13,
    )
    add_section_tag(s, "Diagnosis", ORANGE)

    # 9
    s = prs.slides.add_slide(blank)
    add_title(s, "Attention / mask 诊断：不是完全看不到 station，而是 readout 不稳定")
    add_image(s, assets["variance"], 0.6, 1.0, w=6.0)
    add_bullets(
        s,
        [
            "PGA query 第一层后期能 attend 到 station，但最后输出 std 仍接近 0。",
            "event token / residual / LayerNorm / FFN 可能形成均值解通道。",
            "mask sanity 中 0.11111 是诊断口径问题：att_mask 表示 key 类型可见性，padding_mask 表示真实有效 station。",
            "已新增 effective key mask 诊断以区分两类 mask。",
        ],
        7.0,
        1.2,
        5.8,
        4.5,
        14,
    )
    add_section_tag(s, "Diagnosis", ORANGE)

    # 10
    s = prs.slides.add_slide(blank)
    add_title(s, "Cross-attention readout 设计", "让 event/PGA query 单向、可诊断地读取 station tokens")
    add_image(s, assets["architecture_cross"], 0.55, 1.0, w=12.2)
    add_bullets(s, ["PGA query = target coordinate embedding + learned PGA query token", "readout values come only from station tokens"], 0.9, 6.05, 11.5, 0.65, 12)
    add_section_tag(s, "Cross-attn", BLUE)

    # 11
    s = prs.slides.add_slide(blank)
    add_title(s, "Cross-attention 实验结果", "PGA readout 明显避免常数化；event 联合训练仍需继续诊断")
    add_image(s, assets["cross_attention"], 0.55, 1.0, w=6.25)
    df = tables["cross_attention"][["Experiment", "Val MAE", "Val Corr", "Val R2", "Val slope"]].head(5)
    add_table(s, df, 6.95, 1.05, 6.0, 2.6, 8)
    add_bullets(
        s,
        [
            "fixed input/target：val corr≈0.62，R²≈0.27。",
            "overfit128：val corr≈0.62，R²≈0.32。",
            "event+pga cross-attention first-inputs：PGA val corr 偏低，event/mag/loc 联合训练还需定位。",
        ],
        7.0,
        4.1,
        5.8,
        1.7,
        12,
    )
    add_section_tag(s, "Cross-attn", BLUE)

    # 12
    s = prs.slides.add_slide(blank)
    add_title(s, "Graph 路线动机", "PGA 是空间场：target PGA 应显式依赖输入台站 PGA、距离和传播关系")
    add_bullets(
        s,
        [
            "目标台站 PGA 与输入台站强度、震源/台站几何、距离衰减密切相关。",
            "纯 self-attention / cross-attention 未显式内置距离衰减或传播先验。",
            "Graph readout 可以用 target-station edge feature 表示 dx/dy/dz、distance、log distance、inverse distance。",
            "但旧 graph readout 仍会坍塌，说明需要把 single-station PGA 信号显式接入。",
        ],
        0.75,
        1.25,
        5.8,
        4.6,
        15,
    )
    add_image(s, assets["sample_event"], 6.75, 1.05, w=5.9)
    add_section_tag(s, "Graph", TEAL)

    # 13
    s = prs.slides.add_slide(blank)
    add_title(s, "旧 graph readout 与失败", "单纯 graph message passing 不能自动解决 station PGA 信息读出")
    old_graph = tables["graph_results"][tables["graph_results"]["Experiment"].str.contains("old graph|exp", regex=True)][
        ["Experiment", "Val MAE", "Val Corr", "Val R2", "Val slope", "Single val PGA corr"]
    ]
    add_table(s, old_graph, 0.65, 1.1, 12.0, 2.4, 8)
    add_bullets(
        s,
        [
            "旧 graph first-inputs：val corr≈-0.06，slope≈0。",
            "exp1/exp2/exp3 多种任务下也接近常数预测。",
            "Single-station 分支仍有明显信号，因此 graph readout 需要更强的先验或监督路径。",
        ],
        1.0,
        4.15,
        11.5,
        1.5,
        15,
    )
    add_section_tag(s, "Graph", TEAL)

    # 14
    s = prs.slides.add_slide(blank)
    add_title(s, "Graph prior-residual 架构", "final PGA = distance-weighted station prior baseline + graph residual")
    add_image(s, assets["architecture_graph"], 0.55, 1.0, w=12.2)
    add_section_tag(s, "Graph", TEAL)

    # 15
    s = prs.slides.add_slide(blank)
    add_title(s, "Graph prior-residual 实验结果", "显式 single-station prior 和距离 baseline 明显解除常数化")
    add_image(s, assets["graph_results"], 0.55, 1.0, w=6.25)
    df = tables["graph_results"].iloc[:3][["Experiment", "Val MAE", "Val Corr", "Val R2", "Val slope"]]
    add_table(s, df, 6.95, 1.15, 6.0, 1.9, 8)
    add_bullets(
        s,
        [
            "first-inputs prior-residual：val MAE≈0.278，corr≈0.612，R²≈0.359。",
            "random-inputs prior-residual：val corr≈0.560，R²≈0.209。",
            "旧 graph readout：val corr≈-0.059，R²<0，slope≈0。",
        ],
        7.0,
        3.7,
        5.8,
        1.8,
        13,
    )
    add_section_tag(s, "Graph", TEAL)

    # 16
    s = prs.slides.add_slide(blank)
    add_title(s, "预测方差与保守性", "prior-residual 后 pred std 恢复，但仍小于 target std")
    add_image(s, assets["variance"], 0.65, 1.0, w=7.0)
    add_bullets(
        s,
        [
            "常数化模型 pred std 接近 0。",
            "direct station、cross-attention、graph prior-residual 都能恢复一定预测方差。",
            "graph prior-residual 的 pred std 仍低于 target std，说明 PGA 动态范围仍偏保守。",
            "下一步需要拆 baseline 和 residual 的贡献。",
        ],
        8.05,
        1.2,
        4.7,
        4.3,
        14,
    )
    add_section_tag(s, "Problem", ORANGE)

    # 17
    s = prs.slides.add_slide(blank)
    add_title(s, "两条路线对比")
    compare = pd.DataFrame(
        [
            ["query_transformer", "统一 token 模型", "常数化，readout 不稳"],
            ["direct_station", "证明 station feature/head 可用", "不能处理任意 target 和空间传播"],
            ["cross-attention", "target/event 显式读 station", "空间先验不足，event 联合训练待解"],
            ["graph prior-residual", "single-station prior + 距离传播", "residual 贡献未拆清，方差偏保守"],
        ],
        columns=["Route", "What it solves", "Open issue"],
    )
    add_table(s, compare, 0.65, 1.1, 12.0, 2.5, 10)
    add_bullets(
        s,
        [
            "当前更像是 readout 和 inductive bias 问题，而不是 DiTing backbone 完全无效。",
            "Cross-attention 适合任意 target readout 诊断；graph prior-residual 适合引入空间传播先验。",
            "最终路线可能是二者融合：target-to-station readout + distance/GMPE prior 或 graph residual。",
        ],
        0.9,
        4.15,
        11.8,
        1.7,
        15,
    )
    add_section_tag(s, "Synthesis", BLUE)

    # 18
    s = prs.slides.add_slide(blank)
    add_title(s, "当前主要问题与下一步")
    add_bullets(
        s,
        [
            "当前问题：full event readout 仍不稳定；val event 数少；PGA 输出方差偏保守；graph baseline/residual 贡献未消融。",
            "短期：cross-attention 加 event context、distance bias、relative coordinate MLP。",
            "短期：graph 做 baseline-only / residual-only / baseline+residual 消融。",
            "诊断：记录 residual std、baseline-target corr、residual-target corr、attention top-k distance。",
            "验证：扩大 split 或多 seed / 多 split，避免 4 个 val event 造成误判。",
        ],
        0.75,
        1.1,
        12.1,
        4.9,
        16,
    )
    add_section_tag(s, "Next", ORANGE)

    # 19
    s = prs.slides.add_slide(blank)
    add_title(s, "希望专家重点建议")
    add_bullets(
        s,
        [
            "PGA readout 中 event 信息应如何引入：event query、mag/loc context、GMPE prior，还是显式震源参数？",
            "任意 target PGA 更适合 cross-attention readout，还是 graph message passing？",
            "距离衰减先验更适合作为 baseline、attention bias，还是 residual feature？",
            "小样本 overfit 诊断后，下一步如何设计更稳定的验证集和消融实验？",
            "是否需要引入台站场地项、方位角、路径效应等地震学先验？",
        ],
        0.8,
        1.1,
        11.8,
        5.2,
        17,
    )
    add_section_tag(s, "Discussion", TEAL)

    # Backup slides
    s = prs.slides.add_slide(blank)
    add_title(s, "Backup: Transformer readout ablation full table")
    add_table(s, tables["transformer_ablation"][["Method", "Train MAE", "Train Corr", "Train R2", "Train slope", "Val MAE", "Val Corr", "Val R2", "Val slope"]], 0.35, 1.05, 12.6, 4.4, 7)

    s = prs.slides.add_slide(blank)
    add_title(s, "Backup: Cross-attention full table")
    add_table(s, tables["cross_attention"][["Experiment", "Train MAE", "Train Corr", "Train R2", "Val MAE", "Val Corr", "Val R2", "Val slope", "Val mag MAE"]], 0.28, 1.0, 12.75, 4.7, 7)

    s = prs.slides.add_slide(blank)
    add_title(s, "Backup: Graph full table")
    add_table(s, tables["graph_results"][["Experiment", "Train MAE", "Train Corr", "Train R2", "Val MAE", "Val Corr", "Val R2", "Val slope", "Single val PGA corr"]], 0.28, 1.0, 12.75, 4.8, 7)

    prs.save(PPTX_PATH)
    return PPTX_PATH


def main() -> None:
    assets, tables = generate_assets()
    pptx = build_ppt(assets, tables)
    print(f"PPTX: {pptx}")
    print(f"Assets: {ASSET_DIR}")
    print(f"Tables: {TABLE_DIR}")


if __name__ == "__main__":
    main()
