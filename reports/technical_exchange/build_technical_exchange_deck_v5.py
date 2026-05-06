from __future__ import annotations

from pathlib import Path
from typing import Iterable

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT_DIR = Path(__file__).resolve().parent
ASSET_DIR = OUT_DIR / "assets"
TABLE_DIR = OUT_DIR / "tables"
PPTX_PATH = OUT_DIR / "diting_team_graph_pga_technical_exchange_v5.pptx"

# 16:9 widescreen in inches
SLIDE_W = 13.333
SLIDE_H = 7.5

# Scientific discussion style: calm navy + cyan accent + warm highlight
NAVY = RGBColor(20, 47, 78)
BLUE = RGBColor(31, 78, 121)
CYAN = RGBColor(0, 142, 150)
TEAL = RGBColor(0, 128, 128)
ORANGE = RGBColor(211, 111, 33)
RED = RGBColor(190, 47, 57)
GREEN = RGBColor(63, 140, 89)
PURPLE = RGBColor(112, 85, 166)
BG = RGBColor(247, 250, 252)
PANEL = RGBColor(255, 255, 255)
PANEL_BLUE = RGBColor(236, 246, 250)
PANEL_GREEN = RGBColor(239, 248, 242)
PANEL_ORANGE = RGBColor(255, 245, 236)
PANEL_RED = RGBColor(255, 241, 242)
TEXT = RGBColor(30, 41, 59)
MUTED = RGBColor(100, 116, 139)
LINE = RGBColor(214, 226, 238)
WHITE = RGBColor(255, 255, 255)

FONT = "Microsoft YaHei"
FONT_EN = "Aptos"


def path_asset(name: str) -> Path:
    p = ASSET_DIR / name
    if not p.exists():
        raise FileNotFoundError(f"Missing asset: {p}")
    return p


def read_table(name: str) -> pd.DataFrame:
    p = TABLE_DIR / name
    if not p.exists():
        raise FileNotFoundError(f"Missing table: {p}")
    return pd.read_csv(p)


def fmt_num(v) -> str:
    if pd.isna(v):
        return ""
    if isinstance(v, float):
        return f"{v:.3f}"
    return str(v)


def rgb_hex(c: RGBColor) -> str:
    return f"#{c[0]:02x}{c[1]:02x}{c[2]:02x}"


def set_font(paragraph, size=14, bold=False, color: RGBColor = TEXT, name: str = FONT):
    for run in paragraph.runs:
        run.font.name = name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_textbox(slide, text: str, x, y, w, h, size=14, color: RGBColor = TEXT, bold=False,
                align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.06, font=FONT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_line(slide, x1, y1, x2, y2, color: RGBColor = LINE, width=1.0):
    shape = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    return shape


def add_shape(slide, x, y, w, h, fill: RGBColor, line: RGBColor | None = None,
              radius=True, transparency=0):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.fill.transparency = transparency
    shp.line.color.rgb = line if line is not None else fill
    shp.line.width = Pt(0.8)
    return shp


def add_card(slide, x, y, w, h, title: str, body: str = "", accent: RGBColor = CYAN,
             fill: RGBColor = PANEL, title_size=13.5, body_size=11.5, number: str | None = None):
    add_shape(slide, x, y, w, h, fill, LINE, radius=True)
    add_shape(slide, x, y, 0.08, h, accent, accent, radius=False)
    if number:
        add_shape(slide, x + 0.18, y + 0.18, 0.42, 0.28, accent, accent, radius=True)
        add_textbox(slide, number, x + 0.18, y + 0.17, 0.42, 0.30, 8.5, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, margin=0.0, font=FONT_EN)
        tx = x + 0.68
        tw = w - 0.86
    else:
        tx = x + 0.24
        tw = w - 0.42
    add_textbox(slide, title, tx, y + 0.12, tw, 0.36, title_size, accent, True, margin=0.0)
    if body:
        add_textbox(slide, body, x + 0.24, y + 0.55, w - 0.42, h - 0.66, body_size, TEXT, False, margin=0.0)


def add_pill(slide, text: str, x, y, w, h=0.30, fill: RGBColor = PANEL_BLUE, color: RGBColor = BLUE):
    add_shape(slide, x, y, w, h, fill, fill, radius=True)
    add_textbox(slide, text, x + 0.05, y + 0.02, w - 0.1, h - 0.02, 9.0, color, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, margin=0.0)


def add_bullets(slide, items: Iterable[str], x, y, w, h, size=13.0, color: RGBColor = TEXT,
                bullet_color: RGBColor = CYAN, gap=0.38):
    cur_y = y
    for item in items:
        add_shape(slide, x, cur_y + 0.075, 0.09, 0.09, bullet_color, bullet_color, radius=True)
        add_textbox(slide, item, x + 0.20, cur_y, w - 0.20, gap + 0.12, size, color, False, margin=0.0)
        cur_y += gap


def add_kpi(slide, label: str, value: str, note: str, x, y, w, h, accent: RGBColor):
    add_shape(slide, x, y, w, h, PANEL, LINE, radius=True)
    add_textbox(slide, value, x + 0.15, y + 0.12, w - 0.3, 0.40, 20, accent, True, PP_ALIGN.CENTER, margin=0.0, font=FONT_EN)
    add_textbox(slide, label, x + 0.15, y + 0.58, w - 0.3, 0.28, 10.5, TEXT, True, PP_ALIGN.CENTER, margin=0.0)
    add_textbox(slide, note, x + 0.15, y + 0.92, w - 0.3, h - 0.98, 8.8, MUTED, False, PP_ALIGN.CENTER, margin=0.0)


def add_image_fit(slide, img_path: Path, x, y, w, h, border=True, bg=True):
    with Image.open(img_path) as im:
        iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio >= box_ratio:
        disp_w = w
        disp_h = w / img_ratio
    else:
        disp_h = h
        disp_w = h * img_ratio
    px = x + (w - disp_w) / 2
    py = y + (h - disp_h) / 2
    if bg:
        add_shape(slide, x, y, w, h, WHITE, LINE if border else WHITE, radius=True)
    pic = slide.shapes.add_picture(str(img_path), Inches(px), Inches(py), width=Inches(disp_w), height=Inches(disp_h))
    return pic


def add_table(slide, df: pd.DataFrame, x, y, w, h, font_size=8.2, header_color: RGBColor = BLUE,
              max_rows: int | None = None, cols: list[str] | None = None, widths: list[float] | None = None):
    if cols is not None:
        df = df[cols].copy()
    else:
        df = df.copy()
    if max_rows is not None:
        df = df.head(max_rows).copy()
    for c in df.columns:
        df[c] = df[c].map(fmt_num)
    rows, cols_n = len(df) + 1, len(df.columns)
    table = slide.shapes.add_table(rows, cols_n, Inches(x), Inches(y), Inches(w), Inches(h)).table
    if widths:
        total = sum(widths)
        for j, ww in enumerate(widths):
            table.columns[j].width = Inches(w * ww / total)
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        cell.margin_left = Inches(0.04)
        cell.margin_right = Inches(0.04)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.name = FONT
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
    for i, (_, row) in enumerate(df.iterrows(), start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(248, 251, 253) if i % 2 else WHITE
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                p.font.name = FONT
                p.font.size = Pt(font_size)
                p.font.color.rgb = TEXT
    return table


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_W)
        self.prs.slide_height = Inches(SLIDE_H)
        self.blank = self.prs.slide_layouts[6]
        self.page_no = 0

    def slide(self, title: str, subtitle: str = "", section: str = "", accent: RGBColor = CYAN):
        s = self.prs.slides.add_slide(self.blank)
        self.page_no += 1
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = BG
        # Top brand line
        add_shape(s, 0, 0, SLIDE_W, 0.08, accent, accent, radius=False)
        add_textbox(s, title, 0.50, 0.24, 10.3, 0.45, 22, NAVY, True, margin=0.0)
        if subtitle:
            add_textbox(s, subtitle, 0.52, 0.76, 10.4, 0.25, 9.5, MUTED, False, margin=0.0)
        if section:
            add_pill(s, section, 11.25, 0.26, 1.45, 0.28, fill=PANEL_BLUE, color=BLUE)
        add_line(s, 0.50, 1.05, 12.85, 1.05, LINE, 0.8)
        add_textbox(s, f"{self.page_no:02d}", 12.62, 7.08, 0.28, 0.16, 7.5, MUTED, False, PP_ALIGN.RIGHT, margin=0.0, font=FONT_EN)
        return s

    def title_slide(self):
        s = self.prs.slides.add_slide(self.blank)
        self.page_no += 1
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = BG
        add_shape(s, 0, 0, SLIDE_W, SLIDE_H, BG, BG, radius=False)
        # Decorative panels
        add_shape(s, 0.0, 0.0, 13.333, 0.12, CYAN, CYAN, radius=False)
        add_shape(s, 9.2, 0.7, 3.6, 5.9, PANEL_BLUE, PANEL_BLUE, radius=True)
        add_shape(s, 9.65, 1.10, 2.72, 1.05, WHITE, LINE, radius=True)
        add_shape(s, 9.65, 2.42, 2.72, 1.05, WHITE, LINE, radius=True)
        add_shape(s, 9.65, 3.74, 2.72, 1.05, WHITE, LINE, radius=True)
        add_textbox(s, "DiTing + TEAM/Graph", 0.72, 1.20, 6.5, 0.35, 16, CYAN, True, margin=0.0, font=FONT_EN)
        add_textbox(s, "多台站 PGA / Event 预测\n技术交流", 0.68, 1.72, 7.7, 1.22, 34, NAVY, True, margin=0.0)
        add_textbox(s, "从 readout 坍塌诊断到 graph prior-residual 传播先验", 0.72, 3.12, 7.2, 0.36, 15.5, TEXT, False, margin=0.0)
        add_textbox(s, "科研讨论版 · v5", 0.72, 5.78, 2.2, 0.32, 11, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, margin=0.0)
        add_shape(s, 0.72, 5.75, 2.2, 0.36, CYAN, CYAN, radius=True)
        add_textbox(s, "科研讨论版 · v5", 0.72, 5.78, 2.2, 0.32, 11, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, margin=0.0)
        add_textbox(s, "2026-05-02", 0.72, 6.25, 2.4, 0.24, 10, MUTED, False, margin=0.0, font=FONT_EN)
        # Right chips
        add_textbox(s, "Station signal", 9.85, 1.30, 2.3, 0.24, 11, BLUE, True, PP_ALIGN.CENTER, margin=0.0, font=FONT_EN)
        add_textbox(s, "Readout collapse", 9.85, 2.62, 2.3, 0.24, 11, ORANGE, True, PP_ALIGN.CENTER, margin=0.0, font=FONT_EN)
        add_textbox(s, "Prior-residual", 9.85, 3.94, 2.3, 0.24, 11, GREEN, True, PP_ALIGN.CENTER, margin=0.0, font=FONT_EN)
        add_textbox(s, "有效信号在 station encoder 里；\n瓶颈在 event/target readout 与空间传播先验。", 0.72, 4.30, 7.3, 0.66, 16, NAVY, True, margin=0.0)
        return s

    def save(self, path: Path):
        self.prs.save(path)


def main():
    # Load tables
    all_exp = read_table("all_experiments_summary.csv")
    transformer = read_table("transformer_ablation.csv")
    cross = read_table("cross_attention.csv")
    graph = read_table("graph_results.csv")
    variance = read_table("variance_summary.csv")
    single = read_table("single_station_metrics.csv")
    model_scale = read_table("model_scale_params.csv")
    loss_design = read_table("loss_design.csv")

    # Short labels for dense tables
    transformer_short = transformer[["Method", "Train Corr", "Train slope", "Val Corr", "Val slope"]].copy()
    cross_short = cross[["Experiment", "Train Corr", "Val Corr", "Val MAE", "Val slope"]].copy()
    cross_short["Experiment"] = cross_short["Experiment"].replace({
        "pga15_cross_overfit32": "overfit32 random",
        "pga15_cross_overfit128": "overfit128 random",
        "fixed_inputs_targets": "fixed in/fixed tgt",
        "input_targets": "input as target",
        "event_pga_cross_first_inputs": "event+pga first",
    })
    graph_short = graph[["Experiment", "Train Corr", "Val Corr", "Val MAE", "Val slope", "Single val PGA corr"]].copy()
    graph_short["Experiment"] = graph_short["Experiment"].replace({
        "old graph first-inputs": "old graph first",
        "prior-residual first-inputs": "prior-residual first",
        "prior-residual random-inputs": "prior-residual random",
        "graph exp1 same-station": "exp1 same-station",
        "graph exp2 multi-target": "exp2 multi-target",
        "graph exp3 holdout": "exp3 holdout",
    })
    variance_short = variance.copy()
    variance_short["std ratio"] = variance_short["Pred std"] / variance_short["Target std"]

    d = Deck()
    d.title_slide()

    # 2 Conclusion first
    s = d.slide("先给结论：当前确认了什么？", "把讨论焦点从“哪个实验最好”转向“瓶颈在哪里”", "Conclusion", CYAN)
    add_card(s, 0.65, 1.35, 3.85, 1.35, "1. Station 表征有信号", "single-station PGA / epidist / mag 均有相关性，不是底层表征完全失效。", CYAN, PANEL)
    add_card(s, 4.75, 1.35, 3.85, 1.35, "2. Full self-attention 不稳", "loss 可下降，但 PGA / mag / loc 可能坍塌成近常数。", ORANGE, PANEL)
    add_card(s, 8.85, 1.35, 3.85, 1.35, "3. Cross-attention 改善 readout", "target 显式读取 station tokens 后，overfit_n=32/128 均看到 PGA 信号。", BLUE, PANEL)
    add_card(s, 0.65, 3.05, 3.85, 1.35, "4. 有 graph 不等于有先验", "旧 graph message passing 仍接近常数，说明结构本身不足。", RED, PANEL)
    add_card(s, 4.75, 3.05, 3.85, 1.35, "5. Prior-residual 最有解释性", "single-station prior + distance baseline + residual 明显提升 val corr 与输出方差。", GREEN, PANEL)
    add_card(s, 8.85, 3.05, 3.85, 1.35, "6. 仍是诊断，不是最终结论", "验证集小；需 P-pick 质检、多 seed/split 与结构消融。", PURPLE, PANEL)
    add_shape(s, 0.65, 5.25, 12.05, 0.70, PANEL_BLUE, LINE, radius=True)
    add_textbox(s, "一句话：有效信号在 station encoder 里；瓶颈在 event/target readout 与空间传播先验。", 0.95, 5.40, 11.45, 0.35, 17, NAVY, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, margin=0.0)

    # 3 Questions
    s = d.slide("这次希望专家帮忙看的点", "把后续实验都对应到可讨论的问题", "Discussion", TEAL)
    qs = [
        ("P-pick / 波形窗口是否可靠？", "哪些低 SNR 或 STA/LTA 失败样本应剔除、降权或重新对齐？", RED),
        ("Event 信息如何进入 PGA？", "event query、mag/loc context gate、GMPE prior，还是显式震源参数？", ORANGE),
        ("任意 target PGA 用什么 readout？", "cross-attention、graph prior-residual，还是二者融合？", BLUE),
        ("何时从 overfit 诊断转向大训练？", "需要哪些 sanity、split、seed 和消融结果作为门槛？", GREEN),
    ]
    for i, (t, b, c) in enumerate(qs):
        add_card(s, 0.85 + (i % 2) * 6.05, 1.45 + (i // 2) * 1.95, 5.55, 1.38, t, b, c, PANEL, number=str(i + 1))
    add_textbox(s, "建议讨论方式：先判断数据可靠性，再判断结构先验；否则模型问题和数据问题容易混在一起。", 1.05, 5.85, 11.2, 0.32, 13, MUTED, False, PP_ALIGN.CENTER, margin=0.0)

    # 4 Task transition
    s = d.slide("任务转换：从单台站 DiTing 到多台站 event/PGA", "当前任务不是简单迁移，而是从局部台站任务变为空间场 readout", "Task", BLUE)
    add_image_fit(s, path_asset("task_setup_multistation.png"), 0.65, 1.30, 8.25, 4.65)
    add_card(s, 9.25, 1.30, 3.35, 1.05, "输入", "N 个 station 波形 + 坐标 + station_valid mask", BLUE, PANEL)
    add_card(s, 9.25, 2.65, 3.35, 1.05, "查询", "event query；PGA target query", CYAN, PANEL)
    add_card(s, 9.25, 4.00, 3.35, 1.05, "输出", "event-level mag/loc；target-level PGA", GREEN, PANEL)
    add_textbox(s, "关键难点：target station 与输入 station 不一定重合，PGA 是空间传播场，需要先验。", 9.35, 5.55, 3.1, 0.42, 11.2, TEXT, True, margin=0.0)

    # 5 Data boundary
    s = d.slide("数据与 overfit 诊断边界", "当前结论定位为结构诊断，不作为最终泛化结论", "Data", BLUE)
    add_image_fit(s, path_asset("data_overview.png"), 0.55, 1.25, 6.25, 4.75)
    add_image_fit(s, path_asset("sample_event_geometry.png"), 6.95, 1.25, 5.85, 3.10)
    add_card(s, 7.05, 4.65, 2.75, 1.08, "看什么", "train 是否可记忆；val 是否有信号；输出方差是否摆脱常数。", CYAN, PANEL)
    add_card(s, 10.05, 4.65, 2.75, 1.08, "不看什么", "不能只看 loss；不能用小 val 直接宣称泛化。", ORANGE, PANEL)
    add_pill(s, "指标组合：MAE · Corr · R² · slope · pred std / target std", 7.18, 5.95, 5.35, 0.32, PANEL_BLUE, BLUE)

    # 6 Implementation overview
    s = d.slide("实现级模型总览：共用 encoder + 不同 readout", "后续所有结果都可以映射到这张图", "Model", CYAN)
    add_image_fit(s, path_asset("architecture_implementation_detail.png"), 0.45, 1.18, 12.45, 5.52)

    # 7 Model scale and loss
    s = d.slide("模型规模与训练目标", "当前阶段主要训练 adapter / readout / heads，而不是端到端微调整个 DiTing", "Setup", CYAN)
    scale = model_scale[["Component", "Total params", "Trainable params", "Status"]].copy()
    add_table(s, scale, 0.55, 1.28, 6.1, 2.25, 8.2, widths=[2.4, 1.2, 1.4, 1.2])
    loss = loss_design[["Stage", "Tasks", "Weights", "Implementation note"]].copy()
    add_table(s, loss, 6.95, 1.28, 5.85, 2.25, 7.2, header_color=TEAL, widths=[1.5, 1.3, 1.1, 2.5])
    add_card(s, 0.75, 4.15, 3.75, 1.05, "训练参数", "TEAM cross-attn full model 约 28M 可训练；graph prior-residual 约 38M。", BLUE, PANEL)
    add_card(s, 4.82, 4.15, 3.75, 1.05, "Loss 设计", "Huber point prediction；均值解也可能降低 loss。", ORANGE, PANEL)
    add_card(s, 8.88, 4.15, 3.75, 1.05, "解释原则", "loss 必须和 corr、slope、pred std 一起看。", RED, PANEL)

    # 8 Evidence ladder
    s = d.slide("诊断链条：从表征有效到先验注入", "把多个实验组织成一个递进证据链", "Roadmap", CYAN)
    steps = [
        ("Single-station", "表征有信号", CYAN),
        ("Query Transformer", "full readout 坍塌", ORANGE),
        ("Cross-attention", "显式读取改善", BLUE),
        ("Old Graph", "有 graph 仍不够", RED),
        ("Prior-residual", "先验 + residual 改善", GREEN),
    ]
    x0 = 0.65
    for i, (t, b, c) in enumerate(steps):
        x = x0 + i * 2.45
        add_shape(s, x, 2.15, 1.82, 1.10, PANEL, LINE, radius=True)
        add_shape(s, x + 0.72, 1.62, 0.38, 0.38, c, c, radius=True)
        add_textbox(s, str(i + 1), x + 0.72, 1.64, 0.38, 0.28, 9, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, margin=0.0, font=FONT_EN)
        add_textbox(s, t, x + 0.10, 2.30, 1.62, 0.25, 10.5, c, True, PP_ALIGN.CENTER, margin=0.0, font=FONT_EN)
        add_textbox(s, b, x + 0.12, 2.66, 1.58, 0.32, 10, TEXT, False, PP_ALIGN.CENTER, margin=0.0)
        if i < len(steps) - 1:
            add_line(s, x + 1.82, 2.70, x + 2.45, 2.70, c, 1.4)
    add_shape(s, 1.05, 4.35, 11.25, 0.78, PANEL_BLUE, LINE, radius=True)
    add_textbox(s, "诊断策略：不先追求复杂模型，而是用 overfit / sanity / ablation 判断信息是否能通过 readout 路由到目标。", 1.28, 4.54, 10.8, 0.30, 14, NAVY, True, PP_ALIGN.CENTER, margin=0.0)

    # 9 Single station
    s = d.slide("证据 1：single-station 表征确实有效", "waveform encoder + station adapter 包含 PGA / 震级 / 距离信号", "Evidence", CYAN)
    add_image_fit(s, path_asset("single_station_val_metrics.png"), 0.55, 1.22, 6.25, 4.72)
    single_val = single[single["Split"].eq("val")].copy()
    pga_rows = single_val[single_val["Task"].eq("pga")][["Model", "MAE", "Corr"]].copy()
    add_table(s, pga_rows, 7.05, 1.35, 5.55, 1.45, 8.2, header_color=CYAN, widths=[2.5, 1.0, 1.0])
    add_bullets(s, [
        "Cross overfit128 single val PGA：Corr 0.9313。",
        "Graph prior first-inputs single val PGA：Corr 0.8721。",
        "Graph prior random-inputs single val PGA：Corr 0.7873。",
        "因此 full model 坍塌更像 readout / 空间传播问题。",
    ], 7.20, 3.20, 5.1, 2.1, size=12.3, bullet_color=CYAN)

    # 10 Original transformer collapse
    s = d.slide("证据 2：原始 full self-attention readout 容易坍塌", "PGA token 可以参与 attention，但输出方差仍可能被压扁", "Diagnosis", ORANGE)
    add_image_fit(s, path_asset("architecture_team.png"), 0.55, 1.25, 5.35, 2.65)
    add_image_fit(s, path_asset("result_transformer_ablation.png"), 6.15, 1.25, 6.55, 2.65)
    add_table(s, transformer_short, 0.75, 4.25, 6.2, 1.45, 7.8, header_color=ORANGE, widths=[2.1, 1, 1, 1, 1])
    add_bullets(s, [
        "query_transformer pred std ≈ 3.7e-6，target std ≈ 0.609。",
        "mask_batch1 仍接近常数，跨 batch 泄露不是主因。",
        "mag/loc 也会常数化，说明问题不只在 PGA label。",
    ], 7.25, 4.25, 5.25, 1.45, size=11.5, bullet_color=ORANGE)

    # 11 Ablation conclusions
    s = d.slide("原始 readout ablation：排除了哪些解释？", "用对照实验定位问题，而不是只比较分数", "Diagnosis", ORANGE)
    rows = [
        ["是否 batch/mask 泄露？", "mask_batch1", "仍近常数，不是主因"],
        ["只靠 target 坐标是否足够？", "query_no_transformer", "有相关但不是多台站读取"],
        ["station feature/head 是否无信号？", "direct_station", "可读出信号，feature 可用"],
        ["full transformer 能否自动学 route？", "query_transformer", "不稳定，容易均值解"],
    ]
    add_table(s, pd.DataFrame(rows, columns=["诊断问题", "实验", "结论"]), 0.75, 1.45, 11.85, 2.45, 10.5, header_color=ORANGE, widths=[2.4, 1.5, 3.0])
    add_image_fit(s, path_asset("variance_compression.png"), 0.80, 4.35, 5.25, 1.55)
    add_card(s, 6.40, 4.42, 5.75, 1.05, "关键判断", "不是 station feature 完全无效，而是原始 event/PGA readout 路由与空间传播先验不足。", ORANGE, PANEL_ORANGE)

    # 12 Cross design
    s = d.slide("Cross-attention：把信息路由显式化", "从“让 Transformer 自己发现路径”改为“target query 单向读 station tokens”", "Cross", BLUE)
    add_image_fit(s, path_asset("architecture_cross_attention.png"), 0.55, 1.28, 6.55, 4.65)
    add_card(s, 7.40, 1.45, 4.95, 0.95, "Q", "target coordinate embedding + learned PGA token", BLUE, PANEL)
    add_card(s, 7.40, 2.65, 4.95, 0.95, "K / V", "station_feature_emb；mask = ~station_valid", CYAN, PANEL)
    add_card(s, 7.40, 3.85, 4.95, 0.95, "Readout", "输出不加 query residual，更直接依赖 station tokens", GREEN, PANEL)
    add_textbox(s, "核心作用：降低 readout route 学习难度，避免 event/PGA token 在 full self-attention 中互相干扰。", 7.45, 5.35, 4.75, 0.48, 11.5, TEXT, True, margin=0.0)

    # 13 Cross results
    s = d.slide("Cross-attention 结果：PGA readout 明显改善", "overfit_n=32/128 与 fixed/input target sanity 均看到非平凡信号", "Cross", BLUE)
    add_image_fit(s, path_asset("cross_train_val_detail.png"), 0.55, 1.22, 6.05, 3.60)
    add_table(s, cross_short, 6.85, 1.25, 5.85, 2.65, 7.6, header_color=BLUE, widths=[2.3, 1, 1, 1, 1])
    add_kpi(s, "overfit128 random", "0.618", "Val Corr", 0.90, 5.20, 1.62, 0.95, BLUE)
    add_kpi(s, "fixed targets", "0.621", "Val Corr", 2.75, 5.20, 1.62, 0.95, CYAN)
    add_kpi(s, "input targets", "0.641", "Val Corr", 4.60, 5.20, 1.62, 0.95, GREEN)
    add_bullets(s, [
        "随机 overfit_n=128 比 n=32 更稳。",
        "fixed/input target 设置 train 表现强，route sanity 通过。",
        "event+pga first inputs 的 slope 很低，联合 event 任务仍不稳。",
    ], 7.05, 4.35, 5.25, 1.55, size=11.3, bullet_color=BLUE)

    # 14 Cross limitations
    s = d.slide("Cross-attention 的剩余问题", "显式 readout 有帮助，但空间传播先验仍弱", "Cross", BLUE)
    add_image_fit(s, path_asset("variance_compression.png"), 0.60, 1.25, 5.75, 3.25)
    add_card(s, 6.75, 1.35, 5.50, 0.90, "空间插值仍难", "fixed inputs + random targets 的 val 明显弱，说明任意 target 外推还没解决。", ORANGE, PANEL)
    add_card(s, 6.75, 2.55, 5.50, 0.90, "Event + PGA 联合不稳", "event context 可能没有以正确方式进入 PGA，甚至压扁 PGA slope。", RED, PANEL)
    add_card(s, 6.75, 3.75, 5.50, 0.90, "缺少传播先验", "target 主要靠坐标 query 学空间场，样本效率低。", BLUE, PANEL)
    add_pill(s, "下一步：relative coordinate MLP · distance bias · event context gate · GMPE-style prior", 1.20, 5.38, 10.95, 0.35, PANEL_BLUE, BLUE)

    # 15 Old graph not enough
    s = d.slide("Graph route：为什么“有 graph”还不够？", "旧 graph message passing 仍然可能学成均值解", "Graph", RED)
    add_image_fit(s, path_asset("graph_all_train_val_corr.png"), 0.65, 1.25, 5.75, 4.80)
    old_rows = graph_short[graph_short["Experiment"].isin(["old graph first", "exp1 same-station", "exp2 multi-target", "exp3 holdout"])]
    add_table(s, old_rows[["Experiment", "Train Corr", "Val Corr", "Val slope"]], 6.80, 1.35, 5.55, 1.80, 8.0, header_color=RED, widths=[2.2, 1, 1, 1])
    add_card(s, 6.85, 3.65, 5.45, 1.18, "关键结论", "图结构本身不是先验；如果没有 station prior 和距离传播约束，graph readout 仍可能坍塌。", RED, PANEL_RED)
    add_textbox(s, "因此后续不是简单加 graph 层，而是把可解释的地震传播直觉写进模型。", 6.95, 5.25, 5.25, 0.42, 12.5, TEXT, True, margin=0.0)

    # 16 Graph prior residual design
    s = d.slide("Graph prior-residual：把物理直觉变成可学习结构", "single-station prior 提供局部强度，distance baseline 提供低频传播趋势", "Graph", GREEN)
    add_image_fit(s, path_asset("architecture_graph_prior_residual.png") if (ASSET_DIR / "architecture_graph_prior_residual.png").exists() else path_asset("architecture_graph_detail.png"), 0.55, 1.18, 7.05, 4.88)
    add_card(s, 7.90, 1.32, 4.55, 0.90, "1. Station prior", "station_pga_prior_head 从 single-station pretrain 加载。", CYAN, PANEL)
    add_card(s, 7.90, 2.45, 4.55, 0.90, "2. Distance baseline", "输入台站 prior 的距离加权平均，形成传播初值。", BLUE, PANEL)
    add_card(s, 7.90, 3.58, 4.55, 0.90, "3. Learned residual", "GraphPGAReadout 用 edge features 学路径、方位和非线性修正。", GREEN, PANEL)
    add_shape(s, 7.90, 5.05, 4.55, 0.55, PANEL_GREEN, LINE, radius=True)
    add_textbox(s, "final PGA = distance baseline + learned residual", 8.08, 5.20, 4.20, 0.22, 13, GREEN, True, PP_ALIGN.CENTER, margin=0.0, font=FONT_EN)

    # 17 Graph results
    s = d.slide("Graph prior-residual 结果：当前最有解释性的改善", "改善来自 prior + baseline + residual 的组合，而不是 graph 本身", "Graph", GREEN)
    add_image_fit(s, path_asset("graph_train_val_detail.png"), 0.55, 1.20, 5.95, 3.60)
    key_graph = graph_short[graph_short["Experiment"].isin(["old graph first", "prior-residual first", "prior-residual random"])]
    add_table(s, key_graph, 6.75, 1.25, 5.95, 2.00, 7.4, header_color=GREEN, widths=[2.2, 1, 1, 1, 1, 1.2])
    add_kpi(s, "old graph val corr", "-0.059", "near constant", 0.85, 5.12, 1.70, 0.95, RED)
    add_kpi(s, "prior first val corr", "0.612", "improved", 2.85, 5.12, 1.70, 0.95, GREEN)
    add_kpi(s, "prior random val corr", "0.560", "harder input", 4.85, 5.12, 1.70, 0.95, GREEN)
    add_bullets(s, [
        "old graph pred std≈0.0059，target std≈0.6463。",
        "prior-residual pred std≈0.1853，虽仍偏小但已摆脱近常数。",
    ], 7.00, 4.25, 5.30, 1.05, size=11.3, bullet_color=GREEN)

    # 18 Route comparison
    s = d.slide("两条可行路线的阶段性对比", "cross-attention 适合诊断 readout；graph prior-residual 更适合作为下一阶段主候选", "Synthesis", CYAN)
    comp = pd.DataFrame([
        ["Cross-attention", "query-to-station 路由清晰；PGA sanity 强", "空间传播先验弱；event+PGA 联合不稳", "distance bias / event gate"],
        ["Graph prior-residual", "可解释；利用 single-station prior 和距离 baseline", "仍低估方差；需证明泛化和消融", "baseline-only / prior-only / kNN"],
    ], columns=["路线", "优点", "问题", "下一步"])
    add_table(s, comp, 0.65, 1.40, 12.05, 1.55, 9.2, header_color=CYAN, widths=[1.6, 3.0, 2.8, 2.5])
    add_card(s, 1.00, 3.65, 5.25, 1.30, "推荐策略", "短期：用 cross-attention 继续做 readout 诊断；中期：以 graph prior-residual 为主线做系统消融。", BLUE, PANEL)
    add_card(s, 7.05, 3.65, 5.25, 1.30, "可能融合", "target-to-station attention + distance/GMPE prior + residual graph/message passing。", GREEN, PANEL)
    add_textbox(s, "判断标准：不仅看 val corr，还要看 slope、pred std、baseline-target corr 与 residual 是否真正学习了偏差。", 1.25, 5.55, 10.85, 0.32, 12.5, MUTED, False, PP_ALIGN.CENTER, margin=0.0)

    # 19 P-pick risk
    s = d.slide("P-pick 和数据质量风险", "P-pick 决定波形窗口、station SNR 和 feature/label 对齐", "Data Risk", RED)
    add_image_fit(s, path_asset("ppick_metadata_stats.png"), 0.55, 1.20, 6.00, 3.75)
    add_image_fit(s, path_asset("ppick_audit_waveforms.png"), 6.85, 1.20, 5.95, 3.75)
    add_bullets(s, [
        "当前流程：走时曲线粗定位 + STA/LTA 搜索窗 refine。",
        "低 SNR 或错误 pick 会污染 station feature 和 single-station prior。",
        "部分 p_picks 是 aligned/global sample，超出当前 100s 训练窗口，需回原始 aligned 记录抽查。",
    ], 0.75, 5.25, 11.7, 1.05, size=11.8, bullet_color=RED)

    # 20 caveats
    s = d.slide("当前主要问题与解释边界", "哪些结论可以说，哪些还不能说", "Caveats", ORANGE)
    caveats = [
        ("Overfit 诊断", "当前不是最终泛化评估；验证集小，需要多 seed / 多 split。", ORANGE),
        ("输出方差", "PGA pred std 仍小于 target std，模型偏保守。", RED),
        ("Event readout", "mag/loc 联合训练仍会常数化，需要 event-only / staged 诊断。", BLUE),
        ("抽样策略", "first/random/fixed/input target 不能混为一谈。", PURPLE),
    ]
    for i, (t, b, c) in enumerate(caveats):
        add_card(s, 0.75 + (i % 2) * 6.05, 1.35 + (i // 2) * 1.70, 5.55, 1.10, t, b, c, PANEL)
    add_image_fit(s, path_asset("variance_compression.png"), 2.10, 4.95, 9.15, 1.20)

    # 21 Next experiments
    s = d.slide("下一步实验路线图", "先做可信度闭环，再做结构消融，最后扩大训练", "Plan", GREEN)
    tracks = [
        ("P0 可信度", "P-pick 人工质检；低 SNR/失败类型统计；统一指标表。", RED),
        ("P1 Graph 消融", "baseline-only / residual-only / prior-only / prior+baseline+residual；distance power p；kNN。", GREEN),
        ("P2 Cross 增强", "relative coordinate MLP；distance bias；event context gate；target/input 分组评估。", BLUE),
        ("P3 Event 重建", "event-only cross-attention；mag/loc 与 PGA 分阶段训练；event token 作为 PGA context。", PURPLE),
    ]
    for i, (t, b, c) in enumerate(tracks):
        add_card(s, 0.90, 1.35 + i * 1.23, 11.55, 0.88, t, b, c, PANEL, title_size=12.5, body_size=10.5)
    add_pill(s, "阶段门槛：P-pick 质检通过 + prior-residual 消融清楚 + 多 seed/split 稳定后，再转向大训练集", 1.05, 6.25, 11.25, 0.36, PANEL_GREEN, GREEN)

    # 22 closing questions
    s = d.slide("希望专家重点建议", "欢迎围绕数据、先验和实验节奏直接指出风险", "Discussion", CYAN)
    questions = [
        "P-pick 的 STA/LTA refine 是否合理？哪些样本应该剔除或降权？",
        "PGA target 是否必须显式使用 event 信息？event 信息应如何进入模型？",
        "距离衰减先验更适合作为 baseline、attention bias，还是 residual feature？",
        "是否需要加入台站场地项、方位角、路径效应、区域衰减参数？",
        "当前 overfit 诊断做到哪一步可以转向更大训练集？",
    ]
    for i, q in enumerate(questions):
        add_card(s, 1.00, 1.30 + i * 0.88, 11.35, 0.58, q, "", CYAN if i % 2 == 0 else BLUE, PANEL, title_size=12.0, number=str(i + 1))
    add_shape(s, 1.00, 6.00, 11.35, 0.58, PANEL_BLUE, LINE, radius=True)
    add_textbox(s, "最终建议：先把数据质检和 prior-residual 消融做扎实，再扩大训练集验证 cross-attention / graph 融合模型的泛化。", 1.25, 6.16, 10.80, 0.24, 12.8, NAVY, True, PP_ALIGN.CENTER, margin=0.0)

    d.save(PPTX_PATH)
    print(f"PPTX: {PPTX_PATH}")


if __name__ == "__main__":
    main()
