from __future__ import annotations

import html
import json
import re
from pathlib import Path

import h5py
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import build_technical_exchange_ppt as base


OUT_DIR = base.OUT_DIR
ASSET_DIR = base.ASSET_DIR
TABLE_DIR = base.TABLE_DIR
HTML_PATH = OUT_DIR / "diting_team_graph_pga_technical_exchange_v4.html"
PPTX_PATH = OUT_DIR / "diting_team_graph_pga_technical_exchange_v4.pptx"

BLUE = RGBColor(31, 78, 121)
TEAL = RGBColor(0, 128, 128)
ORANGE = RGBColor(198, 89, 17)
GRAY = RGBColor(80, 80, 80)
LIGHT = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)


def savefig(path: Path) -> Path:
    plt.savefig(path, dpi=230, bbox_inches="tight")
    plt.close()
    return path


def plot_ppick_audit() -> Path:
    sample_rate = 100.0
    pre_s, post_s = 5.0, 15.0
    pre, post = int(pre_s * sample_rate), int(post_s * sample_rate)
    records = []
    with h5py.File(base.H5_PATH, "r") as f:
        for event_key in sorted(f["data"].keys()):
            g = f["data"][event_key]
            picks = np.asarray(g["p_picks"])
            pga = np.asarray(g["pga"])
            valid = np.where((picks >= pre + 100) & (picks < g["waveforms"].shape[1] - post - 100))[0]
            if valid.size == 0:
                continue
            order = valid[np.argsort(pga[valid])]
            chosen = np.unique(np.concatenate([order[:2], order[len(order) // 2:len(order) // 2 + 2], order[-2:]]))
            for station_idx in chosen:
                records.append((event_key, int(station_idx), int(picks[station_idx]), float(pga[station_idx])))
    if not records:
        raise RuntimeError("No P-pick-in-window samples found in japan_overfit.hdf5.")

    records = records[:12]
    fig, axes = plt.subplots(len(records), 1, figsize=(12.2, 9.2), sharex=True)
    if len(records) == 1:
        axes = [axes]
    colors = ["#4C78A8", "#F58518", "#54A24B"]
    names = ["EW", "NS", "UD"]
    with h5py.File(base.H5_PATH, "r") as f:
        for ax, (event_key, station_idx, pick, pga) in zip(axes, records):
            wave = np.asarray(f["data"][event_key]["waveforms"][station_idx])
            seg = wave[pick - pre:pick + post, :]
            t = (np.arange(seg.shape[0]) - pre) / sample_rate
            for c in range(3):
                trace = seg[:, c] - np.nanmedian(seg[:, c])
                scale = np.nanpercentile(np.abs(trace), 98)
                if not np.isfinite(scale) or scale <= 0:
                    scale = np.nanmax(np.abs(trace)) + 1e-8
                ax.plot(t, trace / scale + c * 1.45, lw=0.65, color=colors[c], label=names[c])
            ax.axvline(0, color="crimson", lw=1.2, ls="--")
            ax.axvspan(-0.5, 0.5, color="crimson", alpha=0.07, lw=0)
            ax.set_yticks([])
            ax.set_xlim(-pre_s, post_s)
            ax.text(
                -4.85,
                3.28,
                f"{event_key}  station {station_idx}  pick={pick}  PGA={pga:.3g}",
                fontsize=8.2,
                va="top",
                bbox=dict(facecolor="white", alpha=0.72, edgecolor="none", pad=1.0),
            )
            ax.grid(axis="x", color="#E6E6E6", lw=0.5)
    axes[0].legend(loc="upper right", ncol=3, frameon=False, fontsize=8)
    axes[-1].set_xlabel("Time relative to P pick (s)")
    fig.suptitle(
        "P-pick visual audit samples from japan_overfit.hdf5",
        fontsize=15,
        fontweight="bold",
    )
    fig.text(
        0.02,
        0.01,
        "Only records whose stored P pick falls inside the 100 s waveform window are shown; expert review should also inspect raw aligned records where picks are outside this cut window.",
        fontsize=9,
        color="#666666",
    )
    return savefig(ASSET_DIR / "ppick_audit_waveforms.png")


def experiment_label(name: str) -> str:
    label = name.replace("weights_japan_overfit_", "")
    replacements = {
        "pga15_cross_attention_": "pga15_cross_",
        "event_pga_cross_attention": "event+pga_cross",
        "exp1_": "exp1_",
        "graph_pga_": "graph_",
        "graph_event_pga_": "graph_event+pga_",
    }
    for old, new in replacements.items():
        label = label.replace(old, new)
    return label


def infer_overfit_n(name: str) -> str:
    m = re.search(r"overfit(\d+)", name)
    if m:
        return m.group(1)
    return "32-like" if "exp" in name or "first_inputs" in name else "unknown"


def infer_setting(name: str) -> str:
    if "fixed_inputs_random_targets" in name:
        return "fixed inputs; random targets"
    if "fixed_inputs_targets" in name:
        return "fixed inputs; fixed targets"
    if "input_targets" in name:
        return "input stations also used as targets"
    if "first_inputs" in name:
        return "first input stations"
    if "random_inputs" in name:
        return "random input stations"
    if "same_station" in name:
        return "single input; same-station target"
    if "multi_target" in name:
        return "single input; multiple targets"
    if "holdout" in name:
        return "SNR-filtered holdout targets"
    if "mask_batch1" in name:
        return "batch=1 mask sanity"
    if "query_no_transformer" in name:
        return "target query without transformer"
    if "direct_station" in name:
        return "pooled station embedding readout"
    if "attn_diag" in name:
        return "query_transformer + attention diagnostics"
    return "random inputs/targets"


def infer_route(name: str, cfg: dict, group: str) -> str:
    mode = cfg.get("model_params", {}).get("pga_readout_mode", "")
    event_mode = cfg.get("model_params", {}).get("event_readout_mode", "")
    if group == "Graph":
        if "prior_residual" in name:
            return "graph prior-residual"
        return "graph message passing"
    if mode == "target_cross_attention" and event_mode == "event_cross_attention":
        return "PGA+event cross-attn"
    if mode == "target_cross_attention":
        return "PGA cross-attn"
    if mode == "query_transformer":
        return "full self-attn query"
    if mode == "direct_station":
        return "direct station"
    if mode == "query_no_transformer":
        return "query only"
    return mode or "unknown"


def find_eval_file(exp_dir: Path) -> Path | None:
    files = sorted(exp_dir.glob("eval_results*.txt"))
    return files[0] if files else None


def build_all_experiment_table() -> pd.DataFrame:
    rows = []
    for group, root in [("TEAM", base.TEAM_RES), ("Graph", base.GRAPH_RES)]:
        for exp_dir in sorted([p for p in root.iterdir() if p.is_dir()]):
            cfg_files = sorted(exp_dir.glob("config*.json"))
            cfg = json.loads(cfg_files[0].read_text()) if cfg_files else {}
            mp = cfg.get("model_params", {})
            eval_file = find_eval_file(exp_dir)
            ev = base.extract_eval_metrics(eval_file) if eval_file else {}

            def g(split: str, task: str, key: str):
                return ev.get(split, {}).get(task, {}).get(key)

            rows.append(
                {
                    "Group": group,
                    "Experiment": experiment_label(exp_dir.name),
                    "Route": infer_route(exp_dir.name, cfg, group),
                    "Setting": infer_setting(exp_dir.name),
                    "overfit_n": infer_overfit_n(exp_dir.name),
                    "n_pga": mp.get("n_pga_targets"),
                    "Train MAE": g("train", "pga", "mae"),
                    "Train Corr": g("train", "pga", "corr"),
                    "Train R2": g("train", "pga", "r2"),
                    "Val MAE": g("val", "pga", "mae"),
                    "Val Corr": g("val", "pga", "corr"),
                    "Val R2": g("val", "pga", "r2"),
                    "Val slope": g("val", "pga", "slope"),
                    "Val mag MAE": g("val", "mag", "mae"),
                    "Status": "ok" if g("train", "pga", "mae") is not None or g("val", "pga", "mae") is not None else "eval failed/no PGA metrics",
                }
            )
    df = pd.DataFrame(rows)
    df.to_csv(TABLE_DIR / "all_experiments_summary.csv", index=False)
    return df


def plot_task_setup() -> Path:
    fig, ax = plt.subplots(figsize=(13, 6.2))
    ax.axis("off")
    ax.set_xlim(0, 13)
    ax.set_ylim(0, 6.2)

    _box(ax, 0.35, 4.65, 2.0, 0.75, "DiTing usual use\nsingle station", "#E7E6E6")
    _box(ax, 2.95, 4.65, 2.4, 0.75, "one station waveform\n3-component trace", "#E2F0D9")
    _box(ax, 6.0, 4.65, 2.2, 0.75, "DiTing encoder", "#DDEBF7")
    _box(ax, 8.9, 4.65, 3.4, 0.75, "single-station outputs\nphase pick / mag proxy / station tasks", "#FFF2CC", fs=8.5)
    _arrow(ax, 2.35, 5.02, 2.95, 5.02)
    _arrow(ax, 5.35, 5.02, 6.0, 5.02)
    _arrow(ax, 8.2, 5.02, 8.9, 5.02)

    _box(ax, 0.35, 2.45, 2.0, 0.85, "Current task\nmulti-station event", "#E7E6E6")
    _box(ax, 2.95, 2.45, 2.5, 0.85, "N station waveforms\n+ station coords\n+ station_valid", "#E2F0D9", fs=8.5)
    _box(ax, 6.0, 2.45, 2.3, 0.85, "shared DiTing encoder\nper station embedding", "#DDEBF7", fs=8.5)
    _box(ax, 8.95, 3.05, 3.55, 0.78, "event-level outputs\nmagnitude + location / distance", "#FFF2CC", fs=8.5)
    _box(ax, 8.95, 1.85, 3.55, 0.78, "target-station outputs\nPGA at 1 or 15 target sites", "#FCE4D6", fs=8.5)
    _arrow(ax, 2.35, 2.88, 2.95, 2.88)
    _arrow(ax, 5.45, 2.88, 6.0, 2.88)
    _arrow(ax, 8.3, 2.95, 8.95, 3.42)
    _arrow(ax, 8.3, 2.75, 8.95, 2.25)

    ax.text(0.35, 5.83, "What task are we doing?", fontsize=17, fontweight="bold", color="#1F4E79")
    ax.text(
        0.35,
        0.68,
        "Core difference from standard DiTing use: not one station -> one local prediction, but many stations -> event summary and spatial PGA field readout.",
        fontsize=11,
        color="#666666",
    )
    return savefig(ASSET_DIR / "task_setup_multistation.png")


def plot_ppick_stats() -> Path:
    df = pd.read_csv(base.DATA128 / "split_stations.csv")
    fig = plt.figure(figsize=(12, 6.1))
    gs = fig.add_gridspec(1, 3, width_ratios=[1.0, 1.0, 1.1])

    ax0 = fig.add_subplot(gs[0, 0])
    counts = df["p_pick_refined_source"].fillna("NA").value_counts()
    ax0.bar(counts.index.astype(str), counts.values, color=["#4C78A8", "#F58518", "#54A24B"][: len(counts)])
    ax0.set_title("Refined pick source")
    ax0.set_ylabel("Station records")
    ax0.tick_params(axis="x", rotation=20)
    for i, v in enumerate(counts.values):
        ax0.text(i, v, str(int(v)), ha="center", va="bottom", fontsize=9)

    ax1 = fig.add_subplot(gs[0, 1])
    delta = (df["p_pick_refined_aligned"] - df["p_pick_predicted_aligned"]) / df["sampling_rate_hz"].replace(0, np.nan)
    ax1.hist(delta.clip(-20, 20), bins=45, color="#59A14F", alpha=0.85)
    ax1.axvline(0, color="black", lw=1)
    ax1.set_title("Refined pick - travel-time pick")
    ax1.set_xlabel("Seconds, clipped to [-20, 20]")
    ax1.set_ylabel("Count")

    ax2 = fig.add_subplot(gs[0, 2])
    sc = ax2.scatter(
        df["hypocentral_distance_km"],
        df["stalta_ratio_at_pick"],
        c=np.log10(df["pga_norm_resampled_mps2"].clip(lower=1e-4)),
        s=7,
        alpha=0.35,
        cmap="viridis",
        linewidths=0,
    )
    ax2.set_title("STA/LTA ratio at final pick")
    ax2.set_xlabel("Hypocentral distance (km)")
    ax2.set_ylabel("STA/LTA ratio")
    cb = fig.colorbar(sc, ax=ax2, fraction=0.045, pad=0.02)
    cb.set_label("log10(PGA)")

    fig.suptitle("P-pick metadata diagnostics for expert review", fontsize=15, fontweight="bold")
    return savefig(ASSET_DIR / "ppick_metadata_stats.png")


def _box(ax, x, y, w, h, text, fc="#EAF2F8", ec="#1F4E79", fs=9):
    rect = plt.Rectangle((x, y), w, h, facecolor=fc, edgecolor=ec, linewidth=1.3)
    ax.add_patch(rect)
    ax.text(x + w / 2, y + h / 2, text, ha="center", va="center", fontsize=fs, wrap=True)


def _arrow(ax, x1, y1, x2, y2, color="#555555"):
    ax.annotate("", xy=(x2, y2), xytext=(x1, y1), arrowprops=dict(arrowstyle="->", lw=1.25, color=color))


def plot_architecture_detail() -> Path:
    fig, ax = plt.subplots(figsize=(14, 7.2))
    ax.axis("off")
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 7.2)

    _box(ax, 0.3, 5.75, 2.0, 0.72, "waveform\n(B,S,3,10000)", "#E2F0D9")
    _box(ax, 2.7, 5.75, 2.0, 0.72, "std normalize\nstation_valid mask", "#E2F0D9")
    _box(ax, 5.1, 5.75, 2.1, 0.72, "DiTing encoder\n+ station adapter", "#E7E6E6")
    _box(ax, 7.7, 5.75, 2.2, 0.72, "raw_station_emb\n(B,S,D)", "#DDEBF7")
    _box(ax, 10.4, 5.75, 2.5, 0.72, "LayerNorm + scale gate\nlog std/rms/peak", "#DDEBF7")
    for x in [2.3, 4.7, 7.2, 9.9]:
        _arrow(ax, x, 6.11, x + 0.4, 6.11)

    _box(ax, 0.8, 4.55, 2.7, 0.72, "station coords\nabs/relative to valid-center", "#FFF2CC")
    _box(ax, 4.1, 4.55, 2.6, 0.72, "sin/cos position emb\nor concat coord MLP", "#FFF2CC")
    _box(ax, 7.3, 4.55, 2.8, 0.72, "coord fusion\nadd or concat+proj+norm", "#FFF2CC")
    _box(ax, 10.8, 4.55, 2.5, 0.72, "station_feature_emb\n(B,S,D)", "#D9EAD3")
    for x in [3.5, 6.7, 10.1]:
        _arrow(ax, x, 4.91, x + 0.6, 4.91)
    _arrow(ax, 11.65, 5.75, 11.65, 5.27)

    _box(ax, 0.35, 2.8, 3.25, 1.05, "Original query_transformer\nconcat [event, station, PGA]\npadding_mask=[1, station_valid, pga_valid]\natt_mask keys=[event, station, not PGA]", "#FCE4D6", fs=7.5)
    _box(ax, 4.05, 2.8, 3.25, 1.05, "PGA cross-attention\nQ=target_pos_emb + learned PGA token\nK/V=station_feature_emb\nkey_padding_mask=~station_valid", "#E2F0D9", fs=7.5)
    _box(ax, 7.75, 2.8, 2.95, 1.05, "Event cross-attention\nQ=learned event token\nK/V=station_feature_emb\nmag/loc heads", "#E2F0D9", fs=7.5)
    _box(ax, 11.15, 2.8, 2.45, 1.05, "PGA output head\nMLP + point/MDN\nvalid targets only", "#DDEBF7", fs=7.6)
    _arrow(ax, 12.05, 4.55, 1.95, 3.85)
    _arrow(ax, 12.05, 4.55, 5.65, 3.85)
    _arrow(ax, 12.05, 4.55, 9.2, 3.85)
    _arrow(ax, 7.3, 3.32, 11.15, 3.32)

    _box(ax, 1.35, 1.15, 3.2, 0.9, "Graph route\nstation-to-target messages", "#FFF2CC", fs=8.0)
    _box(ax, 5.35, 1.15, 3.2, 0.9, "single-station PGA prior\nloaded into station_pga_prior_head", "#D9EAD3", fs=8.0)
    _box(ax, 9.35, 1.15, 3.2, 0.9, "final PGA\n= distance baseline + residual", "#DDEBF7", fs=8.0)
    _arrow(ax, 4.55, 1.6, 5.35, 1.6)
    _arrow(ax, 8.55, 1.6, 9.35, 1.6)

    ax.text(0.3, 6.85, "Implementation-level model flow", fontsize=17, fontweight="bold", color="#1F4E79")
    ax.text(0.3, 0.38, "Main diagnosis target: station_feature_emb is usable, but the readout route and inductive bias determine whether outputs collapse to constants.", fontsize=11, color="#666666")
    return savefig(ASSET_DIR / "architecture_implementation_detail.png")


def plot_graph_detail() -> Path:
    fig, ax = plt.subplots(figsize=(13, 6.5))
    ax.axis("off")
    ax.set_xlim(0, 13)
    ax.set_ylim(0, 6.5)
    _box(ax, 0.4, 4.8, 2.2, 0.75, "station_feature_emb\n(B,S,D)", "#DDEBF7")
    _box(ax, 0.4, 3.55, 2.2, 0.75, "target position emb\n(B,T,D)", "#FFF2CC")
    _box(ax, 3.05, 4.65, 2.45, 0.9, "station_pga_prior_head\nprior_i from station token", "#D9EAD3", fs=8.5)
    _box(ax, 3.05, 3.35, 2.45, 0.9, "edge feature e_ti\nrel xyz, dist,\nlog dist, inv dist", "#FFF2CC", fs=8.5)
    _box(ax, 6.0, 3.8, 2.5, 1.1, "score_ti = MLP(\ntarget_t, station_i, edge_ti)\nmasked softmax over valid stations", "#FCE4D6", fs=8.0)
    _box(ax, 9.0, 4.15, 1.6, 0.75, "message_t\nsum_i w_ti value_i", "#FCE4D6", fs=8.2)
    _box(ax, 11.0, 4.15, 1.55, 0.75, "residual_t\nPGA head", "#DDEBF7")
    _box(ax, 6.25, 1.65, 2.2, 0.85, "distance baseline\nsum_i dist^-p * prior_i", "#D9EAD3", fs=8.5)
    _box(ax, 9.35, 1.65, 2.6, 0.85, "PGA_t = baseline_t + residual_t\nvalid target mask applied", "#DDEBF7", fs=8.5)
    for a in [(2.6, 5.18, 3.05, 5.1), (2.6, 3.93, 3.05, 3.8), (5.5, 5.1, 6.0, 4.55), (5.5, 3.8, 6.0, 4.25), (8.5, 4.35, 9.0, 4.52), (10.6, 4.52, 11.0, 4.52), (5.5, 5.0, 6.25, 2.05), (8.45, 2.08, 9.35, 2.08), (11.0, 4.15, 10.7, 2.5)]:
        _arrow(ax, *a)
    ax.text(0.4, 5.98, "Graph route implementation detail", fontsize=16, fontweight="bold", color="#008080")
    ax.text(0.4, 0.65, "This is why prior-residual helped: the model no longer has to rediscover the single-station PGA scale and distance decay from scratch.", fontsize=11, color="#666666")
    return savefig(ASSET_DIR / "architecture_graph_detail.png")


def plot_train_val_corr(df: pd.DataFrame, group: str, filename: str, title: str, max_rows: int | None = None) -> Path:
    data = df[(df["Group"] == group) & df["Train Corr"].notna()].copy()
    if max_rows is not None:
        data = data.head(max_rows)
    data["Display"] = data["Experiment"].str.replace("_", "\n")
    height = max(5.8, 0.34 * len(data) + 1.5)
    fig, ax = plt.subplots(figsize=(12, height))
    y = np.arange(len(data))
    ax.barh(y - 0.18, data["Train Corr"], height=0.34, label="Train Corr", color="#4C78A8")
    ax.barh(y + 0.18, data["Val Corr"], height=0.34, label="Val Corr", color="#F58518")
    ax.axvline(0, color="#333333", lw=0.8)
    ax.set_yticks(y)
    ax.set_yticklabels(data["Display"], fontsize=8)
    ax.invert_yaxis()
    ax.set_xlabel("PGA correlation")
    ax.set_title(title)
    ax.legend(frameon=False, ncol=2, loc="lower right")
    ax.grid(axis="x", color="#E6E6E6", lw=0.6)
    return savefig(ASSET_DIR / filename)


def plot_train_val_r2_mae(df: pd.DataFrame, experiments: list[str], filename: str, title: str) -> Path:
    data = df[df["Experiment"].isin(experiments)].copy()
    order = {name: i for i, name in enumerate(experiments)}
    data = data.sort_values("Experiment", key=lambda s: s.map(order))
    labels = data["Experiment"].str.replace("_", "\n")
    x = np.arange(len(data))
    fig, axes = plt.subplots(1, 2, figsize=(12, 4.8))
    axes[0].bar(x - 0.18, data["Train Corr"], width=0.36, label="Train", color="#4C78A8")
    axes[0].bar(x + 0.18, data["Val Corr"], width=0.36, label="Val", color="#F58518")
    axes[0].axhline(0, color="black", lw=0.8)
    axes[0].set_title("PGA Corr")
    axes[0].set_xticks(x)
    axes[0].set_xticklabels(labels, rotation=18, ha="right", fontsize=8)
    axes[0].legend(frameon=False)
    axes[1].bar(x - 0.18, data["Train MAE"], width=0.36, label="Train", color="#4C78A8")
    axes[1].bar(x + 0.18, data["Val MAE"], width=0.36, label="Val", color="#F58518")
    axes[1].set_title("PGA MAE")
    axes[1].set_xticks(x)
    axes[1].set_xticklabels(labels, rotation=18, ha="right", fontsize=8)
    fig.suptitle(title, fontsize=14, fontweight="bold")
    return savefig(ASSET_DIR / filename)


def plot_all_experiment_matrix(df: pd.DataFrame) -> Path:
    data = df.copy()
    data["Display"] = data["Group"] + ": " + data["Experiment"]
    data = data.sort_values(["Group", "Route", "Experiment"])
    fig, ax = plt.subplots(figsize=(12.8, 8.2))
    y = np.arange(len(data))
    colors = np.where(data["Status"].eq("ok"), "#4C78A8", "#B0B0B0")
    ax.scatter(data["Train Corr"], y, s=54, marker="o", color=colors, label="Train Corr")
    ax.scatter(data["Val Corr"], y, s=54, marker="s", color="#F58518", label="Val Corr")
    ax.axvline(0, color="#333333", lw=0.8)
    ax.set_yticks(y)
    ax.set_yticklabels(data["Display"], fontsize=7.5)
    ax.invert_yaxis()
    ax.set_xlabel("PGA correlation")
    ax.set_title("All experiments in chaosuan_res and team_pytorch2/chaosuan_res")
    ax.grid(axis="x", color="#E6E6E6", lw=0.6)
    ax.legend(frameon=False, loc="lower right")
    return savefig(ASSET_DIR / "all_experiments_train_val_corr.png")


def setting_table(kind: str) -> pd.DataFrame:
    if kind == "ablation":
        rows = [
            ["query_transformer", "PGA target token enters full Transformer with event/station tokens", "tests original readout; tends to constant"],
            ["mask_batch1", "same as query_transformer but batch size 1 for mask sanity", "checks cross-batch leakage/mask issue"],
            ["query_no_transformer", "target coordinate + learned query token only, no station readout", "control for target-position-only shortcut"],
            ["direct_station", "masked mean station_feature_emb -> PGA head", "checks whether station feature/head contains signal"],
            ["target_cross_attention", "PGA target query cross-attends station_feature_emb", "intended fix; earlier eval failed/no metrics"],
        ]
    elif kind == "cross":
        rows = [
            ["pga15_cross_overfit32", "overfit_n=32; 15 PGA targets; random inputs/targets", "baseline cross-attn small sample"],
            ["pga15_cross_overfit128", "overfit_n=128; 15 PGA targets; random inputs/targets", "larger overfit sample"],
            ["fixed_inputs_targets", "overfit_n=32; same input stations and same targets", "easiest memorization/route sanity"],
            ["input_targets", "overfit_n=32; input stations also used as targets", "same-station PGA readout sanity"],
            ["fixed_inputs_random_targets", "overfit_n=32; fixed inputs, random targets", "tests spatial interpolation beyond fixed targets"],
            ["event+pga_cross_first_inputs", "event and PGA both cross-attend; first input stations", "joint event/PGA readout test"],
            ["mag/loc_cross_overfit32", "event task emphasized with cross-attention", "checks event readout collapse"],
        ]
    else:
        rows = [
            ["graph first-inputs", "GraphPGAReadout without station prior/distance baseline", "old graph; still collapses"],
            ["graph event+pga first-inputs", "graph PGA plus event task", "tests joint event/PGA route"],
            ["prior-residual first-inputs", "single-station PGA prior + distance baseline + residual", "best current graph setting"],
            ["prior-residual random-inputs", "same prior-residual, random input stations", "harder input selection"],
            ["exp1 same-station", "single input and same-station target", "old graph sanity"],
            ["exp2 multi-target", "single input and multiple targets", "old graph spatial readout"],
            ["exp3 holdout", "SNR-filtered holdout targets", "old graph harder holdout task"],
        ]
    return pd.DataFrame(rows, columns=["Model", "Setting", "Purpose"])


def build_model_scale_table() -> pd.DataFrame:
    rows = [
        [
            "DiTing encoder",
            "1,234.0M",
            "0",
            "Frozen",
            "Exact count from local TEAM cross-attention instantiation.",
        ],
        [
            "Station adapter",
            "3.59M",
            "3.59M",
            "Trainable",
            "Adapts frozen waveform encoder output to station embedding.",
        ],
        [
            "TEAM cross-attn full model",
            "1,262.0M",
            "28.0M",
            "Trainable heads/readout only",
            "Exact count for pga15_cross_overfit128 config.",
        ],
        [
            "Graph prior-residual extra",
            "~10.0M",
            "~10.0M",
            "Trainable",
            "Analytic estimate: GraphPGAReadout ~9.02M + prior head ~1.00M.",
        ],
        [
            "Graph prior-residual full model",
            "~1,272M",
            "~38M",
            "Approx.",
            "Exact local instantiation blocked by missing dtbench in team_pytorch2 env.",
        ],
    ]
    df = pd.DataFrame(rows, columns=["Component", "Total params", "Trainable params", "Status", "Note"])
    df.to_csv(TABLE_DIR / "model_scale_params.csv", index=False)
    return df


def build_loss_design_table() -> pd.DataFrame:
    rows = [
        [
            "single-station pretrain",
            "mag, epidist, pga",
            "Huber(delta=1)",
            "0.3 / 0.3 / 0.4",
            "Freeze encoder; train adapter + single-station heads.",
        ],
        [
            "full PGA-only runs",
            "pga",
            "Huber(delta=1)",
            "1.0",
            "Loss only on valid PGA targets.",
        ],
        [
            "full event+PGA runs",
            "mag, loc, pga",
            "Huber(delta=1)",
            "0.2 / 0.2 / 1.0",
            "Joint readout test; mag/loc collapse is part of diagnosis.",
        ],
        [
            "graph prior-residual runs",
            "pga",
            "Huber(delta=1)",
            "1.0",
            "Final PGA = distance baseline + learned residual.",
        ],
    ]
    df = pd.DataFrame(rows, columns=["Stage", "Tasks", "Loss", "Weights", "Implementation note"])
    df.to_csv(TABLE_DIR / "loss_design.csv", index=False)
    return df


def _single_station_counts(path: Path) -> dict[str, int]:
    counts: dict[str, int] = {}
    if not path.exists():
        return counts
    for line in path.read_text(errors="ignore").splitlines():
        m = re.search(r"SINGLE-STATION\s+(TRAIN|VAL)\s+set:\s+(\d+)\s+station samples", line)
        if m:
            counts[m.group(1).lower()] = int(m.group(2))
    return counts


def build_single_station_table() -> pd.DataFrame:
    experiments = [
        ("Cross overfit128", base.TEAM_RES / "weights_japan_overfit_pga15_cross_attention_overfit128"),
        ("Cross fixed inputs/targets", base.TEAM_RES / "weights_japan_overfit_pga15_cross_attention_overfit32_fixed_inputs_targets"),
        ("Graph prior first-inputs", base.GRAPH_RES / "weights_japan_overfit_graph_pga_first_inputs_prior_residual"),
        ("Graph prior random-inputs", base.GRAPH_RES / "weights_japan_overfit_graph_pga_random_inputs_prior_residual"),
    ]
    rows = []
    task_names = [("single_mag", "mag"), ("single_epidist", "epidist"), ("single_pga", "pga")]
    for label, exp_dir in experiments:
        eval_file = find_eval_file(exp_dir)
        if eval_file is None:
            continue
        ev = base.extract_eval_metrics(eval_file)
        counts = _single_station_counts(eval_file)
        for split_key, split_name in [("single_train", "train"), ("single_val", "val")]:
            for task_key, task_name in task_names:
                m = ev.get(split_key, {}).get(task_key, {})
                if not m:
                    continue
                rows.append(
                    {
                        "Model": label,
                        "Split": split_name,
                        "Task": task_name,
                        "Samples": counts.get(split_name),
                        "MAE": m.get("mae"),
                        "Corr": m.get("corr"),
                    }
                )
    df = pd.DataFrame(rows)
    df.to_csv(TABLE_DIR / "single_station_metrics.csv", index=False)
    return df


def plot_single_station_metrics(df: pd.DataFrame) -> Path:
    data = df[df["Split"].eq("val")].copy()
    task_order = {"mag": 0, "epidist": 1, "pga": 2}
    data["_task_order"] = data["Task"].map(task_order).fillna(99)
    data = data.sort_values(["Model", "_task_order"])
    fig, axes = plt.subplots(1, 2, figsize=(12.5, 5.2), sharey=True)
    labels = data["Model"] + "\n" + data["Task"]
    y = np.arange(len(data))
    axes[0].barh(y, data["Corr"], color="#4C78A8")
    axes[0].set_title("Single-station val Corr")
    axes[0].set_xlabel("Correlation")
    axes[0].axvline(0, color="#333333", lw=0.8)
    axes[0].set_yticks(y)
    axes[0].set_yticklabels(labels, fontsize=8)
    axes[0].invert_yaxis()
    axes[1].barh(y, data["MAE"], color="#F58518")
    axes[1].set_title("Single-station val MAE")
    axes[1].set_xlabel("MAE")
    axes[1].grid(axis="x", color="#E6E6E6", lw=0.6)
    axes[0].grid(axis="x", color="#E6E6E6", lw=0.6)
    fig.suptitle("Single-station pretrain/readout sanity: waveform station features carry signal", fontsize=14, fontweight="bold")
    return savefig(ASSET_DIR / "single_station_val_metrics.png")


def _read_loss_csv(path: Path) -> pd.DataFrame | None:
    if not path.exists() or path.stat().st_size == 0:
        return None
    df = pd.read_csv(path)
    if not {"step", "value"}.issubset(df.columns):
        return None
    df = df[["step", "value"]].copy()
    df["value"] = pd.to_numeric(df["value"], errors="coerce")
    return df.dropna()


def plot_representative_loss_curves() -> Path:
    runs = [
        ("query_transformer", base.TEAM_RES / "weights_japan_overfit_exp1_attn_diag", "#7F7F7F"),
        ("cross overfit32", base.TEAM_RES / "weights_japan_overfit_pga15_cross_attention_overfit32", "#4C78A8"),
        ("cross overfit128", base.TEAM_RES / "weights_japan_overfit_pga15_cross_attention_overfit128", "#F58518"),
        ("graph prior first", base.GRAPH_RES / "weights_japan_overfit_graph_pga_first_inputs_prior_residual", "#54A24B"),
        ("graph prior random", base.GRAPH_RES / "weights_japan_overfit_graph_pga_random_inputs_prior_residual", "#B279A2"),
    ]
    fig, axes = plt.subplots(1, 2, figsize=(12.5, 4.9), sharey=False)
    for label, exp_dir, color in runs:
        train = _read_loss_csv(exp_dir / "train_epoch_loss.csv")
        val = _read_loss_csv(exp_dir / "val_epoch_loss.csv")
        if train is not None:
            axes[0].plot(train["step"], train["value"], label=label, color=color, lw=1.5)
        if val is not None:
            axes[1].plot(val["step"], val["value"], label=label, color=color, lw=1.5)
    axes[0].set_title("Full-model train loss")
    axes[1].set_title("Full-model val loss")
    for ax in axes:
        ax.set_xlabel("Epoch")
        ax.set_ylabel("Huber loss")
        ax.grid(color="#E6E6E6", lw=0.6)
        ax.legend(frameon=False, fontsize=8)
    fig.suptitle("Representative full-model loss curves: optimization generally converges", fontsize=14, fontweight="bold")
    return savefig(ASSET_DIR / "representative_full_loss_curves.png")


def plot_single_station_loss_curves() -> Path:
    runs = [
        ("cross overfit128", base.TEAM_RES / "weights_japan_overfit_pga15_cross_attention_overfit128", "#4C78A8"),
        ("graph prior first", base.GRAPH_RES / "weights_japan_overfit_graph_pga_first_inputs_prior_residual", "#54A24B"),
        ("graph prior random", base.GRAPH_RES / "weights_japan_overfit_graph_pga_random_inputs_prior_residual", "#B279A2"),
    ]
    fig, axes = plt.subplots(1, 2, figsize=(12.5, 4.9), sharey=False)
    for label, exp_dir, color in runs:
        train = _read_loss_csv(exp_dir / "single_station" / "single_station_train_epoch_loss.csv")
        val = _read_loss_csv(exp_dir / "single_station" / "single_station_val_epoch_loss.csv")
        if train is not None:
            axes[0].plot(train["step"], train["value"], label=label, color=color, lw=1.5)
        if val is not None:
            axes[1].plot(val["step"], val["value"], label=label, color=color, lw=1.5)
    axes[0].set_title("Single-station train loss")
    axes[1].set_title("Single-station val loss")
    for ax in axes:
        ax.set_xlabel("Epoch")
        ax.set_ylabel("Weighted Huber loss")
        ax.grid(color="#E6E6E6", lw=0.6)
        ax.legend(frameon=False, fontsize=8)
    fig.suptitle("Single-station pretrain convergence", fontsize=14, fontweight="bold")
    return savefig(ASSET_DIR / "single_station_loss_curves.png")


def rel(path: Path) -> str:
    return html.escape(str(path.relative_to(OUT_DIR)))


def card(title: str, body: str) -> str:
    return f"<div class='card'><h3>{html.escape(title)}</h3><p>{html.escape(body)}</p></div>"


def bullets(items: list[str]) -> str:
    return "<ul>" + "".join(f"<li>{html.escape(item)}</li>" for item in items) + "</ul>"


def img(path: Path, klass: str = "") -> str:
    return f"<img class='{klass}' src='{rel(path)}' />"


def table_html(df: pd.DataFrame, cols: list[str] | None = None, rows: int | None = None) -> str:
    if cols is not None:
        df = df[cols]
    if rows is not None:
        df = df.head(rows)
    df = df.copy()
    for c in df.columns:
        if pd.api.types.is_float_dtype(df[c]):
            df[c] = df[c].map(lambda x: "" if pd.isna(x) else f"{x:.3f}")
    return df.to_html(index=False, classes="dataframe", border=0, escape=True)


def write_html_deck(assets: dict[str, Path], tables: dict[str, pd.DataFrame]) -> Path:
    slides = []

    def slide(title: str, body: str, subtitle: str = ""):
        slides.append(
            f"<section class='slide'><header><h1>{html.escape(title)}</h1>"
            f"{f'<p>{html.escape(subtitle)}</p>' if subtitle else ''}</header>{body}</section>"
        )

    slide(
        "DiTing 表征接入多台站 PGA / Event 预测",
        "<div class='hero'>" + card("交流目标", "说明做了什么、结果如何、问题在哪里，并请专家对数据、结构和先验设计提意见。") + card("主线判断", "station feature 有信号；坍塌主要发生在 full-model readout 和空间先验不足。") + "</div>",
        "Transformer readout 与 graph prior-residual 的实验诊断",
    )
    slide(
        "这次希望专家帮忙看的点",
        "<div class='grid2'><div>" + bullets([
            "P-pick 是否足够准确，是否会污染 station feature 和 PGA label 对齐。",
            "event 信息如何进入 PGA：event query、mag/loc context、GMPE prior，还是显式震源参数。",
            "任意 target PGA 更适合 cross-attention readout，还是 graph message passing。",
            "小样本 overfit 诊断后，下一步验证集和消融实验如何设计。"]) + "</div><div>" + card("特别新增", "P-pick 波形抽样页和实现级模型结构图，方便现场诊断。") + "</div></div>",
    )
    slide("当前任务：从单台站 DiTing 到多台站事件/PGA 预测", img(assets["task_setup"], "full"))
    slide(
        "任务设置",
        "<div class='grid2'><div>" + bullets([
            "输入：N 个台站的三分量波形、台站坐标、station_valid mask。",
            "目标查询：event query 用于震级/位置；PGA target query 用于给定 target station 的 PGA。",
            "输出：event-level magnitude / location 或震中距；target-level PGA。",
            "核心区别：不是单个台站独立预测，而是多个台站共同约束一个事件和一个空间 PGA 场。"]) + "</div><div>" + img(assets["task_setup"]) + "</div></div>",
    )
    slide("数据覆盖与抽样样本", "<div class='grid2 wideimg'><div>" + img(assets["data_overview"]) + "</div><div>" + img(assets["sample_event"]) + "</div></div>")
    slide("P-pick 生成流程与元数据诊断", "<div class='grid2'><div>" + bullets([
        "当前 P-pick：先用走时曲线给粗定位，再用 STA/LTA 在搜索窗内 refine。",
        "左图：最终 pick 的来源，主要是 STA/LTA threshold，少量退化为 STA/LTA argmax。",
        "中图：STA/LTA refine 后 pick 与走时粗定位 pick 的时间差，检查系统性早/晚偏差。",
        "右图：P pick 处 STA/LTA ratio 随震中距和 PGA 的分布，检查低信噪比远台站风险。"]) + "</div><div>" + img(assets["ppick_stats"]) + "</div></div>")
    slide("P-pick 波形抽样：请专家现场看", img(assets["ppick_audit"], "full") + "<p class='note'>红虚线为当前存储 P-pick。这里只展示 pick 落在 100s HDF5 窗口内的样本；pick 在固定训练窗口外的记录需要回到原始对齐数据继续抽查。</p>")
    slide("实现级模型总览：和后面表格的对应关系", "<div class='grid2'><div>" + img(assets["architecture_detail"]) + "</div><div>" + bullets([
        "表格中的 Route 对应底部三条 readout 分支：query_transformer、cross-attention、graph prior-residual。",
        "所有路线共用上半部分 station_feature_emb：DiTing encoder、station adapter、amplitude scale gate、coordinate fusion。",
        "后续结果表按 readout 分支分组；同一分支内再比较 input/target 选择、overfit_n、是否有 prior/baseline。"]) + "</div></div>")
    slide("模型规模与可训练参数", table_html(tables["model_scale"]) + "<p class='note'>关键点：总参数量主要来自 frozen DiTing encoder；当前 overfit 诊断实际训练的是 adapter、readout 和 task heads。Graph 参数量为实现结构的近似拆分，TEAM cross-attention 为本地精确计数。</p>")
    slide("Loss 函数设计", "<div class='grid2'><div>" + table_html(tables["loss_design"]) + "</div><div>" + bullets([
        "single-station 阶段：L = 0.3 L_mag + 0.3 L_epidist + 0.4 L_pga。",
        "PGA-only full model：只在有效 target PGA 上计算 Huber loss。",
        "event+PGA full model：L = 0.2 L_mag + 0.2 L_loc + 1.0 L_pga。",
        "这些实验使用 point prediction；收敛曲线下降不等于 readout 不坍塌，仍需看 corr、slope 和 pred std。"]) + "</div></div>")
    slide("Single-station model 效果：waveform 表征有效", "<div class='grid2'><div>" + img(assets["single_station_metrics"]) + "</div><div>" + table_html(tables["single_station"][tables["single_station"]["Split"].eq("val")], ["Model", "Task", "Samples", "MAE", "Corr"], 12) + "<p class='note'>single-station mag / epidist / PGA 在 val 上均有明显相关性，说明 waveform encoder + station adapter 不是完全无效；full model 的常数化更像 readout/空间传播问题。</p></div></div>")
    slide("全部实验总览：不是只做了少数代表实验", "<div class='grid2'><div>" + img(assets["all_experiments"]) + "</div><div>" + bullets([
        "TEAM 路线当前有 14 个实验；graph 路线当前有 7 个实验。",
        "图中圆点为 train PGA corr，方块为 val PGA corr。",
        "灰色表示 eval 失败或没有 PGA metrics，不能作为性能结论。",
        "后续主线页只展示代表实验，完整表格在 backup/CSV。"]) + "</div></div>")
    slide("原始 TEAM / query_transformer readout", "<div class='grid2'><div>" + img(assets["architecture_team"]) + "</div><div>" + bullets([
        "PGA query = target coordinate embedding + learned PGA query token。",
        "concat 后进入 full Transformer；PGA key 被 mask，PGA 作为 query-only token。",
        "event token 与 station token 共用 self-attention readout。",
        "风险：读出路径要靠训练自己发现，容易走向均值解。"]) + "</div></div>")
    slide("原始 full self-attention readout：几个 ablation 是什么", table_html(tables["ablation_settings"]) + "<p class='note'>这些实验共同定位：是 readout/路由问题，还是 station feature/PGA head 本身无效。</p>")
    slide("原始 full self-attention readout 的坍塌", "<div class='grid2'><div>" + img(assets["team_train_val"]) + "</div><div>" + table_html(tables["transformer_ablation"], ["Method", "Train Corr", "Train slope", "Val Corr", "Val slope"]) + "</div></div>")
    slide("Attention / mask 诊断结论", "<div class='grid2'><div>" + img(assets["variance"]) + "</div><div>" + bullets([
        "query_transformer：PGA token 与 event/station 一起进 full self-attention，是原始方案。",
        "mask_batch1：只改 batch=1 和 mask sanity，仍常数，跨 batch 泄露不是主因。",
        "query_no_transformer：只有 target query，没有 station value，作为负对照。",
        "direct_station：直接读 pooled station feature，可检验 station embedding 和 head 是否有信号。",
        "PGA query 后期可以 attend 到 station，但输出方差仍被压扁。",
        "0.11111 诊断 bug 已修：区分 key-type mask 与 padding 后 effective key mask。"]) + "</div></div>")
    slide("Cross-attention readout 设计用意", "<div class='grid2'><div>" + img(assets["architecture_cross"]) + "</div><div>" + bullets([
        "设计目标：让 target/event query 单向读取 station tokens，避免 full self-attention 中 event/PGA token 互相干扰。",
        "PGA Q：target coordinate embedding + learned PGA token。",
        "K/V：station_feature_emb；key_padding_mask=~station_valid。",
        "输出不加 query residual，使输出值更直接来自 station tokens。",
        "可选 distance_bias：由 target-station rel xyz + distance 产生 attention bias。"]) + "</div></div>")
    slide("Cross-attention 各实验设置", table_html(tables["cross_settings"]) + "<p class='note'>overfit_n=32/128 的区别在样本事件数；fixed/input/random target 用来拆分记忆、同台站读出和空间泛化。</p>")
    slide("Cross-attention 实验结果：train 和 val 都要看", "<div class='grid2'><div>" + img(assets["cross_train_val_detail"]) + "</div><div>" + table_html(tables["all_experiments"][tables["all_experiments"]["Route"].str.contains("cross", case=False, na=False)], ["Experiment", "overfit_n", "Setting", "Train Corr", "Val Corr", "Train MAE", "Val MAE"], 9) + "</div></div>")
    slide("Graph prior-residual：实现细节与为什么有效", "<div class='grid2'><div>" + img(assets["graph_detail"]) + "</div><div>" + bullets([
        "station_pga_prior_head 从 single-station 预训练加载，提供输入台站 PGA scale。",
        "distance baseline 用 dist^-p 对输入台站 prior 加权，给出物理上更合理的初值。",
        "graph message passing 只学习 residual，降低从零学习空间场的难度。",
        "旧 graph 仍坍塌，说明“有 graph”本身不够，关键是 prior + baseline + residual 路径。"]) + "</div></div>")
    slide("Graph 路线各实验设置", table_html(tables["graph_settings"]) + "<p class='note'>prior-residual 和 old graph 的区别不是“有没有 graph”，而是是否显式接入 single-station PGA prior 与距离 baseline。</p>")
    slide("Graph 路线实验结果：train 和 val 都要看", "<div class='grid2'><div>" + img(assets["graph_train_val_detail"]) + "</div><div>" + table_html(tables["all_experiments"][tables["all_experiments"]["Group"].eq("Graph")], ["Experiment", "Setting", "Train Corr", "Val Corr", "Train MAE", "Val MAE"], 7) + "</div></div>")
    slide("Loss 曲线与收敛情况", "<div class='grid2'><div>" + img(assets["full_loss_curves"]) + "</div><div>" + bullets([
        "代表 full-model 实验 train/val loss 基本随 epoch 下降，优化过程没有明显发散。",
        "query_transformer 的 loss 也会下降，但输出仍可坍塌到均值解。",
        "因此收敛页要和 corr、slope、pred std 一起看：loss 下降只能说明训练目标被优化，不能单独证明空间 readout 学对了。"]) + "</div></div>")
    slide("Single-station loss 曲线", "<div class='grid2'><div>" + img(assets["single_station_loss_curves"]) + "</div><div>" + bullets([
        "single-station 预训练的 train/val loss 稳定下降。",
        "这支持先把单台站 waveform 表征和 PGA prior 训练好，再把它作为 graph prior-residual 的输入。",
        "接下来需要保存并汇报分任务 loss，避免 mag/epidist/PGA 之间被加权总 loss 掩盖。"]) + "</div></div>")
    slide("两条路线对比", "<div class='cards4'>" + card("query_transformer", "统一 token 模型，但 PGA/event 都容易常数化。") + card("direct_station", "证明 station feature/head 可用，但不能处理任意 target 空间传播。") + card("cross-attention", "target 显式读 station，适合诊断 readout。") + card("graph prior-residual", "引入 single-station prior 和距离传播，当前结果最好。") + "</div>")
    slide("当前主要问题", "<div class='grid2'><div>" + bullets([
        "P-pick 准确性还没有人工抽查闭环，可能影响波形窗口和 station SNR。",
        "event mag/loc 联合训练仍会常数化，需要独立定位 event readout。",
        "PGA pred std 仍小于 target std，模型偏保守。",
        "验证集太小，多 seed / 多 split 后结论才稳。"]) + "</div><div>" + img(assets["variance"]) + "</div></div>")
    slide("下一步实验建议", "<div class='grid2'><div>" + bullets([
        "数据：做 P-pick 人工质检样本集，统计 pick 偏差和 STA/LTA 失败类型。",
        "Cross：event cross-attention + PGA event context gate + distance bias 消融。",
        "Graph：baseline-only / residual-only / prior-only / kNN / power p 消融。",
        "评估：扩大 split，多 seed，记录 baseline-target corr、residual-target corr、attention top-k distance。"]) + "</div><div>" + card("优先级", "先确认 P-pick 与 label 对齐，再做 graph/cross 的结构消融。否则模型问题和数据问题会混在一起。") + "</div></div>")
    slide("希望专家重点建议", bullets([
        "P-pick 的 STA/LTA refine 是否合理，哪些样本应该剔除或降权？",
        "PGA target 是否必须显式使用 event 信息？event 信息应如何进入模型？",
        "距离衰减先验更适合作为 baseline、attention bias，还是 residual feature？",
        "是否需要加入台站场地项、方位角、路径效应、区域衰减参数？",
        "目前 overfit 诊断到哪一步可以转向更大训练集？"]))

    body = "\n".join(slides)
    html_text = f"""<!doctype html>
<html lang="zh-CN">
<head>
<meta charset="utf-8" />
<meta name="viewport" content="width=device-width, initial-scale=1" />
<title>DiTing TEAM Graph PGA Technical Exchange</title>
<style>
html, body {{ margin:0; background:#111827; color:#1f2937; font-family: Arial, "Microsoft YaHei", sans-serif; }}
.deck {{ height:100vh; overflow:hidden; }}
.track {{ display:flex; height:100vh; transition: transform .25s ease; }}
.slide {{ box-sizing:border-box; flex:0 0 100vw; height:100vh; background:#f8fafc; padding:3.8vh 4.3vw; display:flex; flex-direction:column; gap:2.0vh; }}
header h1 {{ margin:0; color:#1f4e79; font-size:3.4vh; line-height:1.15; }}
header p {{ margin:.7vh 0 0; color:#64748b; font-size:1.8vh; }}
ul {{ margin:0; padding-left:1.4em; font-size:2.05vh; line-height:1.55; }}
li {{ margin:.45vh 0; }}
p {{ font-size:1.8vh; line-height:1.45; }}
.hero, .grid2, .cards4 {{ flex:1; min-height:0; }}
.hero {{ display:grid; grid-template-columns:1fr 1fr; gap:2vw; align-items:center; }}
.grid2 {{ display:grid; grid-template-columns:1fr 1fr; gap:2vw; align-items:center; }}
.grid2.wideimg {{ align-items:stretch; }}
.cards4 {{ display:grid; grid-template-columns:repeat(4, 1fr); gap:1.4vw; align-items:stretch; }}
.card {{ background:white; border:1px solid #d9e2ec; border-radius:8px; padding:2.2vh 1.4vw; box-shadow:0 1px 4px rgba(15,23,42,.08); }}
.card h3 {{ margin:0 0 1vh; color:#008080; font-size:2.15vh; }}
.card p {{ margin:0; font-size:1.85vh; }}
img {{ max-width:100%; max-height:70vh; object-fit:contain; display:block; margin:auto; }}
img.full {{ width:100%; max-height:78vh; }}
.note {{ margin:.4vh 0 0; color:#64748b; font-size:1.55vh; }}
table.dataframe {{ width:100%; border-collapse:collapse; font-size:1.45vh; background:white; }}
table.dataframe th {{ background:#1f4e79; color:white; padding:.7vh .45vw; text-align:left; }}
table.dataframe td {{ border-bottom:1px solid #e2e8f0; padding:.58vh .45vw; }}
.counter {{ position:fixed; right:1.2vw; bottom:1.5vh; color:white; font-size:1.8vh; opacity:.8; }}
.hint {{ position:fixed; left:1.2vw; bottom:1.5vh; color:white; font-size:1.55vh; opacity:.7; }}
@media print {{
  body {{ background:white; }}
  .deck, .track {{ height:auto; overflow:visible; display:block; }}
  .slide {{ page-break-after:always; width:100vw; height:56.25vw; }}
  .counter, .hint {{ display:none; }}
}}
</style>
</head>
<body>
<div class="deck"><main class="track" id="track">{body}</main></div>
<div class="hint">←/→ 翻页，F 全屏</div><div class="counter" id="counter"></div>
<script>
const slides = [...document.querySelectorAll('.slide')];
let idx = 0;
function show(i) {{
  idx = Math.max(0, Math.min(slides.length - 1, i));
  document.getElementById('track').style.transform = `translateX(${{-idx * 100}}vw)`;
  document.getElementById('counter').textContent = `${{idx + 1}} / ${{slides.length}}`;
}}
document.addEventListener('keydown', e => {{
  if (e.key === 'ArrowRight' || e.key === 'PageDown' || e.key === ' ') show(idx + 1);
  if (e.key === 'ArrowLeft' || e.key === 'PageUp') show(idx - 1);
  if (e.key.toLowerCase() === 'f') document.documentElement.requestFullscreen?.();
}});
show(0);
</script>
</body>
</html>
"""
    HTML_PATH.write_text(html_text, encoding="utf-8")
    return HTML_PATH


def ppt_title(slide, title: str, subtitle: str = ""):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.22), Inches(12.4), Inches(0.58))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(23)
    p.font.bold = True
    p.font.color.rgb = BLUE
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.48), Inches(0.82), Inches(12.1), Inches(0.32))
        p2 = box2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = GRAY


def ppt_bullets(slide, items: list[str], x, y, w, h, size=14):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(35, 35, 35)
        p.space_after = Pt(5)


def ppt_img(slide, path: Path, x, y, w=None, h=None):
    kwargs = {}
    if w is not None:
        kwargs["width"] = Inches(w)
    if h is not None:
        kwargs["height"] = Inches(h)
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), **kwargs)


def ppt_table(slide, df: pd.DataFrame, x, y, w, h, size=7.5):
    base.add_table(slide, df, x, y, w, h, size)


def add_footer(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(7.05), Inches(12.2), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(8.5)
    p.font.color.rgb = RGBColor(120, 120, 120)
    p.alignment = PP_ALIGN.RIGHT


def write_ppt_v2(assets: dict[str, Path], tables: dict[str, pd.DataFrame]) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def new_slide(title: str, subtitle: str = ""):
        s = prs.slides.add_slide(blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = LIGHT
        ppt_title(s, title, subtitle)
        return s

    s = new_slide("DiTing 表征接入多台站 PGA / Event 预测", "Transformer readout 与 graph prior-residual 的实验诊断")
    ppt_bullets(s, ["技术交流：当前进展、问题定位、下一步请专家建议", "新增：P-pick 波形抽样质检页 + 实现级模型结构图"], 0.9, 2.7, 11.5, 1.5, 20)
    add_footer(s, "2026-05-01")

    s = new_slide("这次希望专家帮忙看的点")
    ppt_bullets(s, ["P-pick 是否足够准确，是否会污染 station feature 和 PGA label 对齐。", "event 信息如何进入 PGA：event query、mag/loc context、GMPE prior，还是显式震源参数。", "任意 target PGA 更适合 cross-attention readout，还是 graph message passing。", "小样本 overfit 诊断后，下一步验证集和消融实验如何设计。"], 0.9, 1.4, 11.6, 4.7, 19)

    s = new_slide("当前任务：从单台站 DiTing 到多台站事件/PGA 预测")
    ppt_img(s, assets["task_setup"], 0.45, 1.0, w=12.4)

    s = new_slide("任务设置")
    ppt_bullets(s, ["输入：N 个台站的三分量波形、台站坐标、station_valid mask。", "目标查询：event query 用于震级/位置；PGA target query 用于给定 target station 的 PGA。", "输出：event-level magnitude / location 或震中距；target-level PGA。", "核心区别：多个台站共同约束一个事件和一个空间 PGA 场。"], 0.65, 1.25, 5.2, 4.6, 15)
    ppt_img(s, assets["task_setup"], 6.0, 1.05, w=6.6)

    s = new_slide("数据覆盖与抽样样本")
    ppt_img(s, assets["data_overview"], 0.45, 1.05, w=6.35)
    ppt_img(s, assets["sample_event"], 6.9, 1.05, w=5.95)

    s = new_slide("P-pick 生成流程与元数据诊断")
    ppt_bullets(s, ["当前 P-pick：走时曲线粗定位，再用 STA/LTA 在搜索窗内 refine。", "左图：最终 pick 来源。", "中图：STA/LTA refine 后 pick 与走时粗定位 pick 的时间差。", "右图：P pick 处 STA/LTA ratio 随震中距和 PGA 的分布。"], 0.65, 1.25, 4.9, 4.7, 13)
    ppt_img(s, assets["ppick_stats"], 5.8, 1.05, w=7.0)

    s = new_slide("P-pick 波形抽样：请专家现场看")
    ppt_img(s, assets["ppick_audit"], 0.45, 1.0, w=12.4)
    add_footer(s, "红虚线为当前 P-pick；仅展示 pick 落在当前 100s HDF5 窗口内的样本")

    s = new_slide("实现级模型总览：和后面表格的对应关系")
    ppt_img(s, assets["architecture_detail"], 0.35, 0.95, w=12.65)

    s = new_slide("模型规模与可训练参数")
    ppt_table(s, tables["model_scale"], 0.35, 1.0, 12.65, 3.25, 7.3)
    ppt_bullets(s, ["总参数量主要来自 frozen DiTing encoder；当前实验实际训练 adapter、readout 和 task heads。", "TEAM cross-attention 参数量为本地精确计数；graph prior-residual 参数量为按实现结构的近似估计。"], 0.7, 4.65, 11.8, 1.25, 13.5)

    s = new_slide("Loss 函数设计")
    ppt_table(s, tables["loss_design"], 0.35, 1.0, 12.65, 2.7, 7.4)
    ppt_bullets(s, ["single-station：L = 0.3 L_mag + 0.3 L_epidist + 0.4 L_pga。", "PGA-only full model：只在有效 target PGA 上计算 Huber loss。", "event+PGA full model：L = 0.2 L_mag + 0.2 L_loc + 1.0 L_pga。", "loss 下降不等于 readout 不坍塌，仍需同时看 corr、slope 和 pred std。"], 0.75, 4.15, 11.9, 2.0, 12.8)

    s = new_slide("Single-station model 效果：waveform 表征有效")
    ppt_img(s, assets["single_station_metrics"], 0.45, 1.0, w=6.55)
    single_val = tables["single_station"][tables["single_station"]["Split"].eq("val")]
    ppt_table(s, single_val[["Model", "Task", "MAE", "Corr"]].head(12), 7.25, 1.0, 5.65, 3.75, 7.2)
    ppt_bullets(s, ["single-station mag / epidist / PGA 在 val 上均有明显相关性。", "这说明 waveform encoder + station adapter 不是完全无效；full model 常数化更像 readout/空间传播问题。"], 7.35, 5.0, 5.45, 1.2, 12)

    s = new_slide("全部实验总览：TEAM 14 个，Graph 7 个")
    ppt_img(s, assets["all_experiments"], 0.45, 0.9, w=7.2)
    ppt_bullets(s, ["圆点为 train PGA corr，方块为 val PGA corr。", "灰色表示 eval 失败或没有 PGA metrics。", "后续页面展示代表实验；完整表格输出在 tables/all_experiments_summary.csv。", "overfit_n=32/128 已在表中单独标出。"], 7.95, 1.2, 4.9, 4.2, 14)

    s = new_slide("原始 TEAM / query_transformer readout")
    ppt_img(s, assets["architecture_team"], 0.55, 1.05, w=7.0)
    ppt_bullets(s, ["PGA query = target coordinate embedding + learned PGA query token。", "concat 后进入 full Transformer；PGA key 被 mask，PGA 作为 query-only token。", "风险：读出路径要靠训练自己发现，容易走向均值解。"], 7.75, 1.35, 4.9, 3.4, 14)

    s = new_slide("原始 readout ablation：各模型设置")
    ppt_table(s, tables["ablation_settings"], 0.45, 1.0, 12.45, 3.4, 8.5)

    s = new_slide("原始 full self-attention readout 的坍塌")
    ppt_img(s, assets["team_train_val"], 0.55, 1.05, w=6.2)
    ppt_table(s, tables["transformer_ablation"][["Method", "Train Corr", "Train slope", "Val Corr", "Val slope"]], 6.95, 1.25, 5.95, 2.45, 8)
    ppt_bullets(s, ["query_transformer slope≈0；batch=1 仍常数。", "direct_station / query_no_transformer 是 readout 对照，用来判断 station feature、target query 和 head 的作用。"], 7.05, 4.45, 5.65, 1.2, 13)

    s = new_slide("Attention / mask 诊断结论")
    ppt_img(s, assets["variance"], 0.65, 1.05, w=6.2)
    ppt_bullets(s, ["query_transformer：原始 full self-attention readout。", "mask_batch1：只改 batch=1 和 mask sanity，仍常数。", "query_no_transformer：target query-only 负对照。", "direct_station：直接读 pooled station feature。", "0.11111 诊断 bug 已修：区分 key-type mask 与 effective key mask。"], 7.1, 1.25, 5.4, 4.6, 12.5)

    s = new_slide("Cross-attention readout 设计用意")
    ppt_img(s, assets["architecture_cross"], 0.45, 1.0, w=7.3)
    ppt_bullets(s, ["设计目标：让 target/event query 单向读取 station tokens。", "PGA Q：target coordinate embedding + learned PGA token。", "K/V：station_feature_emb；key_padding_mask=~station_valid。", "输出不加 query residual，使输出值更直接来自 station tokens。"], 8.0, 1.25, 4.6, 4.1, 13)

    s = new_slide("Cross-attention 各实验设置")
    ppt_table(s, tables["cross_settings"], 0.35, 0.95, 12.65, 4.8, 7.2)

    s = new_slide("Cross-attention 实验结果：train 和 val")
    ppt_img(s, assets["cross_train_val_detail"], 0.55, 1.05, w=6.25)
    cross_rows = tables["all_experiments"][tables["all_experiments"]["Route"].str.contains("cross", case=False, na=False)]
    ppt_table(s, cross_rows[["Experiment", "overfit_n", "Train Corr", "Val Corr", "Train MAE", "Val MAE"]].head(9), 6.95, 1.0, 6.0, 3.55, 7.0)

    s = new_slide("Graph prior-residual：实现细节与为什么有效")
    ppt_img(s, assets["graph_detail"], 0.45, 1.0, w=7.3)
    ppt_bullets(s, ["station_pga_prior_head 从 single-station 预训练加载，提供输入台站 PGA scale。", "distance baseline 用 dist^-p 对输入台站 prior 加权。", "graph message passing 只学习 residual，降低从零学习空间场的难度。", "旧 graph 仍坍塌，说明关键不只是 graph，而是 prior + baseline + residual 路径。"], 8.05, 1.25, 4.6, 4.6, 12.5)

    s = new_slide("Graph 路线各实验设置")
    ppt_table(s, tables["graph_settings"], 0.35, 0.95, 12.65, 4.7, 7.4)

    s = new_slide("Graph 路线实验结果：train 和 val")
    ppt_img(s, assets["graph_train_val_detail"], 0.55, 1.05, w=6.25)
    graph_rows = tables["all_experiments"][tables["all_experiments"]["Group"].eq("Graph")]
    ppt_table(s, graph_rows[["Experiment", "Train Corr", "Val Corr", "Train MAE", "Val MAE"]].head(7), 6.95, 1.0, 6.0, 3.1, 7.0)

    s = new_slide("Loss 曲线与收敛情况")
    ppt_img(s, assets["full_loss_curves"], 0.45, 1.0, w=7.25)
    ppt_bullets(s, ["代表 full-model 实验 train/val loss 基本随 epoch 下降，优化过程没有明显发散。", "query_transformer 的 loss 也会下降，但输出仍可坍塌到均值解。", "因此 loss 曲线必须和 corr、slope、pred std 一起解释。"], 8.0, 1.35, 4.6, 3.8, 13)

    s = new_slide("Single-station loss 曲线")
    ppt_img(s, assets["single_station_loss_curves"], 0.45, 1.0, w=7.25)
    ppt_bullets(s, ["single-station 预训练的 train/val loss 稳定下降。", "这支持先训练单台站 waveform 表征和 PGA prior，再作为 graph prior-residual 的输入。", "后续应保存分任务 loss，避免 mag/epidist/PGA 被加权总 loss 掩盖。"], 8.0, 1.35, 4.6, 3.8, 13)

    s = new_slide("两条路线对比")
    compare = pd.DataFrame([
        ["query_transformer", "统一 token 模型", "PGA/event 均易常数化"],
        ["direct_station", "证明 station feature 可用", "不能处理任意 target"],
        ["cross-attention", "target 显式读 station", "空间先验不足"],
        ["graph prior-residual", "single-station prior + 距离传播", "需拆 baseline/residual 贡献"],
    ], columns=["Route", "What it solves", "Open issue"])
    ppt_table(s, compare, 0.7, 1.25, 12.0, 2.35, 10)
    ppt_bullets(s, ["当前更像是 readout 和 inductive bias 问题，而不是 DiTing backbone 完全无效。", "最终路线可能是二者融合：target-to-station readout + distance/GMPE prior 或 graph residual。"], 0.9, 4.3, 11.7, 1.3, 15)

    s = new_slide("当前主要问题")
    ppt_bullets(s, ["P-pick 准确性还没有人工抽查闭环，可能影响波形窗口和 station SNR。", "event mag/loc 联合训练仍会常数化，需要独立定位 event readout。", "PGA pred std 仍小于 target std，模型偏保守。", "验证集太小，多 seed / 多 split 后结论才稳。"], 0.75, 1.3, 5.3, 4.6, 15)
    ppt_img(s, assets["variance"], 6.4, 1.1, w=6.25)

    s = new_slide("下一步实验建议")
    ppt_bullets(s, ["数据：做 P-pick 人工质检样本集，统计 pick 偏差和 STA/LTA 失败类型。", "Cross：event cross-attention + PGA event context gate + distance bias 消融。", "Graph：baseline-only / residual-only / prior-only / kNN / power p 消融。", "评估：扩大 split，多 seed，记录 baseline-target corr、residual-target corr、attention top-k distance。"], 0.75, 1.25, 11.9, 4.9, 16)

    s = new_slide("希望专家重点建议")
    ppt_bullets(s, ["P-pick 的 STA/LTA refine 是否合理，哪些样本应该剔除或降权？", "PGA target 是否必须显式使用 event 信息？event 信息应如何进入模型？", "距离衰减先验更适合作为 baseline、attention bias，还是 residual feature？", "是否需要加入台站场地项、方位角、路径效应、区域衰减参数？", "目前 overfit 诊断到哪一步可以转向更大训练集？"], 0.8, 1.2, 11.7, 5.1, 17)

    prs.save(PPTX_PATH)
    return PPTX_PATH


def main() -> None:
    base.ensure_dirs()
    assets, tables = base.generate_assets()
    all_experiments = build_all_experiment_table()
    tables["all_experiments"] = all_experiments
    tables["ablation_settings"] = setting_table("ablation")
    tables["cross_settings"] = setting_table("cross")
    tables["graph_settings"] = setting_table("graph")
    tables["model_scale"] = build_model_scale_table()
    tables["loss_design"] = build_loss_design_table()
    tables["single_station"] = build_single_station_table()
    assets["task_setup"] = plot_task_setup()
    assets["ppick_audit"] = plot_ppick_audit()
    assets["ppick_stats"] = plot_ppick_stats()
    assets["architecture_detail"] = plot_architecture_detail()
    assets["graph_detail"] = plot_graph_detail()
    assets["all_experiments"] = plot_all_experiment_matrix(all_experiments)
    assets["team_train_val"] = plot_train_val_corr(
        all_experiments,
        "TEAM",
        "team_all_train_val_corr.png",
        "TEAM route experiments: train and val PGA correlation",
    )
    assets["graph_train_val"] = plot_train_val_corr(
        all_experiments,
        "Graph",
        "graph_all_train_val_corr.png",
        "Graph route experiments: train and val PGA correlation",
    )
    assets["cross_train_val_detail"] = plot_train_val_r2_mae(
        all_experiments,
        [
            "pga15_cross_overfit32",
            "pga15_cross_overfit128",
            "pga15_cross_overfit32_fixed_inputs_targets",
            "pga15_cross_overfit32_input_targets",
            "pga15_cross_overfit32_fixed_inputs_random_targets",
        ],
        "cross_train_val_detail.png",
        "Cross-attention PGA experiments: train vs val",
    )
    assets["graph_train_val_detail"] = plot_train_val_r2_mae(
        all_experiments,
        [
            "graph_first_inputs",
            "graph_first_inputs_prior_residual",
            "graph_random_inputs_prior_residual",
            "exp1_graph_single_input_same_station_pga",
            "exp2_graph_single_input_multi_target_pga",
            "exp3_graph_snr_filtered_holdout_pga",
        ],
        "graph_train_val_detail.png",
        "Graph PGA experiments: train vs val",
    )
    assets["single_station_metrics"] = plot_single_station_metrics(tables["single_station"])
    assets["full_loss_curves"] = plot_representative_loss_curves()
    assets["single_station_loss_curves"] = plot_single_station_loss_curves()
    html_path = write_html_deck(assets, tables)
    pptx_path = write_ppt_v2(assets, tables)
    print(f"HTML: {html_path}")
    print(f"PPTX: {pptx_path}")
    print(f"Assets: {ASSET_DIR}")


if __name__ == "__main__":
    main()
